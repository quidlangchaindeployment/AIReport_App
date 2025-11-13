# ---
# app.py (AI Data Analysis App - Refactored Version)
#
# このコードは、商用利用が容易な寛容な（permissive）ライセンス
# (例: MIT, Apache License 2.0, BSD) の下で利用可能な
# ライブラリ、またはライセンスに依存しないコードのみを使用して実装されています。
# GPL, AGPL, SSPLなどのコピーレフト効果を持つライブラリは使用していません。
# ---

# --- 1. ライブラリのインポート ---
import streamlit as st
import pandas as pd
import numpy as np
import os
import re
import json
import logging
import time
import spacy
import altair as alt
import networkx as nx
from networkx.algorithms import community
from pyvis.network import Network
import streamlit.components.v1 as components
from sklearn.decomposition import PCA
from sklearn.preprocessing import StandardScaler
from io import StringIO, BytesIO
from typing import Optional, Dict, List, Any, Union, Set
from dotenv import load_dotenv  # (★) .env読み込みのために追加
import matplotlib
matplotlib.use('Agg') # (★) Streamlitのバックエンドで動作させるためのおまじない
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import base64
from wordcloud import WordCloud

# (★) --- matplotlib 日本語フォント設定 ---
# DockerfileでインストールしたIPAフォントのパスを指定
# (★) ご注意: ローカル環境で実行する場合、このフォントパスが異なる場合があります。
# (★) Docker環境 (Debianベース) を前提としています。
try:
    # (★) Dockerfile内のパス
    font_path = '/usr/share/fonts/opentype/ipafont-gothic/ipagp.ttf' 
    if os.path.exists(font_path):
        fm.fontManager.addfont(font_path)
        plt.rcParams['font.family'] = 'IPAGothic'
        # logger.info(f"日本語フォント '{font_path}' を読み込みました。")
    else:
        # (★) フォールバック (環境によって調整が必要)
        logger.warning(f"日本語フォント '{font_path}' が見つかりません。グラフの日本語が文字化けする可能性があります。")
        # (★) 代替フォントの検索 (やや時間がかかるが堅牢)
        try:
            jp_font = fm.findfont(fm.FontProperties(family='IPAexGothic'))
            plt.rcParams['font.family'] = 'IPAexGothic'
            logger.info(f"代替フォント '{jp_font}' を使用します。")
        except:
             logger.error("代替の日本語フォントも見つかりません。")
             plt.rcParams['font.family'] = 'sans-serif'
except Exception as e:
    logger.error(f"matplotlib日本語フォント設定エラー: {e}")

# (★) LangChain / Google Generative AI のインポート
# ライセンス: Apache License 2.0
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.prompts import PromptTemplate
from langchain_core.output_parsers import StrOutputParser

# (★) Step D (PowerPoint生成) で必要となるライブラリ
# ライセンス: MIT License
try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    st.error(
        "PowerPoint生成ライブラリ(python-pptx)が見つかりません。"
        "pip install python-pptx を実行してください。"
    )

# (★) Step D (ドラッグ＆ドロップUI) で必要となるライブラリ
# ライセンス: MIT License
try:
    from streamlit_sortables import sort_items
except ImportError:
    st.error(
        "UIライブラリ(streamlit-sortables)が見つかりません。"
        "pip install streamlit-sortables を実行してください。"
    )

# 既存のライブラリ (openpyxl, ja_core_news_sm) のインポート
try:
    import openpyxl
except ImportError:
    st.error("Excel (openpyxl) がインストールされていません。`pip install openpyxl` してください。")
try:
    import ja_core_news_sm
except ImportError:
    st.error("spaCy日本語モデル (ja_core_news_sm) が見つかりません。`python -m spacy download ja_core_news_sm` してください。")

# --- 2. (★) 定数定義 ---

# (★) 要件に基づき、使用するAIモデルを定数として定義
MODEL_FLASH_LITE = "gemini-2.5-flash-lite" # Step A, B (高速・効率的)
MODEL_FLASH = "gemini-2.5-flash"         # Step D (代替)
MODEL_PRO = "gemini-2.5-pro"             # Step C, D (高品質)

# バッチサイズと待機時間 (KISS)
FILTER_BATCH_SIZE = 50
FILTER_SLEEP_TIME = 6.1  # Rate Limit 対策 (10 requests per 60 seconds)
TAGGING_BATCH_SIZE = 50  
TAGGING_SLEEP_TIME = 6.1  # Rate Limit 対策

# 地名辞書
try:
    from geography_db import JAPAN_GEOGRAPHY_DB
except ImportError:
    st.error("地名辞書ファイル (geography_db.py) が見つかりません。")
    JAPAN_GEOGRAPHY_DB = {}

COLOR_PALETTE = [
    "#FF5733", "#33FF57", "#3357FF", "#FF33A1", "#33FFF6",
    "#F3FF33", "#FF8C33", "#8C33FF", "#33FF8C", "#FF338C"
]

# --- 3. ロガー設定 ---
class StreamlitLogHandler(logging.Handler):
    """Streamlitのセッションステートにログメッセージを追加するハンドラ"""
    def __init__(self):
        super().__init__()
        if 'log_messages' not in st.session_state:
            st.session_state.log_messages = []

    def emit(self, record):
        log_entry = self.format(record)
        st.session_state.log_messages.append(log_entry)
        st.session_state.log_messages = st.session_state.log_messages[-500:]

logger = logging.getLogger(__name__)
if not logger.handlers:
    logger.setLevel(logging.INFO)
    handler = StreamlitLogHandler()
    handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
    logger.addHandler(handler)


# --- 4. (★) AIモデル・NLPモデルのキャッシュ管理 ---

# (★) 要件に基づき、異なるモデル名を指定してLLMをロードする関数に刷新
@st.cache_resource(ttl=3600)
def get_llm(
    model_name: str, 
    temperature: float = 0.0,
    timeout_seconds: int = 120  # (★) --- 修正: タイムアウト引数を追加 (デフォルト120秒) ---
) -> Optional[ChatGoogleGenerativeAI]:
    """
    指定されたモデル名、温度、タイムアウトでLLM (Google Gemini) モデルをロード・キャッシュする。
    """
    try:
        api_key = os.getenv("GOOGLE_API_KEY")
        if not api_key:
            logger.error(f"get_llm: GOOGLE_API_KEY がありません (Model: {model_name})")
            return None

        llm = ChatGoogleGenerativeAI(
            model=model_name,
            temperature=temperature,
            convert_system_message_to_human=True,
            api_key=api_key,
            request_timeout=timeout_seconds # (★) --- 修正: タイムアウトを渡す ---
        )
        logger.info(f"LLM Model ({model_name}) loaded successfully (Timeout: {timeout_seconds}s).")
        return llm
    except Exception as e:
        logger.error(f"LLM ({model_name}) の初期化に失敗: {e}", exc_info=True)
        st.error(f"AIモデル ({model_name}) のロードに失敗しました: {e}")
        return None

@st.cache_resource
def load_spacy_model() -> Optional[spacy.language.Language]:
    """spaCyの日本語モデル(ja_core_news_sm)をロード・キャッシュする"""
    try:
        logger.info("Loading spaCy model (ja_core_news_sm)...")
        nlp = spacy.load("ja_core_news_sm")
        logger.info("spaCy model loaded successfully.")
        return nlp
    except Exception as e:
        logger.error(f"Failed to load spaCy model: {e}", exc_info=True)
        return None

@st.cache_data(ttl=3600) # 1時間キャッシュ
def get_analysis_tips_list_from_ai() -> List[str]:
    """
    (★) 待機時間中に表示する「データ分析に関するTips」をAIで生成する。
    モデル: MODEL_FLASH_LITE (gemini-2.5-flash-lite)
    """
    logger.info("get_analysis_tips_list_from_ai: AI (Flash Lite) でTIPSを生成します。")
    llm = get_llm(model_name=MODEL_FLASH_LITE, temperature=0.5)
    if llm is None:
        return ["AIモデルの読み込みに失敗しました。"]

    prompt = PromptTemplate.from_template(
        """
        あなたはデータサイエンティストのメンターです。
        データ分析の初心者から中級者に向けて、役立つ「ヒントやTIPS」を【10個】、JSONのリスト形式で生成してください。
        
        # 指示:
        1. 各TIPSは、1〜2文の簡潔な日本語の文字列にすること。
        2. 例: 「'平均値'だけでなく'中央値'も見ることで、外れ値の影響を把握できます。」
        3. 例: 「データを可視化する前に、まずデータの'欠損値'と'型'を確認しましょう。」
        4. 出力はJSONリスト形式（ ["TIPS1", "TIPS2", ...] ）のみ。
        
        # 回答 (JSONリスト形式のみ):
        """
    )
    chain = prompt | llm | StrOutputParser()
    
    try:
        response_str = chain.invoke({})
        logger.debug(f"AI TIPS生成(生): {response_str}")
        
        # マークダウンや不要なテキストを除去し、JSONのみを抽出
        match = re.search(r'\[.*\]', response_str, re.DOTALL)
        if not match:
            logger.warning("AIがTIPSリスト(JSON)の生成に失敗しました。")
            return ["分析TIPSの取得に失敗しました。"]
        
        json_str = match.group(0)
        tips_list = json.loads(json_str)
        
        if isinstance(tips_list, list) and all(isinstance(tip, str) for tip in tips_list):
            logger.info(f"AI TIPS {len(tips_list)}件の生成に成功。")
            return tips_list
        else:
            raise Exception("AIの回答が文字列のリスト形式ではありません。")
            
    except Exception as e:
        logger.error(f"AI TIPS生成中にエラー: {e}", exc_info=True)
        return [
            "TIPSの読み込みに失敗しました。",
            "データ分析は、まず目的（KGI/KPI）を明確にすることから始まります。",
            "「なぜ？」を5回繰り返すことで、分析の真の目的にたどり着くことがあります。",
            "良い分析は、良い「問い」から生まれます。",
            "データは「集める」ことより「どう使うか」が重要です。"
        ]

# --- 5. ファイル読み込みヘルパー ---
def read_file(file: st.runtime.uploaded_file_manager.UploadedFile) -> (Optional[pd.DataFrame], Optional[str]):
    """アップロードされたファイル(Excel/CSV)をPandas DataFrameとして読み込む"""
    file_name = file.name
    logger.info(f"ファイル読み込み開始: {file_name}")
    try:
        if file_name.endswith('.csv'):
            # 文字コードを自動判別
            try:
                content = file.getvalue().decode('utf-8-sig')
            except UnicodeDecodeError:
                logger.warning(f"UTF-8-SIGデコード失敗。CP932で再試行: {file_name}")
                content = file.getvalue().decode('cp932')
            df = pd.read_csv(StringIO(content))

        elif file_name.endswith(('.xlsx', '.xls')):
            df = pd.read_excel(BytesIO(file.getvalue()), engine='openpyxl')
        else:
            msg = f"サポート外のファイル形式: {file_name}"
            logger.warning(msg)
            return None, msg

        logger.info(f"ファイル読み込み成功: {file_name}")
        return df, None
    except Exception as e:
        logger.error(f"ファイル読み込みエラー ({file_name}): {e}", exc_info=True)
        st.error(f"ファイル「{file_name}」の読み込み中にエラー: {e}")
        return None, f"読み込みエラー: {e}"


# --- 6. (★) Step A: AIタグ付け関連関数 ---
# (要件: Step Aは gemini-2.5-flash-lite を使用)

def get_dynamic_categories(analysis_prompt: str) -> Optional[Dict[str, str]]:
    """
    (Step A) ユーザーの分析指針に基づき、AIが動的なカテゴリをJSON形式で生成する。
    (★) モデル: MODEL_FLASH_LITE
    """
    llm = get_llm(model_name=MODEL_FLASH_LITE, temperature=0.0)
    if llm is None:
        logger.error("get_dynamic_categories: LLM (Flash Lite) が利用できません。")
        st.error("AIモデル(Flash Lite)が利用できません。サイドバーでAPIキーを設定してください。")
        return None

    logger.info("動的カテゴリ生成AI (Flash Lite) を呼び出し...")
    
    prompt = PromptTemplate.from_template(
        """
        あなたはデータ分析のスキーマ設計者です。「分析指針」を読み、テキストから抽出するべき「カテゴリ」を考案してください。
        
        # 分析指針:
        {user_prompt}

        # 指示:
        1.  「分析指針」を注意深く読みます。
        2.  もし「分析指針」が【ユーザー定義のカテゴリ名と説明】（例: 「①話題カテゴリ：...」や「②観光地：...」）を明示的に指定している場合、**その指示に厳密に従い**、指定されたカテゴリ名（例: "話題カテゴリ"）と説明（例: "どの話題に関する言及か..."）を抽出してください。
        3.  もし「分析指針」がカテゴリを明示的に指定していない場合（例: 「広島の観光について分析したい」）、分析指針のトピックを元に、抽出するべきカテゴリ（キー）とそ
            の説明（値）を【あなた自身で考案】してください。
        4.  「市区町村」や「地名」に関するカテゴリは、必須カテゴリとして自動で追加されるため、**絶対に考案・抽出しないでください**。
        5.  出力は【厳格なJSON辞書形式】 `{{ "カテゴリ名1": "カテゴリの説明1", "カテゴリ名2": "カテゴリの説明2" }}` のみとします。
        6.  該当するカテゴリが（地名以外に）無い場合は、空のJSON `{{}}` を返してください。

        # 回答 (JSON辞書形式のみ):
        """
    )
    
    chain = prompt | llm | StrOutputParser()
    try:
        response_str = chain.invoke({"user_prompt": analysis_prompt})
        logger.debug(f"AIカテゴリ定義(生): {response_str}")

        # マークダウンや不要なテキストを除去し、JSONのみを抽出
        match = re.search(r'\{.*\}', response_str, re.DOTALL)
        if not match:
            logger.warning("AIがJSON形式で応答しませんでした。")
            return None
        
        json_str = match.group(0).replace("'", '"')
        try:
            categories = json.loads(json_str)
            return categories
        except json.JSONDecodeError as json_e:
            logger.error(f"AI応答のJSONパース失敗: {json_e} - Raw: {json_str}")
            return None
            
    except Exception as e:
        logger.error(f"AIカテゴリ生成中にエラー: {e}", exc_info=True)
        st.error(f"AIカテゴリ生成中にエラーが発生しました: {e}")
        return None

def filter_relevant_data_by_ai(df_batch: pd.DataFrame, analysis_prompt: str) -> pd.DataFrame:
    """
    (Step A) AIを使い、分析指針と無関係な行をフィルタリングする (relevant: true/false)。
    (★) モデル: MODEL_FLASH_LITE (gemini-2.5-flash-lite)
    (★) 要件: 進捗表示 (この関数はバッチ処理の一部として呼ばれ、呼び出し元の
          `render_step_a` 内の `update_progress_ui` で進捗が表示される)
    """
    # (★) Step A の要件に基づき、FLASH_LITE モデルを明示的に指定
    llm = get_llm(model_name=MODEL_FLASH_LITE, temperature=0.0)
    if llm is None:
        logger.error("filter_relevant_data_by_ai: LLM (Flash Lite) が利用できません。")
        st.error("AIモデル(Flash Lite)が利用できません。APIキーを確認してください。")
        return pd.DataFrame()  # 空のDF

    logger.debug(f"{len(df_batch)}件 AI関連性フィルタリング (Flash Lite) 開始...")

    # テキストが長すぎる場合、先頭500文字に切り詰める
    input_texts_jsonl = df_batch.apply(
        lambda row: json.dumps(
            {"id": row['id'], "text": str(row['ANALYSIS_TEXT_COLUMN'])[:500]},
            ensure_ascii=False
        ),
        axis=1
    ).tolist()

    prompt = PromptTemplate.from_template(
        """
        あなたはデータ分析のキュレーターです。「分析指針」に基づき、「テキストデータ(JSONL)」の各行が分析対象として【関連しているか (relevant: true)】、【無関係か (relevant: false)】を判定してください。
        # 分析指針 (Analysis Scope):
        {analysis_prompt}
        # 指示:
        1. 「分析指針」と【強く関連】する投稿のみを `true` とする。
        2. 単なる宣伝、挨拶のみ、指針と無関係な地域の言及は `false` とする。
        3. 出力は【JSONL形式のみ】（id と relevant (boolean) を含む辞書）。
        # テキストデータ (JSONL):
        {text_data_jsonl}
        # 回答 (JSONL形式のみ):
        """
    )
    chain = prompt | llm | StrOutputParser()
    
    try:
        invoke_params = {
            "analysis_prompt": analysis_prompt,
            "text_data_jsonl": "\n".join(input_texts_jsonl)
        }
        response_str = chain.invoke(invoke_params)
        
        results = []
        # マークダウン ````jsonl ... ``` を除去
        match = re.search(r'```(?:jsonl|json)?\s*([\s\S]*?)\s*```', response_str, re.DOTALL)
        jsonl_content = match.group(1).strip() if match else response_str.strip()

        for line in jsonl_content.strip().split('\n'):
            cleaned_line = line.strip()
            if not cleaned_line: continue
            try:
                data = json.loads(cleaned_line)
                # relevant が "true" (str) や true (bool) など揺らぎがあるため堅牢に処理
                is_relevant = False
                if isinstance(data.get("relevant"), bool):
                    is_relevant = data.get("relevant")
                elif isinstance(data.get("relevant"), str):
                    is_relevant = data.get("relevant").lower() == 'true'
                
                results.append({"id": data.get("id"), "relevant": is_relevant})
            except (json.JSONDecodeError, AttributeError) as json_e:
                logger.warning(f"AIフィルタリング回答パース失敗: {cleaned_line} - Error: {json_e}")
                # パース失敗時は、IDが特定できれば関連あり(True)としてフォールバック
                id_match = re.search(r'"id":\s*(\d+)', cleaned_line)
                if id_match:
                    results.append({"id": int(id_match.group(1)), "relevant": True})

        return pd.DataFrame(results) if results else pd.DataFrame(columns=['id', 'relevant'])
        
    except Exception as e:
        logger.error(f"AIフィルタリングバッチ処理中エラー: {e}", exc_info=True)
        st.error(f"AIフィルタリング処理エラー: {e}")
        # エラー時は安全側に倒し、すべて関連あり(True)として返す
        return df_batch[['id']].copy().assign(relevant=True)
@st.cache_data(ttl=3600)
def get_location_normalization_maps(
    db: Dict[str, List[str]], 
    analysis_prompt_str: str
) -> (Dict[str, str], Set[str], Set[str]): # (★) 型ヒントに Set[str] を追加
    """
    (★) Step A 改善: 地名正規化用の辞書を動的生成する
    JAPAN_GEOGRAPHY_DB 全体をスキャンし、エイリアス辞書と曖昧な単語セットを作成する
    """
    if not db:
        return {}, set(), set() # (★) 3つの値を返す

    logger.info("地名正規化マップの動的生成開始...")
    alias_to_city_map = {} # {"日光": "日光市", "尾道": "尾道市"}
    ambiguous_keys = set() # {"広島", "東京", "札幌"}
    prefectures = set() # {"広島県", "東京都"}
    all_cities_wards = set() # {"広島市", "中区", "日光市"}
    
    # 1. DB全体をスキャン
    for key, values in db.items():
        if not isinstance(values, list): continue

        key_normalized = key.replace("県", "").replace("都", "").replace("府", "").replace("市", "")
        
        # 1a. 都道府県/政令市キーの処理
        if "県" in key or "都" in key or "府" in key:
            prefectures.add(key)
            ambiguous_keys.add(key_normalized) # "北海道", "東京", "大阪"
        elif "市" in key and values and "区" in values[0]: # 政令市
            ambiguous_keys.add(key_normalized) # "札幌"
            all_cities_wards.add(key) # "札幌市"
        
        # 1b. 値リスト (市区町村) の処理
        for city_or_ward in values:
            all_cities_wards.add(city_or_ward) # "函館市", "中央区"
            
            # (★) "日光市" -> "日光" のようなエイリアスを動的生成
            alias = city_or_ward.replace("市", "").replace("区", "").replace("町", "").replace("村", "")
            
            if alias != city_or_ward:
                # "中央" や "南" のような汎用的な区名は、曖昧キーとして処理
                if "区" in city_or_ward and len(alias) <= 2: 
                     ambiguous_keys.add(alias)
                # "日光" -> "日光市" のマッピング
                elif alias not in alias_to_city_map:
                    alias_to_city_map[alias] = city_or_ward
                else:
                    # (★) "府中" (東京都/広島県) のような重複エイリアスは曖昧キーに
                    if alias in alias_to_city_map: # (★) 修正: del の前に存在確認
                        del alias_to_city_map[alias]
                    ambiguous_keys.add(alias)

    # 2. 曖昧キーから、分析指針で特定できるものを救出
    prompt_lower = analysis_prompt_str.lower()
    relevant_cities = []
    for city_key in db.keys():
        if "市" in city_key and db[city_key] and "区" in db[city_key][0]:
             city_name_only = city_key.replace("市", "") # "広島"
             if city_name_only in prompt_lower:
                 relevant_cities.append(city_key) # "広島市"

    if relevant_cities:
        logger.info(f"分析指針から関連都市を特定: {relevant_cities}")
        for city in relevant_cities:
            for ward in db[city]: # "中区", "南区" ...
                # (★) "中区" -> "広島市 中区" というマッピングを作成
                alias_to_city_map[ward] = f"{city} {ward}"
                
                # "中" のようなエイリアスも "広島市 中区" に
                ward_alias = ward.replace("区", "")
                if ward_alias in ambiguous_keys:
                    alias_to_city_map[ward_alias] = f"{city} {ward}"

    # 3. 曖昧なキー (市/県/都/府 を取ったもの) と都道府県名は除外対象
    final_ambiguous_set = ambiguous_keys.union(prefectures)
    
    logger.info(f"地名正規化マップ動的生成完了。エイリアス: {len(alias_to_city_map)}件, 曖昧キー: {len(final_ambiguous_set)}件")
    
    # (★) --- [修正] 3つの値を返す ---
    return alias_to_city_map, final_ambiguous_set, all_cities_wards

def perform_ai_tagging(
    df_batch: pd.DataFrame,
    categories_to_tag: Dict[str, str],
    analysis_prompt: str = ""
) -> pd.DataFrame:
    """
    (Step A) テキストのバッチを受け取り、AIが【指定されたカテゴリ定義】に基づいて直接タグ付けを行う
    (★) モデル: MODEL_FLASH_LITE
    (★) 改善: AIはキーワード抽出に専念し、Python側で地名を正規化する
    """
    llm = get_llm(model_name=MODEL_FLASH_LITE, temperature=0.0)
    if llm is None:
        logger.error("perform_ai_tagging: LLM (Flash Lite) が利用できません。")
        st.error("AIモデル(Flash Lite)が利用できません。APIキーを確認してください。")
        return pd.DataFrame()

    logger.info(f"{len(df_batch)}件 AIタグ付け (Flash Lite) 開始 (カテゴリ: {list(categories_to_tag.keys())})")

    # 地名辞書のコンテキストを準備 (分析指針に関連する地名のみをAIに渡す)
    geo_context_str = "{}"
    if JAPAN_GEOGRAPHY_DB and "市区町村キーワード" in categories_to_tag:
        try:
            relevant_geo_db = {}
            prompt_lower = analysis_prompt.lower()
            
            hints = ["広島", "福岡", "大阪", "東京", "北海道", "愛知", "宮城", "札幌", "横浜", "名古屋", "京都", "神戸", "仙台"]
            keys_found = [
                key for key in JAPAN_GEOGRAPHY_DB.keys()
                if any(h in key.lower() for h in hints) and any(h in prompt_lower for h in hints)
            ]
            if "広島" in prompt_lower: keys_found.extend(["広島県", "広島市"])
            if "東京" in prompt_lower: keys_found.extend(["東京都", "東京23区"])
            if "大阪" in prompt_lower: keys_found.extend(["大阪府", "大阪市"])

            for key in set(keys_found): 
                if key in JAPAN_GEOGRAPHY_DB:
                    relevant_geo_db[key] = JAPAN_GEOGRAPHY_DB[key]
            
            if not relevant_geo_db:
                logger.warning("地名辞書の絞り込みヒントなし。主要都市のみ渡します。")
                default_keys = ["東京都", "東京23区", "大阪府", "大阪市", "広島県", "広島市", "福岡県", "福岡市"]
                for key in default_keys:
                    if key in JAPAN_GEOGRAPHY_DB:
                        relevant_geo_db[key] = JAPAN_GEOGRAPHY_DB[key]

            geo_context_str = json.dumps(relevant_geo_db, ensure_ascii=False, indent=2)
            
            if len(geo_context_str) > 5000:
                logger.warning(f"地名辞書が大きすぎ ({len(geo_context_str)}B)。キーのみに縮小。")
                geo_context_str = json.dumps(list(relevant_geo_db.keys()), ensure_ascii=False)
                
            logger.info(f"AIに渡す地名辞書(絞込済): {list(relevant_geo_db.keys())}")
            
            # (★) --- [修正] 3つの値を受け取る ---
            alias_map, ambiguous_set, all_cities_wards = get_location_normalization_maps(JAPAN_GEOGRAPHY_DB, analysis_prompt)
            # (★) --- ここまで ---

        except Exception as e:
            logger.error(f"地名辞書の準備中にエラー: {e}", exc_info=True)
            geo_context_str = "{}" 
            # (★) --- [修正] 3つの変数を初期化 ---
            alias_map, ambiguous_set, all_cities_wards = {}, set(), set()
            
    else:
        # (★) --- [修正] 3つの変数を初期化 ---
        alias_map, ambiguous_set, all_cities_wards = {}, set(), set() 

    # テキストが長すぎる場合、先頭500文字に切り詰める
    input_texts_jsonl = df_batch.apply(
        lambda row: json.dumps(
            {"id": row['id'], "text": str(row['ANALYSIS_TEXT_COLUMN'])[:500]},
            ensure_ascii=False
        ),
        axis=1
    ).tolist()

    # (★) プロンプトは変更なし (AIはキーワードを抽出するだけ)
    prompt = PromptTemplate.from_template(
        """
        あなたは高精度データ分析アシスタントです。「カテゴリ定義」「地名辞書」「分析指針」に基づき、キーワードを抽出します。
        # 分析指針 (Analysis Scope): {analysis_prompt}
        # 地名辞書 (JAPAN_GEOGRAPHY_DB): {geo_context}
        # カテゴリ定義 (categories): {categories}
        # テキストデータ (JSONL): {text_data_jsonl}

        # 指示:
        1. 「テキストデータ(JSONL)」の各行を処理する。
        2. 「カテゴリ定義」のキー名を【厳格に】使用し、全カテゴリを抽出する。
        3. 【すべてのカテゴリ】 ( "市区町村キーワード" を含む ):
           - 値は【単一の文字列】で出力する (該当なければ空文字列 "")。リスト形式は【厳禁】。
           - 文脈から最も関連性の高いものを【1つだけ】選ぶ。
           - 分析指針でカテゴリとその内容の選択肢が提示されている場合は、それに従いラベル付けを行う。
        4. 【"市区町村キーワード" の特別ルール】(★ 変更点):
           - テキストから、最も関連性が高い地名キーワード（例：「宮島」「日光」「尾道」「中区」「広島市」）を【1つだけそのまま】抽出する。
           - 【変換処理は不要】です（例：「宮島」を「廿日市市」に変換しないでください）。
           - 曖昧な表現（例：「広島」）や都道府県名（例：「広島県」）も、もしそれが最も関連性が高いと判断した場合は、そのまま抽出してください。
           - 「分析指針」と無関係な地域の地名は【抽出しない】。
        5. ハルシネーション（情報の捏造）禁止。
        6. 出力は【JSONL形式のみ】（id と categories を含む辞書）。

        # 回答 (JSONL形式のみ):
        """
    )
    chain = prompt | llm | StrOutputParser()
    
    try:
        invoke_params = {
            "categories": json.dumps(categories_to_tag, ensure_ascii=False),
            "geo_context": geo_context_str,
            "text_data_jsonl": "\n".join(input_texts_jsonl),
            "analysis_prompt": analysis_prompt
        }
        response_str = chain.invoke(invoke_params)
        logger.debug(f"AI Tagging - Raw response received.")

        results = []
        expected_keys = list(categories_to_tag.keys())
        
        match = re.search(r'```(?:jsonl|json)?\s*([\s\S]*?)\s*```', response_str, re.DOTALL)
        jsonl_content = match.group(1).strip() if match else response_str.strip()

        for line in jsonl_content.strip().split('\n'):
            cleaned_line = line.strip()
            if not cleaned_line: continue
            try:
                data = json.loads(cleaned_line)
                row_result = {"id": data.get("id")}
                tag_source = data.get('categories', data)
                
                if not isinstance(tag_source, dict):
                    raise json.JSONDecodeError(f"tag_source is not a dict: {tag_source}", "", 0)

                for key in expected_keys:
                    found_key = None
                    for resp_key in tag_source.keys():
                        if str(resp_key).strip() == key:
                            found_key = resp_key
                            break
                    
                    raw_value = tag_source.get(found_key) if found_key else None
                    processed_value = ""
                    if isinstance(raw_value, list) and raw_value:
                        processed_value = str(raw_value[0]).strip()
                    elif raw_value is not None and str(raw_value).strip():
                        processed_value = str(raw_value).strip()
                    
                    if processed_value.lower() in ["該当なし", "none", "null", "", "n/a"]:
                        processed_value = ""
                    
                    # (★) --- [修正] Python 地名正規化ロジック (all_cities_wards が使える) ---
                    if key == "市区町村キーワード" and processed_value:
                        
                        # 1. エイリアスマップで変換 (例: "日光" -> "日光市", "中区" -> "広島市 中区")
                        if processed_value in alias_map:
                            processed_value = alias_map[processed_value]
                        
                        # 2. 曖昧なキー (例: "広島", "東京", "札幌") は破棄
                        elif processed_value in ambiguous_set:
                            logger.debug(f"地名正規化: 曖昧なキー '{processed_value}' を破棄しました。")
                            processed_value = ""
                        
                        # 3. 都道府県名 (例: "広島県") は破棄 (ambiguous_set に含まれる)
                        
                        # 4. DBに存在する正式名称 (例: "広島市") か確認
                        elif processed_value in all_cities_wards:
                            pass # (例: "広島市" はそのまま通す)
                        
                        # 5. それ以外 (例: "アメリカ") は破棄
                        else:
                            # (★) ただし、"広島市 中区" のような「市 区」形式は許可
                            if " " in processed_value and any(s in processed_value for s in ["市", "区"]):
                                pass
                            else:
                                logger.debug(f"地名正規化: 不明なキー '{processed_value}' を破棄しました。")
                                processed_value = ""
                    # (★) --- 正規化ロジックここまで ---

                    row_result[key] = processed_value
                
                results.append(row_result)
                
            except (json.JSONDecodeError, AttributeError) as json_e:
                logger.warning(f"AIタグ付け回答パース失敗: {cleaned_line} - Error: {json_e}")
                id_match = re.search(r'"id":\s*(\d+)', cleaned_line)
                if id_match:
                    results.append({"id": int(id_match.group(1))})
                    
        return pd.DataFrame(results) if results else pd.DataFrame(columns=['id'] + list(expected_keys))

    except Exception as e:
        logger.error(f"AIタグ付けバッチ処理中エラー: {e}", exc_info=True)
        st.error(f"AIタグ付け処理エラー: {e}")
        return pd.DataFrame()

def perform_ai_location_inference(
    df_batch: pd.DataFrame,
    analysis_prompt: str,
    normalization_maps: tuple
) -> pd.DataFrame:
    """
    (★ Pass 2) AIを使い、投稿内容や他カテゴリから間接的に地名を「推論」する
    (★) モデル: MODEL_FLASH_LITE
    """
    llm = get_llm(model_name=MODEL_FLASH_LITE, temperature=0.0)
    if llm is None:
        logger.error("perform_ai_location_inference: LLM (Flash Lite) が利用できません。")
        return pd.DataFrame()

    logger.info(f"{len(df_batch)}件 AI地名推論 (Flash Lite) 開始...")
    
    # (★) 正規化マップをアンパック (AIの回答を検証するため)
    alias_map, ambiguous_set, all_cities_wards = normalization_maps

    # AIに渡すコンテキストをJSONL形式で作成
    # (★) 投稿本文 + Pass 1 でタグ付けされた全カテゴリ をコンテキストにする
    def create_context(row):
        context_data = {
            "id": row['id'],
            "text": str(row['ANALYSIS_TEXT_COLUMN'])[:500]
        }
        # 他のカテゴリ（観光地、農産品など）をヒントとして追加
        other_tags = {
            k: v for k, v in row.items() 
            if k not in ['id', 'ANALYSIS_TEXT_COLUMN', '市区町村キーワード'] and pd.notna(v) and str(v).strip()
        }
        if other_tags:
            context_data["other_tags_context"] = other_tags
        return json.dumps(context_data, ensure_ascii=False)

    input_contexts_jsonl = df_batch.apply(create_context, axis=1).tolist()

    prompt = PromptTemplate.from_template(
        """
        あなたは日本の地理に精通した地名推論AIです。
        「分析指針」と「コンテキストデータ(JSONL)」を読み、各データから最も関連性の高い「市区町村名」を【1つだけ】推論してください。

        # 分析指針 (Analysis Scope): {analysis_prompt}
        # コンテキストデータ (JSONL):
        {text_data_jsonl}

        # 指示:
        1.  `text`（投稿本文）や `other_tags_context`（他のカテゴリのヒント）を注意深く読みます。
        2.  ヒントから、最も可能性の高い「市区町村名」を【1つだけ】推論します。
            (例: "厳島神社" -> "廿日市市")
            (例: "那須の牛乳" -> "那須塩原市" または "那須町")
            (例: "テキストに「札幌」とあり、other_tagsに「中央」とあれば" -> "札幌市 中央区")
        3.  推論した地名キーワード（例: "廿日市市"）を `inferred_location` として返します。
        4.  地名の変換は不要です（例：「廿日市市」を「廿日市」にしないでください）。
        5.  推論できない、または「分析指針」と無関係な場合は `null` を返します。
        6.  出力は【JSONL形式のみ】（id と inferred_location (string or null) を含む辞書）。

        # 回答 (JSONL形式のみ):
        """
    )
    chain = prompt | llm | StrOutputParser()
    
    try:
        invoke_params = {
            "analysis_prompt": analysis_prompt,
            "text_data_jsonl": "\n".join(input_contexts_jsonl)
        }
        response_str = chain.invoke(invoke_params)
        logger.debug(f"AI Location Inference - Raw response received.")

        results = []
        
        match = re.search(r'```(?:jsonl|json)?\s*([\s\S]*?)\s*```', response_str, re.DOTALL)
        jsonl_content = match.group(1).strip() if match else response_str.strip()

        for line in jsonl_content.strip().split('\n'):
            cleaned_line = line.strip()
            if not cleaned_line: continue
            try:
                data = json.loads(cleaned_line)
                inferred_value = data.get("inferred_location")
                
                processed_value = ""
                if inferred_value and isinstance(inferred_value, str):
                    processed_value = inferred_value.strip()
                
                if processed_value.lower() in ["該当なし", "none", "null", "", "n/a"]:
                    processed_value = ""
                
                # (★) --- 推論結果もPython側で厳密に正規化・検証 ---
                if processed_value:
                    # 1. エイリアスマップで変換 (例: AIが "日光" と返した場合)
                    if processed_value in alias_map:
                        processed_value = alias_map[processed_value]
                    
                    # 2. 曖昧なキー (例: "広島") は破棄
                    elif processed_value in ambiguous_set:
                        processed_value = ""
                    
                    # 3. DBに存在する正式名称 (例: "廿日市市") か確認
                    elif processed_value in all_cities_wards:
                        pass # OK
                    
                    # 4. それ以外 (例: "アメリカ") は破棄
                    else:
                        if " " in processed_value and any(s in processed_value for s in ["市", "区"]):
                            pass # "札幌市 中央区" は OK
                        else:
                            processed_value = "" # 不明な地名として破棄
                # (★) --- 正規化ロジックここまで ---
                
                results.append({
                    "id": data.get("id"),
                    "inferred_location": processed_value
                })
                
            except (json.JSONDecodeError, AttributeError) as json_e:
                logger.warning(f"AI地名推論 回答パース失敗: {cleaned_line} - Error: {json_e}")
                id_match = re.search(r'"id":\s*(\d+)', cleaned_line)
                if id_match:
                    results.append({"id": int(id_match.group(1)), "inferred_location": ""})
                    
        return pd.DataFrame(results) if results else pd.DataFrame(columns=['id', 'inferred_location'])

    except Exception as e:
        logger.error(f"AI地名推論 バッチ処理中エラー: {e}", exc_info=True)
        st.error(f"AI地名推論処理エラー: {e}")
        return pd.DataFrame() # 失敗時は空のDF

# --- 7. (★) Step A: UI描画関数 ---

def update_progress_ui(
    progress_placeholder: st.delta_generator.DeltaGenerator,
    log_placeholder: st.delta_generator.DeltaGenerator,
    tip_placeholder: st.delta_generator.DeltaGenerator,  # (★) Tips用プレースホルダ
    processed_rows: int,
    total_rows: int,
    message_prefix: str
):
    """
    (Step A) の進捗バーとログエリアを更新する (DRY)
    (★) 要件: AI読み込み時間の進捗を0～100％で表示
    (★) 要件: AI Tipsをローテーション表示
    """
    try:
        # total_rows が 0 の場合 DivisionByZero を防ぐ
        if total_rows == 0:
            progress_percent = 1.0
        else:
            progress_percent = min(processed_rows / total_rows, 1.0)
            
        progress_text = f"[{message_prefix}] 処理中: {processed_rows}/{total_rows} 件 ({progress_percent:.0%})"
        progress_placeholder.progress(progress_percent, text=progress_text)

        # ログ表示 (最新50件)
        log_text_for_ui = "\n".join(st.session_state.log_messages[-50:])
        log_placeholder.text_area(
            "実行ログ (最新50件):",
            log_text_for_ui,
            height=200,
            key=f"log_update_{message_prefix}_{processed_rows}", # 重複キーを避ける
            disabled=True
        )
        
        # (★) --- AI Tips のローテーション表示 ---
        if 'tips_list' not in st.session_state or not st.session_state.tips_list:
             # 万が一Tipsが空の場合はAI Tips関数を呼び出す
             st.session_state.tips_list = get_analysis_tips_list_from_ai()
             st.session_state.current_tip_index = 0
             st.session_state.last_tip_time = time.time()

        now = time.time()
        # 60秒ごと（またはTIPSが1件しかない場合）にTIPSを更新
        if (now - st.session_state.last_tip_time > 60) or (len(st.session_state.tips_list) == 1):
            if len(st.session_state.tips_list) > 1:
                st.session_state.current_tip_index = (st.session_state.current_tip_index + 1) % len(st.session_state.tips_list)
            st.session_state.last_tip_time = now
        
        # (★) リストが空でないかチェック
        if st.session_state.tips_list:
            current_tip = st.session_state.tips_list[st.session_state.current_tip_index]
            tip_placeholder.info(f"💡 データ分析TIPS: {current_tip}")

    except Exception as e:
        # UIの更新エラーはログに警告のみ残し、処理は続行
        logger.warning(f"UI update failed: {e}")

def render_step_a():
    """(Step A) タグ付け処理のUIを描画する"""
    st.title("🏷️ Step A: AIタグ付け & キュレーション")

    # Step A 固有のセッションステートを初期化
    if 'cancel_analysis' not in st.session_state:
        st.session_state.cancel_analysis = False
    if 'generated_categories' not in st.session_state:
        st.session_state.generated_categories = {}
    if 'selected_categories' not in st.session_state:
        st.session_state.selected_categories = set()
    if 'analysis_prompt_A' not in st.session_state:
        st.session_state.analysis_prompt_A = ""
    if 'selected_text_col' not in st.session_state:
        st.session_state.selected_text_col = {}
    if 'tagged_df_A' not in st.session_state:
        st.session_state.tagged_df_A = pd.DataFrame()

    st.header("Step 1: 分析対象ファイルのアップロード")
    uploaded_files = st.file_uploader(
        "分析したい Excel / CSV ファイル（複数可）",
        type=['csv', 'xlsx', 'xls'],
        accept_multiple_files=True,
        key="uploader_A"
    )

    if not uploaded_files:
        st.info("分析を開始するには、ExcelまたはCSVファイルをアップロードしてください。")
        return

    # ファイル読み込み処理
    valid_files_data = {}
    error_messages = []
    for f in uploaded_files:
        df, err = read_file(f)
        if err:
            error_messages.append(f"**{f.name}**: {err}")
        else:
            valid_files_data[f.name] = df
            
    if error_messages:
        st.error("以下のファイルは読み込めませんでした:\n" + "\n".join(error_messages))
    if not valid_files_data:
        st.warning("読み込み可能なファイルがありません。")
        return

    st.header("Step 2: 分析指針の入力とカテゴリ生成")
    analysis_prompt = st.text_area(
        "AIがタグ付けとキュレーションを行う際の指針を入力してください（必須）:",
        value=st.session_state.analysis_prompt_A,
        height=100,
        placeholder="例: 広島県の観光に関するInstagramの投稿。無関係な地域の投稿や、単なる挨拶・宣伝は除外したい。\n例: ①農産品カテゴリ（牛乳,チーズ,米） ②農産品のイメージ（濃厚,新鮮）",
        key="analysis_prompt_input_A"
    )
    st.session_state.analysis_prompt_A = analysis_prompt
    
    st.markdown(f"（(★) 使用モデル: `{MODEL_FLASH_LITE}`）")
    if st.button("AIにカテゴリ候補を生成させる (Step 2)", key="gen_cat_button", type="primary"):
        if not analysis_prompt.strip():
            st.warning("分析指針は必須です。AIがデータを理解するために目的を入力してください。")
        elif not os.getenv("GOOGLE_API_KEY"):
            st.error("Google APIキーが設定されていません。（.envファイルを確認してください）")
        else:
            with st.spinner(f"AI ({MODEL_FLASH_LITE}) が分析指針を読み解き、カテゴリを考案中..."):
                logger.info("AIカテゴリ生成ボタンクリック")
                st.session_state.generated_categories = {"市区町村キーワード": "地名辞書(JAPAN_GEOGRAPHY_DB)から抽出された市区町村名"}
                
                ai_categories = get_dynamic_categories(analysis_prompt)
                
                if ai_categories:
                    st.session_state.generated_categories.update(ai_categories)
                    logger.info(f"AIカテゴリ生成成功: {list(ai_categories.keys())}")
                    st.success("AIによるカテゴリ候補の生成が完了しました。Step 3 に進んでください。")
                else:
                    st.error("AIによるカテゴリ生成に失敗しました。AIの応答を確認してください。")

    if not analysis_prompt.strip():
        st.warning("分析指針は必須です。AIがデータを理解するために目的を入力してください。")
        return

    st.header("Step 3: 分析カテゴリの選択")
    if not st.session_state.generated_categories:
        st.info("Step 2 で「AIにカテゴリ候補を生成させる」ボタンを押してください。")
        return
        
    st.markdown("タグ付けしたいカテゴリを以下から選択してください（「市区町村キーワード」は必須です）")
    
    selected_cats = []
    cols = st.columns(3)
    categories_to_show = st.session_state.generated_categories.items()
    
    for i, (cat, desc) in enumerate(categories_to_show):
        with cols[i % 3]:
            is_checked = st.checkbox(
                cat,
                value=(cat == "市区町村キーワード" or cat in st.session_state.selected_categories),
                help=str(desc), # (★ TypeError 修正)
                key=f"cat_cb_{cat}",
                disabled=(cat == "市区町村キーワード")
            )
            if is_checked:
                selected_cats.append(cat)
    st.session_state.selected_categories = set(selected_cats)

    st.header("Step 4: 分析対象テキスト列の指定")
    selected_text_col_map = {}
    st.markdown("ファイルごとに、タグ付け対象のテキストが含まれる列を指定してください。")
    for f_name, df in valid_files_data.items():
        cols_list = list(df.columns)
        default_index = 0
        
        if st.session_state.selected_text_col.get(f_name) in cols_list:
            default_index = cols_list.index(st.session_state.selected_text_col.get(f_name))
        elif any(c in cols_list for c in ['text', 'body', 'content', '投稿', '本文']):
            try:
                default_index = next(i for i, c in enumerate(cols_list) if c in ['text', 'body', 'content', '投稿', '本文'])
            except StopIteration:
                default_index = 0
                
        selected_col = st.selectbox(f"**{f_name}** のテキスト列:", cols_list, index=default_index, key=f"col_select_{f_name}")
        selected_text_col_map[f_name] = selected_col
    st.session_state.selected_text_col = selected_text_col_map

    st.header("Step 5: 分析実行")
    st.markdown(f"（(★) 使用モデル: `{MODEL_FLASH_LITE}`）")
    
    col_run, col_cancel = st.columns([1, 1])
    with col_cancel:
        if st.button("キャンセル", key="cancel_button_A", use_container_width=True):
            st.session_state.cancel_analysis = True
            logger.warning("分析キャンセルボタンが押されました。")
            st.warning("次のバッチ処理後に分析をキャンセルします...")
    
    with col_run:
        if st.button("分析実行 (Step 5)", type="primary", key="run_analysis_A", use_container_width=True):
            st.session_state.cancel_analysis = False
            st.session_state.log_messages = []
            st.session_state.tagged_df_A = pd.DataFrame()
            
            tip_placeholder = st.empty()
            try:
                with st.spinner("分析TIPSをAIで生成中..."):
                    if 'tips_list' not in st.session_state or not st.session_state.tips_list:
                        st.session_state.tips_list = get_analysis_tips_list_from_ai()
                
                if not st.session_state.tips_list: # フォールバック
                    st.session_state.tips_list = ["データ分析TIPSの取得に失敗しました。"]

                st.session_state.current_tip_index = random.randint(0, len(st.session_state.tips_list) - 1)
                st.session_state.last_tip_time = time.time()
                tip_placeholder.info(f"💡 データ分析TIPS: {st.session_state.tips_list[st.session_state.current_tip_index]}")
            except Exception as e:
                logger.error(f"Tips初期化エラー: {e}")

            try:
                with st.spinner(f"Step A: AI分析処理中 ({MODEL_FLASH_LITE})..."):
                    logger.info("Step A 分析実行ボタンクリック")
                    progress_placeholder = st.progress(0.0, text="処理待機中...")
                    log_placeholder = st.empty()

                    # --- 1. ファイル結合 ---
                    update_progress_ui(progress_placeholder, log_placeholder, tip_placeholder, 0, 100, "ファイル結合")
                    temp_dfs = []
                    for f_name, df in valid_files_data.items():
                        col_name = selected_text_col_map[f_name]
                        temp_df = df.rename(columns={col_name: 'ANALYSIS_TEXT_COLUMN'})
                        temp_dfs.append(temp_df)
                    
                    master_df = pd.concat(temp_dfs, ignore_index=True, sort=False)
                    master_df['id'] = master_df.index
                    if master_df.empty:
                        raise Exception("分析対象のデータがありません。")

                    # --- 2. 重複削除 ---
                    initial_row_count = len(master_df)
                    master_df.drop_duplicates(subset=['ANALYSIS_TEXT_COLUMN'], keep='first', inplace=True)
                    deduped_row_count = len(master_df)
                    logger.info(f"重複削除 完了。 {initial_row_count}行 -> {deduped_row_count}行")

                    # --- 3. (★) AI関連性フィルタリング (キュレーション) ---
                    total_filter_rows = len(master_df)
                    total_filter_batches = (total_filter_rows + FILTER_BATCH_SIZE - 1) // FILTER_BATCH_SIZE
                    all_filtered_results = []
                    
                    for i in range(0, total_filter_rows, FILTER_BATCH_SIZE):
                        if st.session_state.cancel_analysis:
                            raise Exception("分析がキャンセルされました")
                        
                        batch_df = master_df.iloc[i:i + FILTER_BATCH_SIZE]
                        current_batch_num = (i // FILTER_BATCH_SIZE) + 1
                        
                        update_progress_ui(
                            progress_placeholder, log_placeholder, tip_placeholder,
                            min(i + FILTER_BATCH_SIZE, total_filter_rows), total_filter_rows,
                            f"AIキュレーション (バッチ {current_batch_num}/{total_filter_batches})"
                        )
                        
                        filtered_df = filter_relevant_data_by_ai(batch_df, analysis_prompt)
                        if filtered_df is not None and not filtered_df.empty:
                            all_filtered_results.append(filtered_df)
                        
                        time.sleep(FILTER_SLEEP_TIME)
                    
                    if not all_filtered_results:
                        raise Exception("AIフィルタリング処理に失敗しました。")

                    filter_results_df = pd.concat(all_filtered_results, ignore_index=True)
                    relevant_ids = filter_results_df[filter_results_df['relevant'] == True]['id']
                    filtered_master_df = master_df[master_df['id'].isin(relevant_ids)].copy()
                    filtered_row_count = len(filtered_master_df)
                    logger.info(f"AIフィルタリング 完了。 {deduped_row_count}行 -> {filtered_row_count}行")

                    if filtered_master_df.empty:
                        st.warning("AIキュレーションの結果、分析対象のデータが0件になりました。")
                        st.session_state.tagged_df_A = pd.DataFrame()
                        progress_placeholder.progress(1.0, text="処理完了 (対象データ0件)")
                        return

                    # --- 4. (★) AIタグ付け (Pass 1) ---
                    selected_category_definitions = {
                        cat: desc for cat, desc in st.session_state.generated_categories.items()
                        if cat in st.session_state.selected_categories
                    }
                    
                    master_df_for_tagging = filtered_master_df
                    total_rows = len(master_df_for_tagging)
                    all_tagged_results = []
                    total_batches = (total_rows + TAGGING_BATCH_SIZE - 1) // TAGGING_BATCH_SIZE
                    
                    for i in range(0, total_rows, TAGGING_BATCH_SIZE):
                        if st.session_state.cancel_analysis:
                            raise Exception("分析がキャンセルされました")
                        
                        batch_df = master_df_for_tagging.iloc[i:i + TAGGING_BATCH_SIZE]
                        current_batch_num = (i // TAGGING_BATCH_SIZE) + 1
                        
                        update_progress_ui(
                            progress_placeholder, log_placeholder, tip_placeholder,
                            min(i + TAGGING_BATCH_SIZE, total_rows), total_rows,
                            f"AIタグ付け[1/2] (バッチ {current_batch_num}/{total_batches})"
                        )

                        tagged_df = perform_ai_tagging(batch_df, selected_category_definitions, analysis_prompt)
                        if tagged_df is not None and not tagged_df.empty:
                            all_tagged_results.append(tagged_df)
                        
                        time.sleep(TAGGING_SLEEP_TIME)

                    if not all_tagged_results:
                        raise Exception("AIタグ付け処理(Pass 1)に失敗しました。")
                    
                    tagged_results_df = pd.concat(all_tagged_results, ignore_index=True)

                    # --- (★) [改善 A-2] 5. AI地名推論 (Pass 2) ---
                    
                    # (★) Pass 1 の結果を一時的にマージ
                    temp_merged_df = pd.merge(master_df_for_tagging, tagged_results_df, on='id', how='left')
                    
                    # (★) 「市区町村キーワード」が空の行を抽出
                    rows_needing_inference = temp_merged_df[
                        temp_merged_df['市区町村キーワード'].isnull() | (temp_merged_df['市区町村キーワード'] == '')
                    ]
                    
                    all_inferred_results = []
                    total_inference_rows = len(rows_needing_inference)
                    
                    if total_inference_rows > 0:
                        logger.info(f"AI地名推論(Pass 2) 開始。対象: {total_inference_rows}件")
                        # (★) 地名正規化マップを（キャッシュから）取得
                        norm_maps = get_location_normalization_maps(JAPAN_GEOGRAPHY_DB, analysis_prompt)
                        
                        total_inf_batches = (total_inference_rows + TAGGING_BATCH_SIZE - 1) // TAGGING_BATCH_SIZE
                        
                        for i in range(0, total_inference_rows, TAGGING_BATCH_SIZE):
                            if st.session_state.cancel_analysis:
                                raise Exception("分析がキャンセルされました")
                            
                            batch_df = rows_needing_inference.iloc[i:i + TAGGING_BATCH_SIZE]
                            current_batch_num = (i // TAGGING_BATCH_SIZE) + 1
                            
                            update_progress_ui(
                                progress_placeholder, log_placeholder, tip_placeholder,
                                min(i + TAGGING_BATCH_SIZE, total_inference_rows), total_inference_rows,
                                f"AI地名推論[2/2] (バッチ {current_batch_num}/{total_inf_batches})"
                            )

                            inferred_df = perform_ai_location_inference(batch_df, analysis_prompt, norm_maps)
                            if inferred_df is not None and not inferred_df.empty:
                                all_inferred_results.append(inferred_df)
                            
                            time.sleep(TAGGING_SLEEP_TIME)
                        
                        if all_inferred_results:
                            inferred_results_df = pd.concat(all_inferred_results, ignore_index=True)
                            
                            # (★) Pass 2 の結果を Pass 1 の結果にマージ
                            tagged_results_df = tagged_results_df.set_index('id')
                            inferred_results_df = inferred_results_df.set_index('id')
                            
                            # (★) inferred_location の値で、'市区町村キーワード' の null/空 を埋める
                            tagged_results_df['市区町村キーワード'].fillna(inferred_results_df['inferred_location'], inplace=True)
                            tagged_results_df.loc[tagged_results_df['市区町村キーワード'] == '', '市区町村キーワード'] = inferred_results_df['inferred_location']
                            
                            tagged_results_df = tagged_results_df.reset_index()
                            logger.info("AI地名推論(Pass 2)の結果をマージしました。")

                    # --- 6. 最終マージ ---
                    logger.info("全AIタグ付け結果結合...");
                    
                    logger.info("最終マージ処理開始...");
                    final_df = pd.merge(master_df_for_tagging, tagged_results_df, on='id', how='right')
                    
                    final_cols = list(master_df_for_tagging.columns) + [col for col in tagged_results_df.columns if col not in master_df_for_tagging.columns]
                    final_df = final_df[final_cols]

                    st.session_state.tagged_df_A = final_df
                    logger.info("Step A 分析処理 正常終了");
                    st.success("AIによる分析処理が完了しました。");
                    progress_placeholder.progress(1.0, text="処理完了")
                    
                    update_progress_ui(
                        progress_placeholder, log_placeholder, tip_placeholder, 
                        total_rows, total_rows, "処理完了"
                    )
                    
                    tip_placeholder.empty() # 処理完了後、Tipsを消す

            except Exception as e:
                logger.error(f"Step A 分析実行中にエラー: {e}", exc_info=True)
                st.error(f"分析実行中にエラーが発生しました: {e}")
                if 'progress_placeholder' in locals():
                    progress_placeholder.progress(1.0, text="エラーにより処理中断")
                if 'tip_placeholder' in locals():
                    tip_placeholder.empty() # エラー時もTipsを消す

    # (★) 要件④: エクスポートリンクを表示
    if not st.session_state.tagged_df_A.empty:
        st.header("Step 6: 分析結果の確認とエクスポート")
        st.dataframe(st.session_state.tagged_df_A.head(50))

        @st.cache_data
        def convert_df_to_csv(df: pd.DataFrame) -> bytes:
            """DataFrameをUTF-8-SIGエンコードのCSV (bytes) に変換する"""
            return df.to_csv(encoding="utf-8-sig", index=False).encode("utf-8-sig")

        csv_data = convert_df_to_csv(st.session_state.tagged_df_A)
        st.download_button(
            label="分析結果CSV (Curated_Data.csv) をダウンロード",
            data=csv_data,
            file_name="Curated_Data.csv",
            mime="text/csv",
        )
        st.info("このCSVファイルを、Step B でアップロードして分析を続けてください。")

import networkx as nx # (★) Step B (共起ネットワーク) で必要
from itertools import combinations # (★) Step B (共起ネットワーク) で必要

def find_col(df: pd.DataFrame, patterns: List[str]) -> Optional[str]:
    """DataFrameから、複数のパターンに最初に一致する列名(str)を1つ返す"""
    cols = df.columns
    for pattern in patterns:
        try:
            # 1. 完全一致 (大文字小文字無視)
            for col in cols:
                if col.lower() == pattern.lower():
                    return col
            # 2. 部分一致 (大文字小文字無視)
            for col in cols:
                if re.search(pattern, col, re.IGNORECASE):
                    return col
        except re.error:
            continue # (e.g. invalid regex pattern)
    return None

def find_cols(df: pd.DataFrame, patterns: List[str]) -> List[str]:
    """DataFrameから、複数のパターンに一致する列名(list)をすべて返す"""
    cols = df.columns
    found_cols = set()
    for pattern in patterns:
        try:
            for col in cols:
                if re.search(pattern, col, re.IGNORECASE):
                    found_cols.add(col)
        except re.error:
            continue
    return sorted(list(found_cols))

def find_engagement_cols(df: pd.DataFrame, patterns: List[str]) -> List[str]:
    """DataFrameから、パターンに一致する「数値」列名(list)をすべて返す"""
    numeric_cols = df.select_dtypes(include=np.number).columns
    found_cols = set()
    for pattern in patterns:
        try:
            for col in numeric_cols: # 数値列のみを検索
                if re.search(pattern, col, re.IGNORECASE):
                    found_cols.add(col)
        except re.error:
            continue
    return sorted(list(found_cols))

def suggest_analysis_techniques_py(df: pd.DataFrame) -> List[Dict[str, Any]]:
    """
    (Step B) データフレームを分析し、Pythonで実行可能な基本的な分析手法を提案する。
    (列名に依存しない、汎用的な分析手法を提案するよう修正)
    """
    suggestions = []
    if df is None or df.empty:
        logger.error("suggest_analysis_techniques_py: DFが空です。")
        return suggestions
        
    try:
        # --- 1. 柔軟な列名の特定 ---
        all_cols = list(df.columns)
        
        text_col = find_col(df, ['ANALYSIS_TEXT_COLUMN', 'text', 'content', '本文'])
        location_col = find_col(df, ['市区町村キーワード', 'location', 'city', '地域'])
        sentiment_col = find_col(df, ['sent', 'センチメント'])
        
        # 日付列の堅牢な検索
        date_col = None
        object_cols_for_date = df.select_dtypes(include='object').columns.tolist()
        date_patterns = ['date', 'time', '日付', '日時']
        for col in object_cols_for_date:
            if any(re.search(p, col, re.IGNORECASE) for p in date_patterns):
                 if df[col].isnull().all(): continue
                 try:
                     if pd.to_datetime(df[col].dropna().sample(n=min(5, df[col].count())), errors='coerce').notna().any():
                         date_col = col
                         break
                 except Exception:
                     pass
            if date_col:
                break
        
        engagement_cols = find_engagement_cols(df, ['eng', 'like', 'いいね', 'エンゲージメント'])
        
        # 汎用カテゴリ列 (flag_cols) の特定
        base_flag_cols = find_cols(df, ['key', 'keyword', 'キーワード', 'カテゴリ', 'topic', 'ハッシュタグ'])
        flag_cols = sorted(list(set([c for c in base_flag_cols if c is not None and c != location_col])))

        # その他のカテゴリ列
        other_categorical = [
            col for col in df.select_dtypes(include='object').columns
            if col not in flag_cols and col != text_col and col != date_col and col != location_col
        ]
        
        # 全てのカテゴリ列 (汎用 + 地域 + その他)
        all_categorical = flag_cols + ([location_col] if location_col else []) + other_categorical
        
        logger.info(f"提案分析(PY) - Text:{text_col}, Location:{location_col}")
        logger.info(f"提案分析(PY) - FlagCols(汎用カテゴリ):{flag_cols}")
        logger.info(f"提案分析(PY) - Engagement:{engagement_cols}, Sentiment:{sentiment_col}, Date:{date_col}")

        potential_suggestions = []

        # --- 2. 提案ロジック (汎用化版) ---

        # 1. 全体メトリクス
        overall_metric_cols = [c for c in [sentiment_col] + engagement_cols if c is not None]
        potential_suggestions.append({
            "priority": 1, "name": "全体のメトリクス",
            "description": "投稿数、エンゲージメント、センチメント傾向など、データセット全体の概要を計算します。",
            "reason": "データ全体の状況把握に必須です。",
            "suitable_cols": overall_metric_cols,
            "type": "python"
        })

        # 3. 単純集計（頻度分析）
        for col in flag_cols + ([location_col] if location_col else []):
            potential_suggestions.append({
                "priority": 1, 
                "name": f"単純集計: {col}",
                "description": f"「{col}」列の出現頻度（TOP50）を分析します。",
                "reason": f"カテゴリ列({col})の基本指標です。",
                "suitable_cols": [col],
                "type": "python"
            })

        # 2. クロス集計
        if len(all_categorical) >= 2:
            potential_suggestions.append({
                "priority": 2, "name": "クロス集計（カテゴリ間）",
                "description": "2つのカテゴリ列（例: '話題カテゴリ' vs '市区町村'）を選択し、その組み合わせを分析します。",
                "reason": f"複数カテゴリ列({len(all_categorical)}個)あり、関連性の発見に。",
                "suitable_cols": all_categorical, 
                "type": "python"
            })

        # 3. 時系列分析
        if date_col and all_categorical:
            potential_suggestions.append({
                "priority": 3, "name": "時系列キーワード分析",
                "description": f"特定のカテゴリ列（例: '話題カテゴリ'）の出現数が時間（{date_col}）とともにどう変化したか分析します。",
                "reason": f"カテゴリ列と日時列({date_col})あり。",
                "suitable_cols": {"datetime": [date_col], "keywords": all_categorical},
                "type": "python"
            })
            
        # 3. 共起ネットワーク
        if text_col:
            potential_suggestions.append({
                "priority": 3, "name": "共起ネットワーク",
                "description": "投稿テキスト内の単語の出現パターンを分析し、関連性の高い単語のネットワークを構築します。",
                "reason": "テキストデータから隠れたトピックや関連性を発見します。",
                "suitable_cols": [text_col],
                "type": "python"
            })
            
        # 4. テキストマイニング
        if text_col:
            potential_suggestions.append({
                "priority": 4, "name": "テキストマイニング（頻出単語）",
                "description": "原文テキストから頻出する単語を抽出し、どのような言葉が多く使われているか全体像を把握します。",
                "reason": "原文テキストがあり、タグ付け以外のインサイト発見に。",
                "suitable_cols": [text_col],
                "type": "python"
            })

        # 4. カテゴリ列の集計と深掘り (Python + AI)
        if flag_cols and text_col:
            potential_suggestions.append({
                "priority": 4, "name": "カテゴリ列の集計と深掘り",
                "description": "指定したカテゴリ列（例: '話題カテゴリ'）ごとに投稿数を集計し、AIが投稿内容のサマリを生成します。",
                "reason": "カテゴリごとの主要な話題を把握します。",
                "suitable_cols": {'category_cols': flag_cols, 'text_col': [text_col]},
                "type": "python"
            })

        # 4. カテゴリ別 数値列TOP5分析 (Python + AI)
        if flag_cols and text_col and engagement_cols:
            potential_suggestions.append({
                "priority": 4, "name": "カテゴリ別 数値列TOP5分析",
                "description": f"指定したカテゴリ列ごとに、指定した数値列（例: '{engagement_cols[0]}'）が高いTOP5投稿を抽出し、AIがその概要を生成します。",
                "reason": "カテゴリごとに「バズった」投稿の内容を把握します。",
                "suitable_cols": {'category_cols': flag_cols, 'text_col': [text_col], 'numeric_cols': engagement_cols},
                "type": "python"
            })
        
        # 5. A/B 比較分析
        if all_categorical and location_col:
             potential_suggestions.append({
                "priority": 5, "name": "A/B 比較分析",
                "description": "2つの異なる投稿グループ（例：カテゴリA vs B、またはエリアA vs B）を選択し、投稿数や人気観光地（市区町村）の順位変動を比較します。",
                "reason": "グループ間の傾向の違いを明確にし、戦略立案に役立てます。",
                "suitable_cols": {'category_cols': all_categorical, 'location_col': [location_col]},
                "type": "python"
            })

        suggestions = sorted(potential_suggestions, key=lambda x: x['priority'])
        
        final_suggestions = []
        seen_names = set()
        for s in suggestions:
             if s['name'] not in seen_names:
                 final_suggestions.append(s)
                 seen_names.add(s['name'])
                 
        logger.info(f"Pythonベース提案(ソート後): {[s['name'] for s in final_suggestions]}")
        return final_suggestions

    except Exception as e:
        logger.error(f"Python分析手法提案中にエラー: {e}", exc_info=True)
        st.warning(f"分析手法提案中にエラー: {e}")
    return suggestions

def suggest_analysis_techniques_ai(
    user_prompt: str,
    df: pd.DataFrame,
    existing_suggestions: List[Dict[str, Any]]
) -> List[Dict[str, Any]]:
    """
    (Step B) ユーザーの自由記述プロンプトに基づき、AIが追加の分析手法を提案する。
    """
    logger.info("AIプロンプトベースの分析提案 (Flash Lite) を開始...")
    
    llm = get_llm(model_name=MODEL_FLASH_LITE, temperature=0.1)
    if llm is None:
        logger.error("suggest_analysis_techniques_ai: LLM (Flash Lite) が利用できません。")
        return []

    try:
        col_info = []
        for col in df.columns:
            col_info.append(f"- {col} (型: {df[col].dtype}, 例: {df[col].dropna().iloc[0] if not df[col].dropna().empty else 'N/A'})")
        column_info_str = "\n".join(col_info[:15])
        
        existing_names = [s['name'] for s in existing_suggestions]
        
        # (Bug 1.2 / 差分問題) 重複タスクをAIに厳格に禁止する
        forbidden_tasks = [
            "全体のメトリクス", "単純集計", "市区町村別投稿数", "クロス集計", 
            "時系列キーワード分析", "共起ネットワーク", "テキストマイニング",
            "カテゴリ列の集計と深掘り", "カテゴリ別 数値列TOP5分析", 
            "A/B 比較分析", "センチメント", "Sentiment",
            # AIが生成しがちな重複タスク名も明示的に禁止
            "話題カテゴリ別 投稿数とサマリ", 
            "話題カテゴリ別投稿数とサマリ",
            "話題カテゴリ別 エンゲージメントTOP5と概要", 
            "話題カテゴリ別エンゲージメント上位投稿TOP5",
            "市区町村別投稿数集計",
            "全体のセンチメント分析"
        ]
        existing_names_str = ", ".join(list(set(existing_names + forbidden_tasks)))

        prompt = PromptTemplate.from_template(
            """
            あなたはデータ分析の専門家です。ユーザーの「分析指示」と「データ構造」を読み、実行可能な「分析タスク」をJSONリスト形式で提案してください。
            
            # データ構造 (利用可能な列名):
            {column_info}
            
            # 既に提案済みのタスク (これらは提案しないでください):
            {existing_tasks}
            
            # ユーザーの分析指示:
            {user_prompt}
            
            # 指示:
            1. 「ユーザーの分析指示」を解釈し、具体的な分析タスク（例：「広島市と観光地の相関分析」）に分解する。
            2. 【重要】「既に提案済みのタスク」リストにあるタスクや、それに酷似したタスク（例：「単純集計」や「カテゴリ別サマリ」など）は【絶対に】提案しないでください。
            3. 各タスクを以下のJSON形式で定義する。
            4. `name`はタスク名、`description`はAI（あなた自身）がこの後実行するタスクの具体的な指示（プロンプト）とする。
            5. `priority`は 5 固定、`type`は "ai" 固定とする。
            
            # 回答 (JSONリスト形式のみ):
            [
              {{
                "priority": 5,
                "name": "（ユーザー指示に基づくタスク名1）",
                "description": "（このタスクを実行するためのAIへの具体的な指示プロンプト1）",
                "reason": "ユーザー指示に基づく",
                "suitable_cols": [],
                "type": "ai"
              }}
            ]
            """
        )
        chain = prompt | llm | StrOutputParser()
        response_str = chain.invoke({
            "column_info": column_info_str,
            "user_prompt": user_prompt,
            "existing_tasks": existing_names_str
        })

        logger.info(f"AI追加提案(生): {response_str}")
        match = re.search(r'\[.*\]', response_str, re.DOTALL)
        if not match:
            logger.warning("AIがJSONリスト形式で応答しませんでした。")
            return []
            
        json_str = match.group(0)
        ai_suggestions = json.loads(json_str)
        
        for s in ai_suggestions:
            s['type'] = 'ai'
            if 'priority' not in s: s['priority'] = 5
            
        logger.info(f"AI追加提案(パース済): {len(ai_suggestions)}件")
        return ai_suggestions

    except Exception as e:
        logger.error(f"AI追加提案の生成中にエラー: {e}", exc_info=True)
        st.warning(f"AI追加提案の生成中にエラーが発生しました: {e}")
        return []

# --- 8.0. グラフ生成ヘルパー ---
def generate_graph_image(
    df: pd.DataFrame,
    plot_type: str,
    x_col: Optional[str] = None,
    y_col: Optional[str] = None,
    title: str = "分析グラフ"
) -> Optional[str]:
    """
    DataFrameからmatplotlibグラフを生成し、Base64エンコードされた画像文字列を返す。
    (グラフサイズを動的に変更)
    """
    logger.info(f"グラフ生成開始: {title} (タイプ: {plot_type})")
    if df is None or df.empty:
        logger.warning("グラフ生成スキップ: DataFrameが空です。")
        return None

    # プロットタイプに応じてFigureサイズを変更
    if plot_type == 'network':
        plt.figure(figsize=(12, 12)) # 共起ネットワーク: 正方形 (12x12)
    elif plot_type == 'timeseries':
        plt.figure(figsize=(15, 7)) # 時系列: 横長 (15x7)
    else:
        plt.figure(figsize=(10, 7)) # デフォルト (棒グラフなど)
    
    plt.rcParams['font.size'] = 12
    
    try:
        if plot_type == 'bar' and x_col and y_col:
            df_plot = df.nlargest(20, y_col).sort_values(by=y_col, ascending=True)
            if df_plot.empty:
                raise ValueError("グラフ描画対象のデータがありません。")
                
            bars = plt.barh(df_plot[x_col], df_plot[y_col], color='#7280C1')
            plt.xlabel('件数')
            plt.ylabel(x_col)
            plt.grid(axis='x', linestyle='--', alpha=0.6)
            
            for bar in bars:
                plt.text(
                    bar.get_width() + (df_plot[y_col].max() * 0.01),
                    bar.get_y() + bar.get_height() / 2,
                    f' {bar.get_width():.0f}',
                    va='center',
                    ha='left'
                )
        
        elif plot_type == 'timeseries' and x_col and y_col:
            try:
                df[x_col] = pd.to_datetime(df[x_col])
                df_pivot = df.pivot(index=x_col, columns='keyword', values=y_col).fillna(0)
                
                if len(df_pivot.columns) > 6:
                    top_5_keywords = df_pivot.sum().nlargest(5).index
                    df_pivot['その他'] = df_pivot.drop(columns=top_5_keywords).sum(axis=1)
                    df_pivot = df_pivot[list(top_5_keywords) + ['その他']]
                
                df_pivot.plot(kind='line', ax=plt.gca(), linewidth=2.5)
                plt.xlabel('日付')
                plt.ylabel('件数')
                plt.legend(title='キーワード', bbox_to_anchor=(1.05, 1), loc='upper left')
                plt.grid(axis='y', linestyle='--', alpha=0.6)
                
            except Exception as e:
                logger.error(f"時系列ピボット/プロットエラー: {e}")
                plt.plot(df[x_col], df[y_col])
                plt.xlabel(x_col)
                plt.ylabel(y_col)

        elif plot_type == 'network':
            df_plot = df.nlargest(100, 'weight')
            if df_plot.empty:
                raise ValueError("ネットワーク描画対象のデータがありません。")

            G = nx.from_pandas_edgelist(df_plot, 'source', 'target', ['weight'])
            
            try:
                partition = community.best_partition(G)
                num_communities = len(set(partition.values()))
                colors = plt.cm.get_cmap('tab20', num_communities)
                node_colors = [colors(partition.get(node)) for node in G.nodes()]
            except Exception:
                node_colors = '#7280C1'
            
            # 共起ネットワークのレイアウト調整
            k_val = 2.5 / math.sqrt(len(G.nodes())) # k値を調整 (ノードを広げる)
            pos = nx.spring_layout(G, k=max(k_val, 0.5), iterations=50, seed=42)
            
            node_sizes = []
            try:
                for node in G.nodes():
                    total_weight = sum(data['weight'] for _, _, data in G.edges(node, data=True))
                    node_sizes.append(total_weight * 20)
            except Exception:
                node_sizes = 500
            
            edge_weights = [d['weight'] / df_plot['weight'].max() * 8 for u, v, d in G.edges(data=True)]

            nx.draw_networkx_nodes(G, pos, node_size=node_sizes, node_color=node_colors, alpha=0.9)
            nx.draw_networkx_edges(G, pos, width=edge_weights, alpha=0.1, edge_color='grey') # alphaを調整
            nx.draw_networkx_labels(G, pos, font_size=10, font_family='IPAGothic')
            
            plt.axis('off')
        
        elif plot_type == 'wordcloud' and not df.empty:
            if 'word' not in df.columns or 'count' not in df.columns:
                 raise ValueError("ワードクラウドには 'word' と 'count' 列が必要です。")
            
            frequencies = df.set_index('word')['count'].to_dict()
            
            if not frequencies:
                 raise ValueError("ワードクラウド用の単語がありません。")

            wc = WordCloud(
                font_path=font_path,
                width=800,
                height=500,
                background_color='white',
                colormap='viridis',
                max_words=100
            ).generate_from_frequencies(frequencies)
            
            plt.imshow(wc, interpolation='bilinear')
            plt.axis('off')

        else:
            logger.warning(f"未対応のプロットタイプ: {plot_type}")
            return None

        plt.title(title, fontsize=16, pad=20)
        plt.tight_layout()

        buf = BytesIO()
        plt.savefig(buf, format='png', dpi=96)
        buf.seek(0)
        
        image_base64 = base64.b64encode(buf.getvalue()).decode('utf-8')
        logger.info(f"グラフ生成成功: {title}")
        return image_base64

    except Exception as e:
        logger.error(f"グラフ生成 ({title}) 失敗: {e}", exc_info=True)
        return None
    finally:
        plt.clf()
        plt.close('all')

# --- 8.1. (★) Step B: Python分析ヘルパー ---

def run_simple_count(df: pd.DataFrame, suggestion: Dict[str, Any]) -> Dict[str, Any]:
    """(Step B) 単純集計（頻度分析）を実行し、DataFrameとグラフ(Base64)を返す"""
    results = {"data": pd.DataFrame(), "image_base64": None, "summary": ""}
    
    flag_cols = suggestion.get('suitable_cols', [])
    if not flag_cols:
        msg = "集計対象の列が見つかりません。"
        logger.warning(f"run_simple_count: {msg}")
        results["summary"] = msg
        return results
    
    # UIで編集された列を取得 (フォールバックあり)
    col_to_analyze = suggestion.get('ui_selected_col', flag_cols[0])
    
    if col_to_analyze not in df.columns:
        msg = f"列 '{col_to_analyze}' がDFに存在しません。"
        logger.warning(f"run_simple_count: {msg}")
        results["summary"] = msg
        return results
        
    try:
        s = df[col_to_analyze].astype(str).str.split(', ').explode()
        s = s[s.str.strip().isin(['', 'nan', 'None', 'N/A', '該当なし']) == False]
        s = s.str.strip()
        
        if s.empty:
            msg = "集計対象のキーワードがありませんでした。"
            logger.info(f"run_simple_count: {msg}")
            results["summary"] = msg
            return results
            
        counts = s.value_counts().head(50)
        counts_df = counts.reset_index()
        counts_df.columns = [col_to_analyze, 'count']
        
        results["data"] = counts_df
        results["summary"] = f"'{col_to_analyze}' の単純集計（頻度分析）を実行。上位は {counts_df.iloc[0,0]} ({counts_df.iloc[0,1]}件), {counts_df.iloc[1,0]} ({counts_df.iloc[1,1]}件) でした。"
        
        results["image_base64"] = generate_graph_image(
            df=counts_df,
            plot_type='bar',
            x_col=col_to_analyze,
            y_col='count',
            title=f"「{col_to_analyze}」 頻出TOP20"
        )
        return results
            
    except Exception as e:
        logger.error(f"run_simple_count error: {e}", exc_info=True)
        results["summary"] = f"エラー: {e}"
    return results

def run_crosstab(df: pd.DataFrame, suggestion: Dict[str, Any]) -> Dict[str, Any]:
    """(Step B) クロス集計を実行し、DataFrameを返す"""
    results = {"data": pd.DataFrame(), "image_base64": None, "summary": ""}
    
    cols = suggestion.get('suitable_cols', [])
    if len(cols) < 2:
        msg = "クロス集計には2列以上必要です。"
        logger.warning(f"run_crosstab: {msg}")
        results["summary"] = msg
        return results

    # UIで編集された列を取得 (フォールバックあり)
    col1 = suggestion.get('ui_selected_col1', cols[0])
    col2 = suggestion.get('ui_selected_col2', cols[1])

    if col1 not in df.columns or col2 not in df.columns:
        msg = f"選択された列 ({col1}, {col2}) がDFに存在しません。"
        logger.warning(f"run_crosstab: {msg}")
        results["summary"] = msg
        return results
    
    try:
        df_exploded_1 = df.assign(**{col1: df[col1].astype(str).str.split(', ')}).explode(col1)
        df_exploded_2 = df_exploded_1.assign(**{col2: df_exploded_1[col2].astype(str).str.split(', ')}).explode(col2)

        df_exploded_2[col1] = df_exploded_2[col1].str.strip()
        df_exploded_2[col2] = df_exploded_2[col2].str.strip()
        df_exploded_2 = df_exploded_2.replace('', np.nan).replace('nan', np.nan).replace('None', np.nan).dropna(subset=[col1, col2])
        
        crosstab_df = pd.crosstab(df_exploded_2[col1], df_exploded_2[col2])
        
        crosstab_long = crosstab_df.stack().reset_index()
        crosstab_long.columns = [col1, col2, 'count']
        crosstab_long = crosstab_long[crosstab_long['count'] > 0].sort_values(by='count', ascending=False)
        
        results["data"] = crosstab_long.head(100)
        results["summary"] = f"'{col1}' と '{col2}' のクロス集計を実行。最強の組み合わせは {crosstab_long.iloc[0,0]} x {crosstab_long.iloc[0,1]} ({crosstab_long.iloc[0,2]}件) でした。"
        logger.info("run_crosstab: グラフ生成はスキップされました。")
        
        return results
        
    except Exception as e:
        logger.error(f"run_crosstab error: {e}", exc_info=True)
        results["summary"] = f"エラー: {e}"
    return results

def run_timeseries(df: pd.DataFrame, suggestion: Dict[str, Any]) -> Dict[str, Any]:
    """(Step B) 時系列分析を実行し、DataFrameとグラフ(Base64)を返す"""
    results = {"data": pd.DataFrame(), "image_base64": None, "summary": ""}
    
    cols_dict = suggestion.get('suitable_cols', {})
    if not isinstance(cols_dict, dict) or 'datetime' not in cols_dict or 'keywords' not in cols_dict:
        msg = "列情報（datetime, keywords）が不十分です。"
        logger.warning(f"run_timeseries: {msg}")
        results["summary"] = msg
        return results
        
    # UIで編集された列を取得 (フォールバックあり)
    dt_col = suggestion.get('ui_selected_dt_col', cols_dict['datetime'][0])
    kw_col = suggestion.get('ui_selected_kw_col', cols_dict['keywords'][0])

    if dt_col not in df.columns:
        msg = f"日時列 '{dt_col}' が見つかりません。"
        logger.warning(f"run_timeseries: {msg}"); results["summary"] = msg; return results
    if kw_col not in df.columns:
        msg = f"キーワード列 '{kw_col}' が見つかりません。"
        logger.warning(f"run_timeseries: {msg}"); results["summary"] = msg; return results

    try:
        df_copy = df[[dt_col, kw_col]].copy()
        df_copy[dt_col] = pd.to_datetime(df_copy[dt_col], errors='coerce')
        df_copy = df_copy.dropna(subset=[dt_col])
        
        df_exploded = df_copy.assign(**{kw_col: df_copy[kw_col].astype(str).str.split(', ')}).explode(kw_col)
        df_exploded[kw_col] = df_exploded[kw_col].str.strip()
        df_exploded = df_exploded[df_exploded[kw_col].isin(['', 'nan', 'None', 'N/A', '該当なし']) == False]
        
        if df_exploded.empty:
            msg = "有効な日時/キーワードデータがありません。"
            logger.info(f"run_timeseries: {msg}"); results["summary"] = msg; return results

        time_df = df_exploded.groupby([pd.Grouper(key=dt_col, freq='D'), kw_col]).size().rename("count").reset_index()
        
        time_df.columns = ['date', 'keyword', 'count']
        
        top_keywords = df_exploded[kw_col].value_counts().head(50).index
        time_df_filtered = time_df[time_df['keyword'].isin(top_keywords)]
        
        time_df_for_graph = time_df_filtered.copy()
        
        time_df_for_json = time_df_filtered.sort_values(by=['keyword', 'date'])
        time_df_for_json['date'] = time_df_for_json['date'].dt.strftime('%Y-%m-%d')
        
        results["data"] = time_df_for_json
        
        results["image_base64"] = generate_graph_image(
            df=time_df_for_graph,
            plot_type='timeseries',
            x_col='date',
            y_col='count',
            title=f"「{kw_col}」別 時系列トレンド (TOP5)"
        )
        results["summary"] = f"'{dt_col}' と '{kw_col}' で時系列分析を実行。TOP5キーワードのトレンドグラフを生成しました。"
        return results
            
    except Exception as e:
        logger.error(f"run_timeseries error: {e}", exc_info=True)
        results["summary"] = f"エラー: {e}"
    return results

def run_text_mining(df: pd.DataFrame, suggestion: Dict[str, Any]) -> Dict[str, Any]:
    """(Step B) テキストマイニング（頻出単語）を実行し、DataFrameとグラフ(Base64)を返す"""
    results = {"data": pd.DataFrame(), "image_base64": None, "summary": ""}
    
    # UIで編集された列を取得 (フォールバックあり)
    text_col = suggestion.get('ui_selected_text_col', suggestion.get('suitable_cols', ['ANALYSIS_TEXT_COLUMN'])[0])
    
    if text_col not in df.columns or df[text_col].empty:
        msg = f"テキスト列 '{text_col}' がないか、空です。"
        logger.warning(f"run_text_mining: {msg}"); results["summary"] = msg; return results

    nlp = load_spacy_model()
    if nlp is None:
        st.error("spaCy日本語モデルのロードに失敗しました。")
        results["summary"] = "spaCy日本語モデルのロードに失敗しました。"
        return results
            
    try:
        texts = df[text_col].dropna().astype(str)
        if texts.empty:
            results["summary"] = "テキストデータが空です。"
            return results
            
        words = []
        target_pos = {'NOUN', 'PROPN', 'ADJ'}
        stop_words = {
            'の', 'に', 'は', 'を', 'が', 'で', 'て', 'です', 'ます', 'こと', 'もの', 'それ', 'あれ',
            'これ', 'ため', 'いる', 'する', 'ある', 'ない', 'いう', 'よう', 'そう', 'など', 'さん',
            '的', '人', '自分', '私', '僕', '何', 'その', 'この', 'あの'
        }
        
        custom_stop_words_str = suggestion.get('ui_custom_stop_words', '')
        if custom_stop_words_str:
            try:
                # カンマ、空白、改行、読点（、）で区切られた単語をセットに追加
                custom_set = set(
                    word.strip() for word in re.split(r'[\s,、\n]+', custom_stop_words_str) if word.strip()
                )
                stop_words.update(custom_set)
                logger.info(f"テキストマイニング: カスタム除外語 {len(custom_set)}件 を追加。")
            except Exception as e:
                logger.warning(f"カスタム除外語の解析失敗: {e}")

        total_texts = len(texts)
        if 'progress_text' not in st.session_state:
             st.session_state.progress_text = ""
        st.session_state.progress_text = "テキストマイニング (spaCy) 処理中... 0%"

        for i, doc in enumerate(nlp.pipe(texts, disable=["parser", "ner"], batch_size=50)):
            for token in doc:
                if (token.pos_ in target_pos) and (not token.is_stop) and (token.lemma_ not in stop_words) and (len(token.lemma_) > 1):
                    words.append(token.lemma_)
            
            if (i + 1) % 100 == 0:
                percent = (i + 1) / total_texts
                st.session_state.progress_text = f"テキストマイニング (spaCy) 処理中... {percent:.0%}"

        if not words:
            msg = "抽出可能な有効な単語が見つかりませんでした。"
            logger.warning(f"run_text_mining: {msg}"); results["summary"] = msg; return results

        word_counts = pd.Series(words).value_counts().head(100)
        word_counts_df = word_counts.reset_index()
        word_counts_df.columns = ['word', 'count']
        
        st.session_state.progress_text = "テキストマイニング (spaCy) 完了。"
        
        results["data"] = word_counts_df
        results["summary"] = f"'{text_col}' に対するテキストマイニングを実行。頻出単語は '{word_counts_df.iloc[0,0]}' ({word_counts_df.iloc[0,1]}件) でした。"
        
        # グラフ生成 (ワードクラウド)
        results["image_base64"] = generate_graph_image(
            df=word_counts_df,
            plot_type='wordcloud',
            title=f"「{text_col}」 頻出単語 ワードクラウド (TOP100)"
        )
        return results
        
    except Exception as e:
        logger.error(f"run_text_mining error: {e}", exc_info=True)
        results["summary"] = f"エラー: {e}"
    return results

def run_overall_metrics(df: pd.DataFrame, suggestion: Dict[str, Any]) -> Dict[str, Any]:
    """(Step B) データセット全体のメトリクスを計算する (単位追加)"""
    logger.info("run_overall_metrics 実行...")
    metrics = {}
    try:
        # (Bug 1.6) 単位を文字列として追加
        metrics["total_posts"] = f"{len(df):,}件"

        engagement_cols = [col for col in df.columns if any(c in col.lower() for c in ['いいね', 'like', 'エンゲージメント', 'engagement', 'retweet', 'リツイート'])]
        total_engagement = 0
        if engagement_cols:
            for col in engagement_cols:
                if pd.api.types.is_numeric_dtype(df[col]):
                    total_engagement += df[col].sum()
            metrics["total_engagement"] = f"{int(total_engagement):,}件"
        else:
            metrics["total_engagement"] = "N/A"

        sentiment_col = None
        if 'センチメント' in df.columns:
            sentiment_col = 'センチメント'
        elif find_col(df, ['sent', 'センチメント']):
            sentiment_col = find_col(df, ['sent', 'センチメント'])
            
        if sentiment_col:
            pos_count = int(df[df[sentiment_col].astype(str).str.contains('ポジティブ|Positive', case=False, na=False)].shape[0])
            neg_count = int(df[df[sentiment_col].astype(str).str.contains('ネガティブ|Negative', case=False, na=False)].shape[0])
            
            metrics["positive_posts"] = f"{pos_count:,}件"
            metrics["negative_posts"] = f"{neg_count:,}件"
            
            if (pos_count + neg_count) > 0:
                tendency = ((pos_count - neg_count) / (pos_count + neg_count)) * 100
                metrics["sentiment_tendency_percent"] = f"{int(np.floor(tendency))}%" # % を追加
            else:
                metrics["sentiment_tendency_percent"] = "0%"
        else:
            logger.warning("列 'センチメント' が見つかりませんでした。")
            metrics["positive_posts"] = "N/A"
            metrics["negative_posts"] = "N/A"
            metrics["sentiment_tendency_percent"] = "N/A"

        summary = f"全体のメトリクスを計算。総投稿数: {metrics['total_posts']}, 総エンゲージメント: {metrics['total_engagement']}。"
        
        return {"data": metrics, "image_base64": None, "summary": summary}

    except Exception as e:
        logger.error(f"run_overall_metrics error: {e}", exc_info=True)
        return {"data": {"error": str(e)}, "image_base64": None, "summary": f"エラー: {e}"}

def run_cooccurrence_network_pyvis(df: pd.DataFrame, suggestion: Dict[str, Any]) -> Dict[str, Any]:
    """
    (Step B) pyvis を使用し、詳細なパラメータに基づき共起ネットワークを構築する
    """
    logger.info("run_cooccurrence_network (pyvis版) 実行...")
    results = {"data": pd.DataFrame(), "image_base64": None, "html_content": None, "summary": "", "ai_legend": None, "communities": None}
    
    # --- 1. UIから渡されたパラメータの解析 ---
    try:
        # UI (Step 5) で設定されたパラメータを取得 (ui_... で始まるキー)
        flag_col = suggestion.get('ui_selected_flag_col')
        selected_keywords = suggestion.get('ui_selected_keywords')
        text_col = suggestion.get('ui_selected_text_col')
        
        # フォールバック (一括実行時など)
        if not flag_col:
            flag_col = find_col(df, ['市区町村キーワード', 'location', 'city', '地域']) or find_col(df, ['話題カテゴリ', 'topic', 'category'])
        if not text_col:
            text_col = find_col(df, ['ANALYSIS_TEXT_COLUMN', 'text', 'content', '本文'])
        if selected_keywords is None: # None と [] は区別
            try:
                s = df[flag_col].dropna().astype(str).str.split(',').explode().str.strip()
                s = s[~s.isin(['', 'nan', 'Nan', 'NaN'])]
                selected_keywords = s.value_counts().index.tolist()[:10] # デフォルトはTop10
            except Exception:
                selected_keywords = []
        
        solver = suggestion.get('solver', 'barnesHut')
        gravity = suggestion.get('gravity', -2000)
        node_distance = suggestion.get('node_distance', 200)
        spring_length = suggestion.get('spring_length', 250)
        top_n_words_limit = suggestion.get('top_n_words_limit', 100)
        max_degree_cutoff = suggestion.get('max_degree_cutoff', 50)
        min_occurrence = suggestion.get('min_occurrence', 10)
        default_node_size = suggestion.get('default_node_size', 15)
        default_text_size = suggestion.get('default_text_size', 50)
        run_ai_legend = suggestion.get('run_ai_legend', False)
        
    except Exception as e:
        msg = f"UIパラメータの解析エラー: {e}"
        logger.error(f"run_cooccurrence_network: {msg}", exc_info=True)
        results["summary"] = msg
        return results

    if not selected_keywords:
        msg = "絞り込みキーワードが選択されていません。"
        logger.warning(f"run_cooccurrence_network: {msg}"); results["summary"] = msg; return results
    if not flag_col or not text_col:
        msg = "対象列 (絞り込み列またはテキスト列) が見つかりません。"
        logger.warning(f"run_cooccurrence_network: {msg}"); results["summary"] = msg; return results
    if flag_col not in df.columns or text_col not in df.columns:
        msg = f"対象列 ({flag_col} または {text_col}) がDFに存在しません。"
        logger.warning(f"run_cooccurrence_network: {msg}"); results["summary"] = msg; return results

    nlp = load_spacy_model()
    if nlp is None:
        msg = "spaCy日本語モデルのロードに失敗しました。"
        logger.error(msg); results["summary"] = msg; return results

    try:
        # --- 2. spaCy 処理 ---
        target_pos = {'NOUN', 'PROPN', 'ADJ', 'VERB'}
        stop_words = {
            'の', 'に', 'は', 'を', 'が', 'で', 'て', 'です', 'ます', 'こと', 'もの', 'それ', 'あれ',
            'これ', 'ため', 'いる', 'する', 'ある', 'ない', 'いう', 'よう', 'そう', 'など', 'さん',
            '的', '人', '自分', '私', '僕', '何', 'その', 'この', 'あの', 'れる', 'られる',
            'てる', 'なる', '中', 'ところ', 'たち', '人達', '今回', '本当', 'とても', '色々'
        }
        
        custom_stop_words_str = suggestion.get('ui_custom_stop_words', '')
        if custom_stop_words_str:
            try:
                # カンマ、空白、改行、読点（、）で区切られた単語をセットに追加
                custom_set = set(
                    word.strip() for word in re.split(r'[\s,、\n]+', custom_stop_words_str) if word.strip()
                )
                stop_words.update(custom_set)
                logger.info(f"共起ネットワーク: カスタム除外語 {len(custom_set)}件 を追加。")
            except Exception as e:
                logger.warning(f"カスタム除外語の解析失敗: {e}")
        
        G = nx.Graph()
        
        # 1. キーワードでDataFrameをフィルタリング
        escaped_keywords = [re.escape(k) for k in selected_keywords]
        pattern = '|'.join(escaped_keywords)
        df_filtered = df[df[flag_col].astype(str).str.contains(pattern, na=False)]
        
        if df_filtered.empty:
            msg = "選択したキーワードを含む投稿が見つかりませんでした。"
            logger.warning(f"run_cooccurrence_network: {msg}"); results["summary"] = msg; return results

        texts_to_analyze = df_filtered[text_col].dropna().astype(str)
        
        # 2. Top N 単語リストの作成
        st.session_state.progress_text = "共起ネットワーク: (1/3) Top N 単語を計算中..."
        all_words = []
        for text in texts_to_analyze:
            doc = nlp(text)
            for token in doc:
                if (token.pos_ in target_pos) and (not token.is_stop) and (token.lemma_ not in stop_words) and (len(token.lemma_) > 1):
                    if token.lemma_ not in selected_keywords:
                        all_words.append(token.lemma_)
        
        if not all_words:
            msg = "フィルタ結果から分析対象の単語が見つかりませんでした。"
            logger.warning(f"run_cooccurrence_network: {msg}"); results["summary"] = msg; return results
            
        top_n_words_set = set(pd.Series(all_words).value_counts().head(top_n_words_limit).index)
        
        # 3. グラフ(G)の構築
        st.session_state.progress_text = "共起ネットワーク: (2/3) ネットワークを構築中..."
        for text in texts_to_analyze:
            doc = nlp(text)
            words_in_text = set()
            for token in doc:
                if (token.pos_ in target_pos) and (token.lemma_ in top_n_words_set):
                    words_in_text.add(token.lemma_)
            
            for word1, word2 in combinations(sorted(list(words_in_text)), 2):
                if G.has_edge(word1, word2):
                    G[word1][word2]['weight'] += 1
                else:
                    G.add_edge(word1, word2, weight=1)

        # 4. フィルタリング
        edges_to_remove = [(u, v) for u, v, data in G.edges(data=True) if data['weight'] < min_occurrence]
        G.remove_edges_from(edges_to_remove)
        G.remove_nodes_from(list(nx.isolates(G)))

        degrees = dict(G.degree())
        nodes_to_remove = [node for node, degree in degrees.items() if degree > max_degree_cutoff]
        G.remove_nodes_from(nodes_to_remove)
        G.remove_nodes_from(list(nx.isolates(G)))
        
        if G.number_of_nodes() == 0 or G.number_of_edges() == 0:
            msg = f"フィルタ条件 (最小共起: {min_occurrence}, 最大接続: {max_degree_cutoff}) により、表示可能なノードが0件になりました。"
            logger.warning(f"run_cooccurrence_network: {msg}"); results["summary"] = msg; return results

        # 5. pyvis グラフの生成
        st.session_state.progress_text = "共起ネットワーク: (3/3) グラフを描画中..."
        net = Network(height="700px", width="100%", cdn_resources='in_line')
        
        degrees = dict(G.degree())
        min_degree, max_degree = (min(degrees.values()) or 1), (max(degrees.values()) or 1)
        
        # 6. コミュニティ検出と色分け
        community_map = {}
        communities_with_words = {}
        ai_legend_map = {}
        try:
            communities_list = community.greedy_modularity_communities(G)
            communities_with_words = {i: list(comm) for i, comm in enumerate(communities_list)}
            community_map = {node: i for i, comm in communities_with_words.items() for node in comm}
            logger.info(f"コミュニティ検出成功。{len(communities_list)}個のクラスタを発見。")
        except Exception as e:
            logger.warning(f"コミュニティ検出に失敗: {e}。色分けなしで続行します。")
            
        results["communities"] = communities_with_words # 凡例表示用に格納

        for node in G.nodes():
            if node not in degrees: continue
            size_factor = degrees.get(node, 0)
            size = default_node_size + 30 * (size_factor - min_degree) / (max_degree - min_degree + 1e-6)
            group_id = community_map.get(node, 0)
            color = COLOR_PALETTE[group_id % len(COLOR_PALETTE)]
            
            net.add_node(
                node, label=node, size=size, title=f"{node} (クラスタ: {group_id}, 結合数: {size_factor})",
                color=color,
                font={"size": default_text_size}
            )
            
        for u, v, data in G.edges(data=True):
            weight = data['weight']
            net.add_edge(u, v, title=f"共起回数: {weight}", value=weight)

        if solver == 'barnesHut':
            net.barnes_hut(gravity=gravity, overlap=0.1)
        else:
            net.repulsion(node_distance=node_distance, spring_length=spring_length)
        net.solver = solver
        net.show_buttons(filter_=['physics', 'nodes', 'layout'])
        
        # 7. HTMLコンテンツを生成して返す
        html_file = "cooccurrence_network.html"
        net.save_graph(html_file)
        with open(html_file, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        results["html_content"] = html_content
        
        edge_list = pd.DataFrame(G.edges(data=True), columns=["source", "target", "data"])
        edge_list['weight'] = edge_list['data'].apply(lambda x: x['weight'])
        results["data"] = edge_list[['source', 'target', 'weight']].sort_values(by="weight", ascending=False)
        results["summary"] = f"'{flag_col}' ( {', '.join(selected_keywords[:3])}...) で絞り込み、共起ネットワーク (pyvis) を生成。{G.number_of_nodes()}ノード, {G.number_of_edges()}エッジ。"

        # 8. AI凡例生成
        if run_ai_legend and communities_with_words:
            st.session_state.progress_text = "共起ネットワーク: (AI) 凡例を生成中..."
            llm = get_llm(model_name=MODEL_FLASH_LITE, temperature=0.1)
            if llm:
                prompt = PromptTemplate.from_template(
                    "以下の「単語リスト」の共通テーマを【3語以内】で考案してください。\n"
                    "# 単語リスト (上位10件): {word_list_str}\n"
                    "# 回答 (3語以内):"
                )
                chain = prompt | llm | StrOutputParser()
                
                for group_id, words in communities_with_words.items():
                    if not words: continue
                    words_top10 = sorted(words, key=lambda w: degrees.get(w, 0), reverse=True)[:10]
                    words_str = ", ".join(words_top10)
                    try:
                        raw_label = chain.invoke({"word_list_str": words_str})
                        cleaned_label = re.sub(r'^(#|回答)\s*\(.*?\)\s*:\s*', '', raw_label.strip())
                        ai_legend_map[group_id] = cleaned_label
                        time.sleep(1.0) # Rate Limit
                    except Exception as e:
                        logger.error(f"AI凡例生成エラー (Group {group_id}): {e}")
                        ai_legend_map[group_id] = "(AIエラー)"
                results["ai_legend"] = ai_legend_map
                results["summary"] += " AIによる凡例生成完了。"
            else:
                results["summary"] += " AI凡例生成スキップ (LLMロード失敗)。"
        
        st.session_state.progress_text = "共起ネットワーク 完了。"
        return results

    except Exception as e:
        logger.error(f"run_cooccurrence_network (pyvis版) エラー: {e}", exc_info=True)
        results["summary"] = f"共起ネットワーク分析中にエラー: {e}"
        st.session_state.progress_text = f"共起ネットワーク エラー: {e}"
        return results

def run_generic_category_summary(df: pd.DataFrame, suggestion: Dict[str, Any]) -> Dict[str, Any]:
    """
    (★) 汎用: カテゴリ列ごとに投稿数、サマリ(AI)、上位キーワードを分析する
    """
    logger.info("run_generic_category_summary 実行...")
    results = {"data": pd.DataFrame(), "image_base64": None, "summary": ""}
    
    # 1. UI (Step 5) から渡された「分析軸となるカテゴリ列」を取得
    default_topic_col = find_col(df, ['話題カテゴリ', 'topic', 'category'])
    topic_col = suggestion.get('ui_selected_category_col', default_topic_col)
    text_col = find_col(df, ['ANALYSIS_TEXT_COLUMN', 'text', 'content', '本文'])

    if not topic_col or not text_col:
        msg = f"分析に必要な列 (カテゴリ列またはテキスト列) が見つかりません。"
        logger.warning(f"run_generic_category_summary: {msg}")
        return {"data": pd.DataFrame([{"error": msg}]), "image_base64": None, "summary": msg}
    if topic_col not in df.columns or text_col not in df.columns:
        msg = f"指定された列 ('{topic_col}', '{text_col}') がDFに存在しません。"
        logger.warning(f"run_generic_category_summary: {msg}")
        return {"data": pd.DataFrame([{"error": msg}]), "image_base64": None, "summary": msg}
    
    # 2. (Enhancement 2.4) 上位キーワードの候補列を動的に決定
    flag_cols = [col for col in df.columns if col.endswith('キーワード')]
    location_col = find_col(df, ['市区町村キーワード', 'location', 'city', '地域'])
    # location_col と topic_col 自身を除外
    cols_to_use_for_keywords = [col for col in flag_cols if col != location_col and col != topic_col]
    logger.info(f"TOPキーワード集計対象 (地域/トピック除外): {cols_to_use_for_keywords}")
    
    # 3. (Enhancement 2.1) ハードコードせず、列のユニーク値上位10件を対象
    try:
        s = df[topic_col].astype(str).str.split(', ').explode()
        s = s[s.str.strip().isin(['', 'nan', 'None', 'N/A', '該当なし']) == False]
        s = s.str.strip()
        if s.empty:
            raise ValueError(f"カテゴリ列 '{topic_col}' に有効なデータがありません。")
        
        target_categories = s.value_counts().head(10).index.tolist()
        logger.info(f"'{topic_col}' の上位10カテゴリを分析対象とします: {target_categories}")
        
    except Exception as e:
        msg = f"カテゴリ列 '{topic_col}' の値の取得に失敗: {e}"
        logger.error(msg, exc_info=True)
        return {"data": pd.DataFrame([{"error": msg}]), "image_base64": None, "summary": msg}

    results_list = []
    
    total_cats = len(target_categories)
    if 'progress_text' not in st.session_state:
            st.session_state.progress_text = ""
            
    for i, category in enumerate(target_categories):
        # (Bug 1.3) サブ進捗を更新
        st.session_state.progress_text = f"カテゴリ深掘り ({i+1}/{total_cats}): {category}"
        
        df_filtered = df[df[topic_col].astype(str).str.contains(re.escape(category), na=False)]
        post_count = len(df_filtered)
        
        if post_count == 0:
            results_list.append({
                "category": category,
                "post_count": 0,
                "summary_ai": "N/A (投稿なし)",
                "top_keywords": []
            })
            continue
        
        ai_suggestion = {
            "name": f"Summary for {category}",
            "description": f"「{category}」カテゴリに関する以下の投稿サンプルを読み、主要な話題を1～2文で要約してください。"
        }
        summary_ai = run_ai_summary_batch(df_filtered, ai_suggestion)
        
        # 4. (Enhancement 2.4) 上位キーワード (Python)
        top_keywords = []
        if cols_to_use_for_keywords and not df_filtered.empty:
            all_keywords_series = []
            for kw_col in cols_to_use_for_keywords:
                if kw_col in df_filtered.columns:
                    s_kw = df_filtered[kw_col].astype(str).str.split(', ').explode()
                    s_kw = s_kw[s_kw.str.strip().isin(['', 'nan', 'None', 'N/A', '該当なし']) == False]
                    s_kw = s_kw.str.strip()
                    if not s_kw.empty:
                        all_keywords_series.append(s_kw)
            if all_keywords_series:
                combined_s = pd.concat(all_keywords_series)
                top_keywords = combined_s.value_counts().head(5).index.tolist()
        
        results_list.append({
            "category": category,
            "post_count": post_count,
            "summary_ai": summary_ai,
            "top_keywords": top_keywords
        })
        
        time.sleep(max(TAGGING_SLEEP_TIME / 2, 1.0))

    st.session_state.progress_text = "カテゴリ深掘り 完了。"
    results_df = pd.DataFrame(results_list)
    
    image_base64 = generate_graph_image(
        df=results_df,
        plot_type='bar',
        x_col='category',
        y_col='post_count',
        title=f"「{topic_col}」別 投稿数 (Top 10)"
    )
    
    summary = f"「{topic_col}」別の分析を実行。投稿数グラフを生成しました。"
    return {"data": results_df, "image_base64": image_base64, "summary": summary}

# --- 
# [修正版] app.py の L1628-L1715 (run_generic_engagement_top5)
# (str.contains() バグを修正)
# ---
def run_generic_engagement_top5(df: pd.DataFrame, suggestion: Dict[str, Any]) -> Dict[str, Any]:
    """
    (★) 汎用: カテゴリ列別に数値列TOP5投稿と概要(AI)を分析する
    (★) 修正: str.contains() バグを修正し、explode() ベースの集計に変更
    """
    logger.info("run_generic_engagement_top5 実行...")
    results = {"data": pd.DataFrame(), "image_base64": None, "summary": ""}

    # 1. UI (Step 5) から渡された「分析軸となる列」を取得
    default_topic_col = find_col(df, ['話題カテゴリ', 'topic', 'category'])
    default_text_col = find_col(df, ['ANALYSIS_TEXT_COLUMN', 'text', 'content', '本文'])
    default_eng_col = find_engagement_cols(df, ['eng', 'like', 'いいね', 'エンゲージメント'])
    default_eng_col = default_eng_col[0] if default_eng_col else None

    topic_col = suggestion.get('ui_selected_category_col', default_topic_col)
    text_col = suggestion.get('ui_selected_text_col', default_text_col)
    engagement_col = suggestion.get('ui_selected_numeric_col', default_eng_col)
    
    if not topic_col or not text_col or not engagement_col:
        msg = f"分析に必要な列 (カテゴリ列, テキスト列, 数値列) が見つかりません。"
        logger.warning(f"run_generic_engagement_top5: {msg}")
        return {"data": pd.DataFrame([{"error": msg}]), "image_base64": None, "summary": msg}
    if topic_col not in df.columns:
        msg = f"カテゴリ列 '{topic_col}' が見つかりません。"
        return {"data": pd.DataFrame([{"error": msg}]), "image_base64": None, "summary": msg}
    if engagement_col not in df.columns or not pd.api.types.is_numeric_dtype(df[engagement_col]):
        msg = f"数値列 '{engagement_col}' が数値列として存在しません。"
        return {"data": pd.DataFrame([{"error": msg}]), "image_base64": None, "summary": msg}
    if text_col not in df.columns:
        msg = f"テキスト列 '{text_col}' が見つかりません。"
        return {"data": pd.DataFrame([{"error": msg}]), "image_base64": None, "summary": msg}

    # 2. (★) --- 修正: explode ベースの集計ロジック ---
    try:
        # (★) 1. 元のDFを explode する (カンマ区切りを堅牢に処理)
        df_exploded = df.assign(**{topic_col: df[topic_col].astype(str).str.split(',')}).explode(topic_col)
        df_exploded[topic_col] = df_exploded[topic_col].str.strip()

        # (★) 2. 空白・N/A等を除外
        s = df_exploded[topic_col]
        s = s[s.str.strip().isin(['', 'nan', 'None', 'N/A', '該当なし']) == False]
        
        if s.empty:
            raise ValueError(f"カテゴリ列 '{topic_col}' に有効なデータがありません。")
            
        # (★) 3. 上位10カテゴリを決定 (これが正しい母数)
        target_categories = s.value_counts().head(10).index.tolist()
        logger.info(f"'{topic_col}' の上位10カテゴリを分析対象とします: {target_categories}")
    except Exception as e:
        msg = f"カテゴリ列 '{topic_col}' の値の取得に失敗: {e}"
        logger.error(msg, exc_info=True)
        return {"data": pd.DataFrame([{"error": msg}]), "image_base64": None, "summary": msg}

    # 3. (Enhancement 2.3) メディアリンク列を特定
    link_col_candidates = ['link', 'url', 'media_url', '投稿URL', 'URL', 'Link', 'Url']
    found_link_col = find_col(df, link_col_candidates)
    if found_link_col:
        logger.info(f"メディアリンク列: '{found_link_col}' を使用します。")
    else:
        logger.warning(f"メディアリンク列 ({link_col_candidates}) が見つかりませんでした。")

    results_list = []
    
    total_cats = len(target_categories)
    if 'progress_text' not in st.session_state:
            st.session_state.progress_text = ""

    for i, category in enumerate(target_categories):
        st.session_state.progress_text = f"数値列TOP5分析中 ({i+1}/{total_cats}): {category}"
        
        df_filtered = df_exploded[df_exploded[topic_col] == category]
        post_count = len(df_filtered)
        
        if post_count == 0:
            continue
            
        # (★) フィルタ済みDFから nlargest を取得 (これは元のロジックでOK)
        df_top5 = df_filtered.nlargest(5, engagement_col, keep='first')
        top5_posts_data = []
        
        if df_top5.empty:
                results_list.append({
                "category": category,
                "post_count": post_count,
                "top_posts": []
            })
                continue

        # (★) --- [品質向上案 B-2] AI呼び出しをループの外に出す ---
        top_5_texts_list = df_top5[text_col].astype(str).tolist()
        combined_texts_for_ai = "\n---\n".join([f"投稿{idx+1}: {text[:300]}..." for idx, text in enumerate(top_5_texts_list)])
        
        ai_suggestion_combined = {
            "name": f"Summary for Top 5 {category}",
            "description": f"以下の「{category}」カテゴリで「{engagement_col}」が多かった投稿（{len(top_5_texts_list)}件）のサンプルです。これらの投稿に共通する「人気の理由」や「傾向」を1〜2文で要約してください。\n\n# 投稿サンプル:\n{combined_texts_for_ai}"
        }
        common_summary_ai = run_ai_summary_batch(df_filtered, ai_suggestion_combined)
        time.sleep(max(TAGGING_SLEEP_TIME / 2, 1.0)) # 1カテゴリ1コール後のスリープ

        for _, row in df_top5.iterrows():
            post_text = str(row[text_col])
            engagement_value = row[engagement_col]
            
            # (★) 1回だけ呼び出したAIの共通サマリを使用
            summary_ai_for_post = common_summary_ai 
            
            # 4. (Enhancement 2.3) メディアリンクを取得
            link_value = None
            if found_link_col and found_link_col in row and pd.notna(row[found_link_col]):
                link_value = str(row[found_link_col])
            
            top5_posts_data.append({
                "engagement": int(engagement_value),
                "summary_ai": summary_ai_for_post, 
                "original_text_snippet": post_text[:100],
                "media_link": link_value
            })
            
            # (★) ループ内のAIコールと time.sleep を削除

        results_list.append({
            "category": category,
            "post_count": post_count,
            "top_posts": top5_posts_data
        })

    st.session_state.progress_text = "数値列TOP5分析 完了。"
    results_df = pd.DataFrame(results_list)
    
    summary = f"「{topic_col}」別の高「{engagement_col}」投稿TOP5を抽出しました。"
    return {"data": results_df, "image_base64": None, "summary": summary}

# (New Feature 3.1) A/B比較関数
def run_ab_comparison(df: pd.DataFrame, suggestion: Dict[str, Any]) -> Dict[str, Any]:
    """
    (Step B) 2つのグループ(A, B)の投稿数と人気観光地ランキングを比較する
    """
    logger.info("run_ab_comparison 実行...")
    results = {"data": {}, "image_base64": None, "summary": ""}
    
    try:
        # 1. UIからパラメータを取得
        ab_params = suggestion.get('ui_ab_params', {})
        a_col = ab_params.get('a_col')
        a_val = ab_params.get('a_val')
        b_col = ab_params.get('b_col')
        b_val = ab_params.get('b_val')
        
        # 汎用的に列を見つける
        location_col = find_col(df, ['市区町村キーワード', 'location', 'city', '地域'])
        topic_col = find_col(df, ['話題カテゴリ', 'topic', 'category'])

        if not all([a_col, a_val, b_col, b_val, location_col, topic_col]):
            msg = f"A/B比較のパラメータが不足しています (A/B列/値、地域列、トピック列が必要です)"
            logger.warning(f"run_ab_comparison: {msg}")
            return {"data": {"error": msg}, "image_base64": None, "summary": msg}

        # 2. グループA, BのDataFrameを作成
        df_A = df[df[a_col].astype(str).str.contains(re.escape(a_val), na=False)]
        df_B = df[df[b_col].astype(str).str.contains(re.escape(b_val), na=False)]
        
        if df_A.empty or df_B.empty:
             msg = f"グループA ({a_val}: {len(df_A)}件) または グループB ({b_val}: {len(df_B)}件) のデータが0件です。"
             logger.warning(f"run_ab_comparison: {msg}")
             return {"data": {"error": msg}, "image_base64": None, "summary": msg}

        # 3. カテゴリ別投稿数 比較
        cats_A = df_A[topic_col].value_counts().rename(f"Count (A: {a_val})")
        cats_B = df_B[topic_col].value_counts().rename(f"Count (B: {b_val})")
        
        df_cat_compare = pd.concat([cats_A, cats_B], axis=1).fillna(0).astype(int)
        df_cat_compare['Total'] = df_cat_compare.sum(axis=1)
        df_cat_compare.sort_values(by='Total', ascending=False, inplace=True)
        sum_A = df_cat_compare[cats_A.name].sum()
        sum_B = df_cat_compare[cats_B.name].sum()
        df_cat_compare[f"Share (A: {a_val})"] = (df_cat_compare[cats_A.name] / sum_A).map('{:.1%}'.format) if sum_A > 0 else 0
        df_cat_compare[f"Share (B: {b_val})"] = (df_cat_compare[cats_B.name] / sum_B).map('{:.1%}'.format) if sum_B > 0 else 0

        # 4. 観光地別(地域) 順位変動 比較
        locs_A = df_A[location_col].value_counts().rename(f"Count (A: {a_val})")
        locs_B = df_B[location_col].value_counts().rename(f"Count (B: {b_val})")
        
        df_rank_compare = pd.concat([locs_A, locs_B], axis=1).fillna(0).astype(int)
        df_rank_compare[f"Rank (A: {a_val})"] = df_rank_compare[locs_A.name].rank(ascending=False, method='min').astype(int)
        df_rank_compare[f"Rank (B: {b_val})"] = df_rank_compare[locs_B.name].rank(ascending=False, method='min').astype(int)
        
        df_rank_compare['Rank Change (A vs B)'] = (df_rank_compare[f"Rank (B: {b_val})"] - df_rank_compare[f"Rank (A: {a_val})"]).astype(int)
        
        df_rank_compare.sort_values(by=f"Count (B: {b_val})", ascending=False, inplace=True)
        df_rank_compare = df_rank_compare[[
            f"Rank (A: {a_val})", f"Count (A: {a_val})", 
            f"Rank (B: {b_val})", f"Count (B: {b_val})", 
            'Rank Change (A vs B)'
        ]]
        
        summary = f"A/B比較: 「{a_val}」 (A: {len(df_A)}件) vs 「{b_val}」 (B: {len(df_B)}件) を実行。"
        
        results["data"] = {
            "category_comparison": df_cat_compare.reset_index().rename(columns={'index': topic_col}).to_dict(orient='records'),
            "ranking_comparison": df_rank_compare.reset_index().rename(columns={'index': location_col}).head(20).to_dict(orient='records')
        }
        results["summary"] = summary
        
        return results
        
    except Exception as e:
        logger.error(f"run_ab_comparison error: {e}", exc_info=True)
        return {"data": {"error": f"A/B比較エラー: {e}"}, "image_base64": None, "summary": f"A/B比較エラー: {e}"}

# --- 8.2. (★) Step B: AI分析ヘルパー (Bug 1.1 修正) ---
def run_ai_summary_batch(df: pd.DataFrame, suggestion: Dict[str, Any]) -> str:
    """
    (Step B) AI (Flash Lite) を使用して、指定されたタスク(description)を実行する。
    (★) 改善: プロンプトを「要約」から「考察」に変更
    """
    logger.info(f"run_ai_summary_batch 実行 (タスク: {suggestion.get('name', 'N/A')})...")
    
    llm = get_llm(model_name=MODEL_FLASH_LITE, temperature=0.1, timeout_seconds=120)
    if llm is None:
        logger.error("run_ai_summary_batch: LLM (Flash Lite) が利用できません。")
        return "AIモデル(Flash Lite)が利用できませんでした。"

    try:
        ai_prompt_instruction = suggestion.get('description', 'データからインサイトを抽出してください。')
        
        total_rows_count = len(df)
        
        text_col = find_col(df, ['ANALYSIS_TEXT_COLUMN', 'text', 'content', '本文'])
        
        if text_col:
            # テキスト列がある場合、サンプルを抽出
            sample_size = min(50, total_rows_count)
            if sample_size > 0:
                text_samples = df[text_col].dropna().sample(n=sample_size, random_state=1).tolist()
                data_context = "\n".join([f"- {text[:200]}..." for text in text_samples])
            else:
                data_context = "（サンプルデータなし）"
        else:
            # テキスト列がない場合、DFの先頭をJSONで渡す
            data_context = df.head(10).to_json(orient='records', force_ascii=False)

        prompt = PromptTemplate.from_template(
            """
            あなたはデータアナリストです。「指示」と「データサンプル」に基づき、
            単なる要約ではなく、データから読み取れる【インサイト（発見）】や【傾向の背景（仮説）】を
            簡潔に考察してください。

            # 指示 (Task):
            {ai_instruction}

            # データサンプル (Data Sample):
            (分析対象: 全 {total_rows} 件からの抜粋)
            {data_context}

            # 考察のポイント:
            - データが示している「最も重要な事実」は何か？
            - なぜその傾向が起きているのか（背景・原因の仮説）？
            - （もしあれば）データから読み取れる「次のアクションのヒント」は何か？

            # 回答 (分析結果の考察のみをMarkdown形式で):
            """
        )

        chain = prompt | llm | StrOutputParser()
        
        response_str = chain.invoke({
            "ai_instruction": ai_prompt_instruction,
            "data_context": data_context,
            "total_rows": total_rows_count
        })
        
        return response_str.strip()

    except Exception as e:
        logger.error(f"run_ai_summary_batch 実行エラー: {e}", exc_info=True)
        return f"AI分析タスクの実行中にエラーが発生しました: {e}"

# --- 8.3. (★) Step B: 分析実行ルーター (汎用化対応) ---
def execute_analysis(
    analysis_name: str,
    df: pd.DataFrame,
    suggestion: Dict[str, Any]
) -> Dict[str, Any]:
    """
    (Step B) 分析名に基づき、適切なPythonまたはAIの実行関数を呼び出すルーター
    """
    try:
        analysis_type = suggestion.get('type', 'python')
        
        if analysis_type == 'python':
            if analysis_name == "全体のメトリクス":
                return run_overall_metrics(df, suggestion)
            elif analysis_name.startswith("単純集計:"):
                return run_simple_count(df, suggestion)
            elif analysis_name.startswith("クロス集計"):
                return run_crosstab(df, suggestion)
            elif analysis_name == "時系列キーワード分析":
                return run_timeseries(df, suggestion)
            elif analysis_name == "テキストマイニング（頻出単語）":
                return run_text_mining(df, suggestion)
            elif analysis_name == "共起ネットワーク":
                return run_cooccurrence_network_pyvis(df, suggestion)
            # 汎用タスク名に対応
            elif analysis_name == "カテゴリ列の集計と深掘り":
                return run_generic_category_summary(df, suggestion)
            elif analysis_name == "カテゴリ別 数値列TOP5分析":
                return run_generic_engagement_top5(df, suggestion)
            # A/B比較のルーティング
            elif analysis_name == "A/B 比較分析":
                return run_ab_comparison(df, suggestion)
            else:
                logger.warning(f"Python分析 '{analysis_name}' の実行ロジックが定義されていません。AI分析にフォールバックします。")
                suggestion['description'] = f"データサンプルを使い、'{analysis_name}' を実行してください。"
                ai_result_str = run_ai_summary_batch(df, suggestion)
                return {"data": ai_result_str, "image_base64": None, "summary": ai_result_str[:100] + "..."}
        
        elif analysis_type == 'ai':
            ai_result_str = run_ai_summary_batch(df, suggestion)
            return {"data": ai_result_str, "image_base64": None, "summary": ai_result_str[:100] + "..."}
            
        else:
            err_msg = f"不明な分析タイプ ('{analysis_type}') です: {analysis_name}"
            return {"data": err_msg, "image_base64": None, "summary": err_msg}
            
    except Exception as e:
        logger.error(f"execute_analysis ('{analysis_name}') 実行エラー: {e}", exc_info=True)
        err_msg = f"分析 '{analysis_name}' の実行中にエラーが発生しました: {e}"
        return {"data": err_msg, "image_base64": None, "summary": err_msg}

# --- 8.4. (★) Step B: JSON出力ヘルパー (汎用化対応) ---
def convert_results_to_json_string(results_dict: Dict[str, Any]) -> str:
    """
    (Step B) 実行された分析結果(dict)を、Step Cで読み込むためのJSONL文字列に変換する。
    "OverallSummary" (全体のメトリクス) は特別扱いし、他のタスクのサマリ情報も集約する。
    """
    logger.info(f"JSONL変換開始: {len(results_dict)}件の結果を処理...")
    json_lines = []
    overall_summary_data = {}
    task_summaries = {} # 他タスクのサマリを集約するため

    # --- 1. まず "全体のメトリクス" (OverallSummary) を探す ---
    overall_task_name = "全体のメトリクス"
    if overall_task_name in results_dict:
        result = results_dict[overall_task_name]
        overall_summary_data = {
            "analysis_task": "OverallSummary",
            "data": result.get("data", {"error": "data not found"}),
            "summary": result.get("summary", ""),
            "image_base64": None,
            "image_note": "No image",
            "analysis_summaries": {} # プレースホルダ
        }
    else:
        logger.warning("JSONL変換: '全体のメトリクス' (OverallSummary) が見つかりません。")


    # --- 2. "OverallSummary" 以外のタスクを処理 ---
    for task_name, result in results_dict.items():
        if task_name == overall_task_name:
            continue # 後で処理するのでスキップ

        try:
            line_data = {}
            line_data["analysis_task"] = task_name
            line_data["summary"] = result.get("summary", "N/A")

            # data の型に応じてシリアライズ
            data = result.get("data")
            if isinstance(data, pd.DataFrame):
                # 汎用化: TOP5系タスクはDFだが、中身は辞書のリスト
                if task_name == "カテゴリ別 数値列TOP5分析":
                    line_data["data"] = data.to_dict(orient='records')
                else:
                    # 通常のDFはJSON文字列に (トークン数節約)
                    if len(data) > 500:
                        line_data["data"] = data.head(500).to_json(orient='records', force_ascii=False)
                        line_data["note"] = f"Data truncated. Showing 500 of {len(data)} records."
                    else:
                        line_data["data"] = data.to_json(orient='records', force_ascii=False)
            
            elif isinstance(data, pd.Series):
                line_data["data"] = data.to_dict()
            
            elif isinstance(data, dict) or isinstance(data, list): # メトリクス, A/B比較
                line_data["data"] = data

            elif isinstance(data, str): # AIの回答
                line_data["data"] = data
            
            elif data is None or (hasattr(data, 'empty') and data.empty):
                line_data["data"] = None
                record["note"] = "No data returned from analysis."
            
            else:
                line_data["data"] = str(data)

            # 画像 (Base64) と HTML (pyvis) の処理
            html_content = result.get("html_content") # pyvis用
            image_base64 = result.get("image_base64")

            if html_content:
                 line_data["image_base64"] = None
                 line_data["image_note"] = "No image (pyvis HTML)"
            elif image_base64 and len(image_base64) < (1024 * 1024 * 1.0):
                line_data["image_base64"] = image_base64
                line_data["image_note"] = "Base64 encoded PNG image attached."
            elif image_base64:
                line_data["image_base64"] = None
                line_data["image_note"] = "Image was generated but exceeded 1MB and was not included."
            else:
                line_data["image_base64"] = None
                line_data["image_note"] = "No image generated for this task."

            json_lines.append(json.dumps(line_data, ensure_ascii=False, default=str))
            task_summaries[task_name] = line_data["summary"] # サマリを収集
            
        except Exception as e:
            logger.error(f"JSONL変換エラー ({task_name}): {e}", exc_info=True)
            json_lines.append(json.dumps({"analysis_task": task_name, "error": str(e)}))

    # --- 3. OverallSummary に収集したサマリを結合 ---
    if overall_summary_data:
        overall_summary_data["analysis_summaries"] = task_summaries
        # JSONLの *先頭* に OverallSummary を追加
        json_lines.insert(0, json.dumps(overall_summary_data, ensure_ascii=False))
    
    logger.info(f"JSONL変換完了: {len(json_lines)}行のJSONLを生成。")
    return "\n".join(json_lines)

def render_step_b():
    """(Step B) 分析手法の提案・実行・データ出力UIを描画する"""
    st.title("📊 Step B: インタラクティブ分析とデータ出力")

    # セッションステートの初期化
    if 'df_flagged_B' not in st.session_state:
        st.session_state.df_flagged_B = pd.DataFrame()
    if 'suggestions_B' not in st.session_state: # すべての提案 (タスク名 -> 詳細dict)
        st.session_state.suggestions_B = {}
    if 'selected_tasks_B' not in st.session_state: # ユーザーがチェックしたタスク名 (set)
        st.session_state.selected_tasks_B = set()
    if 'step_b_results' not in st.session_state: # 実行結果 (タスク名 -> 結果dict)
        st.session_state.step_b_results = {}
    if 'step_b_json_output' not in st.session_state:
        st.session_state.step_b_json_output = None
    if 'progress_text' not in st.session_state:
         st.session_state.progress_text = ""
    if 'suggestions_attempted_B' not in st.session_state: # 提案ボタンを押したか
        st.session_state.suggestions_attempted_B = False
        
    if 'tips_list' not in st.session_state:
        st.session_state.tips_list = []
    if 'current_tip_index' not in st.session_state:
        st.session_state.current_tip_index = 0
    if 'last_tip_time' not in st.session_state:
        st.session_state.last_tip_time = time.time()
    # A/B比較用のウィジェット値を保持
    if 'step_b_ab_params' not in st.session_state:
        st.session_state.step_b_ab_params = {}
    # 汎用カテゴリ分析用のウィジェット値を保持
    if 'step_b_generic_params' not in st.session_state:
        st.session_state.step_b_generic_params = {}

    # --- 1. ファイルアップロード ---
    st.header("Step 1: キュレーション済みCSVのアップロード")
    st.info(f"Step A でエクスポートした CSV (Curated_Data.csv) をアップロードしてください。")
    uploaded_flagged_file = st.file_uploader(
        "フラグ付け済みCSVファイル",
        type=['csv'],
        key="step_b_uploader"
    )

    if uploaded_flagged_file:
        try:
            current_file_id = f"{uploaded_flagged_file.name}_{uploaded_flagged_file.size}"
            # ファイルが変更された場合のみリロード＆リセット
            if ('df_flagged_B' not in st.session_state or 
                st.session_state.df_flagged_B.empty or 
                st.session_state.get('current_file_id_B') != current_file_id):
                
                logger.info(f"Step B: 新しいファイル {current_file_id} をロードします。")
                df, err = read_file(uploaded_flagged_file)
                if err:
                    st.error(f"ファイル読み込みエラー: {err}")
                    st.session_state.df_flagged_B = pd.DataFrame()
                    st.session_state.current_file_id_B = None
                    return
                
                st.session_state.df_flagged_B = df
                st.session_state.current_file_id_B = current_file_id
                
                # 関連ステートをすべてリセット
                st.session_state.suggestions_B = {} 
                st.session_state.selected_tasks_B = set()
                st.session_state.step_b_results = {}
                st.session_state.step_b_json_output = None
                st.session_state.suggestions_attempted_B = False
                st.session_state.step_b_ab_params = {}
                st.session_state.step_b_generic_params = {}
                
                st.success(f"ファイル「{uploaded_flagged_file.name}」読込完了 ({len(df)}行)")
                with st.expander("データプレビュー (先頭5行)", expanded=True):
                    st.dataframe(df.head())
            
            else:
                if 'df_flagged_B' in st.session_state and not st.session_state.df_flagged_B.empty:
                    with st.expander("データプレビュー (先頭5行)"):
                        st.dataframe(st.session_state.df_flagged_B.head())
                
        except Exception as e:
            logger.error(f"Step B ファイル読込エラー: {e}", exc_info=True)
            st.error(f"ファイル読み込み中にエラー: {e}")
            st.session_state.df_flagged_B = pd.DataFrame()
            st.session_state.current_file_id_B = None
            return
    else:
        # ファイルがクリアされたら、関連ステートもクリア
        st.session_state.df_flagged_B = pd.DataFrame()
        st.session_state.suggestions_B = {}
        st.session_state.selected_tasks_B = set()
        st.session_state.step_b_results = {}
        st.session_state.step_b_json_output = None
        st.session_state.current_file_id_B = None
        st.session_state.suggestions_attempted_B = False
        st.session_state.step_b_ab_params = {}
        st.session_state.step_b_generic_params = {}
        st.warning("分析を続けるには、Step A で生成したCSVファイルをアップロードしてください。")
        return

    df_B = st.session_state.df_flagged_B
    
    all_cols = list(df_B.columns)
    
    # 汎用カテゴリ列: ...キーワード, カテゴリ, topic など
    base_flag_cols = find_cols(df_B, ['key', 'keyword', 'キーワード', 'カテゴリ', 'topic', 'ハッシュタグ'])
    location_col_search = find_col(df_B, ['市区町村キーワード', 'location', 'city', '地域'])
    # 汎用カテゴリ列 (場所列を除外)
    flag_cols = sorted(list(set([c for c in base_flag_cols if c is not None and c != location_col_search])))
    # 場所列
    location_cols = [location_col_search] if location_col_search else []
    # カテゴリ + 場所
    all_categorical_cols = flag_cols + location_cols
    
    # テキスト列
    object_cols = df_B.select_dtypes(include='object').columns.tolist()
    text_cols = [col for col in object_cols if col not in all_categorical_cols]
    main_text_col_search = find_col(df_B, ['ANALYSIS_TEXT_COLUMN'])
    if main_text_col_search and main_text_col_search in text_cols:
         text_cols.insert(0, text_cols.pop(text_cols.index(main_text_col_search)))
    
    # 日付列
    date_col_search = find_col(df_B, ['date', 'time', '日付', '日時'])
    date_cols = [date_col_search] if date_col_search is not None else []
    
    # 数値列
    numeric_cols = df_B.select_dtypes(include=np.number).columns.tolist()
    engagement_cols = find_engagement_cols(df_B, ['eng', 'like', 'いいね', 'エンゲージメント'])

    # --- 2. 分析手法の提案 ---
    st.header("Step 2: 分析手法の提案")
    st.markdown(f"（(★) AI提案モデル: `{MODEL_FLASH_LITE}`）")
    
    analysis_prompt_B = st.text_area(
        "（任意）AIに追加で指示したい分析タスクを入力:",
        placeholder="例: グルメ投稿と自然投稿の傾向を比較したい。",
        key="step_b_prompt"
    )

    if st.button("💡 分析手法を提案させる (Step 2)", key="suggest_button_B", type="primary"):
        st.session_state.suggestions_attempted_B = True # 提案を実行したフラグ
        
        if not st.session_state.tips_list:
            with st.spinner("分析TIPSをAIで生成中..."):
                st.session_state.tips_list = get_analysis_tips_list_from_ai()
                st.session_state.current_tip_index = random.randint(0, len(st.session_state.tips_list) - 1) if st.session_state.tips_list else 0
                st.session_state.last_tip_time = time.time()
            
        with st.spinner(f"データ構造と指示内容を分析し、手法を提案中 ({MODEL_FLASH_LITE})..."):
            st.session_state.step_b_results = {}
            st.session_state.step_b_json_output = None
            
            base_suggestions = suggest_analysis_techniques_py(df_B)
            ai_suggestions = []
            if analysis_prompt_B.strip():
                ai_suggestions = suggest_analysis_techniques_ai(
                    analysis_prompt_B, df_B, base_suggestions
                )
            base_names = {s['name'] for s in base_suggestions}
            filtered_ai_suggestions = [s for s in ai_suggestions if s['name'] not in base_names]
            all_suggestions = sorted(base_suggestions + filtered_ai_suggestions, key=lambda x: x['priority'])
            
            if not all_suggestions:
                st.session_state.suggestions_B = {}
                st.session_state.selected_tasks_B = set()
            else:
                # 提案を辞書として保存
                st.session_state.suggestions_B = {s['name']: s for s in all_suggestions}
                # デフォルトですべて選択状態にする
                st.session_state.selected_tasks_B = set(st.session_state.suggestions_B.keys())
                st.success(f"分析手法の提案が完了しました ({len(all_suggestions)}件)。Step 3 で実行するタスクを選択してください。")
            
            st.rerun() # 提案後、UIを再描画して Step 3 を表示

    # --- 3. (★) 実行タスクの選択（チェックボックス） ---
    if not st.session_state.suggestions_attempted_B:
        st.info("Step 2 で「分析手法を提案させる」ボタンを押してください。")
        return
    
    if st.session_state.suggestions_attempted_B and not st.session_state.suggestions_B:
        st.warning(
            "分析手法の提案が 0件 でした。\n"
            "アップロードしたCSVファイルに、分析可能な列（`...キーワード`で終わる列や、`ANALYSIS_TEXT_COLUMN`など）が"
            "正しく含まれているか確認してください。"
        )
        return

    st.markdown("---")
    st.header("Step 3: 実行する分析タスクの選択")
    st.info("一括実行したい分析タスクを選択してください。")

    selected_tasks = set()
    
    def select_all_analyses():
        st.session_state.selected_tasks_B = set(st.session_state.suggestions_B.keys())
        
    def deselect_all_analyses():
        st.session_state.selected_tasks_B = set()

    col_select, col_deselect = st.columns(2)
    with col_select:
        st.button("すべて選択", key="select_all_b", use_container_width=True, on_click=select_all_analyses)
    with col_deselect:
        st.button("すべて解除", key="deselect_all_b", use_container_width=True, on_click=deselect_all_analyses)

    st.markdown("---")
    
    cols = st.columns(3)
    i = 0
    sorted_suggestions = sorted(
        st.session_state.suggestions_B.items(), 
        key=lambda item: item[1].get('priority', 99)
    )
    
    for task_name, details in sorted_suggestions:
        with cols[i % 3]:
            is_checked = st.checkbox(
                task_name,
                value=(task_name in st.session_state.selected_tasks_B),
                key=f"cb_{task_name}",
                help=details.get('description', '')
            )
            if is_checked:
                selected_tasks.add(task_name)
        i += 1
    
    st.session_state.selected_tasks_B = selected_tasks


    # --- 4. (★) 選択項目の一括実行 (Bug 1.3 UIフリーズ対応) ---
    if st.button(f"🏃 選択した {len(st.session_state.selected_tasks_B)} 件の分析を実行 (Step 4)", type="primary", use_container_width=True):
        st.session_state.progress_text = "選択項目の実行を開始します..."
        
        progress_text_placeholder_bulk = st.empty()
        progress_text_placeholder_bulk.info(st.session_state.progress_text)

        tasks_to_run = st.session_state.selected_tasks_B
        cleared_results_count = 0
        for task_name in list(st.session_state.step_b_results.keys()):
            if task_name not in tasks_to_run:
                del st.session_state.step_b_results[task_name]
                cleared_results_count += 1
        if cleared_results_count > 0:
            logger.info(f"選択解除された {cleared_results_count} 件の古い結果をクリアしました。")
        
        total_tasks = len(tasks_to_run)
        if total_tasks == 0:
            st.warning("実行するタスクが選択されていません。")
            progress_text_placeholder_bulk.empty()
            st.rerun()

        progress_bar = st.progress(0.0, text="一括実行 待機中...")
        tip_placeholder_b_bulk = st.empty()
        
        with st.spinner(f"全 {total_tasks} 件の分析を実行中..."):
            i = 0
            for task_name in tasks_to_run:
                if task_name not in st.session_state.suggestions_B:
                    continue
                
                suggestion_details = st.session_state.suggestions_B[task_name]
                i += 1
                st.session_state.progress_text = f"({i}/{total_tasks}) 「{task_name}」を実行中..."
                
                progress_bar.progress(i / total_tasks, text=f"実行中: {task_name}")
                progress_text_placeholder_bulk.info(st.session_state.progress_text)
                
                try:
                    result_data = execute_analysis(task_name, df_B, suggestion_details)
                    st.session_state.step_b_results[task_name] = result_data
                except Exception as e:
                    logger.error(f"一括実行エラー ({task_name}): {e}", exc_info=True)
                    st.session_state.step_b_results[task_name] = {
                        "data": f"一括実行中にエラーが発生しました: {e}",
                        "image_base64": None,
                        "summary": f"エラー: {e}"
                    }
            
            st.session_state.progress_text = "全分析の実行が完了しました。"
            progress_bar.progress(1.0, text="実行 完了")
            tip_placeholder_b_bulk.empty()
            progress_text_placeholder_bulk.empty() # 完了したら消す
            st.success("選択された分析の実行が完了しました。Step 5 で結果を確認してください。")
            st.rerun()


    # --- 5. (★) 分析のプレビューとパラメータ修正 (汎用化) ---
    st.markdown("---")
    st.header("Step 5: 分析のプレビューとパラメータ修正")
    
    if not st.session_state.step_b_results:
        st.info("Step 4 で分析を実行すると、ここにプレビューが表示されます。")
        return
    
    st.info("各分析項目の「▼」を開き、プレビューを確認してください。パラメータを修正して、個別に「再実行/更新」も可能です。")
    
    tip_placeholder = st.empty()
    if st.session_state.tips_list:
        try:
            current_tip = st.session_state.tips_list[st.session_state.current_tip_index]
            tip_placeholder.info(f"💡 データ分析TIPS: {current_tip}")
        except IndexError:
            st.session_state.current_tip_index = 0

    progress_text_placeholder = st.empty()
    if st.session_state.progress_text:
         progress_text_placeholder.info(st.session_state.progress_text)
         
    sorted_executed_tasks = sorted(
        st.session_state.step_b_results.keys(),
        key=lambda task_name: st.session_state.suggestions_B.get(task_name, {}).get('priority', 99)
    )

    # 汎用パラメータUIのためのヘルパー
    def get_generic_param(task_name, param_key, default_value):
        return st.session_state.step_b_generic_params.get(task_name, {}).get(param_key, default_value)

    def set_generic_param(task_name, param_key, value):
        if task_name not in st.session_state.step_b_generic_params:
            st.session_state.step_b_generic_params[task_name] = {}
        st.session_state.step_b_generic_params[task_name][param_key] = value

    for task_name in sorted_executed_tasks:
        if task_name not in st.session_state.suggestions_B:
            continue
            
        suggestion_details = st.session_state.suggestions_B[task_name].copy()
        result = st.session_state.step_b_results[task_name]
        
        st.markdown("---")
        
        # プレビュー表示
        st.subheader(f"✅ プレビュー: {task_name}")
        
        if task_name == "カテゴリ別 数値列TOP5分析" and isinstance(result.get('data'), pd.DataFrame):
            for _, row in result['data'].iterrows():
                st.markdown(f"**{row['category']}** (投稿数: {row['post_count']})")
                if row['top_posts']:
                     for post in row['top_posts']:
                         st.markdown(f"  - **EG: {post['engagement']}** - {post['summary_ai']}")
                         if post.get('media_link'):
                             st.markdown(f"    [Link]({post['media_link']})")
                st.markdown("---")
        elif task_name == "A/B 比較分析" and isinstance(result.get('data'), dict):
            if "category_comparison" in result["data"]:
                st.markdown("##### カテゴリ別 投稿数比較")
                st.dataframe(pd.DataFrame(result["data"]["category_comparison"]))
            if "ranking_comparison" in result["data"]:
                st.markdown("##### 地域別 順位変動 (Top 20)")
                st.dataframe(pd.DataFrame(result["data"]["ranking_comparison"]))
        
        # 共起ネットワーク (pyvis HTML)
        elif task_name == "共起ネットワーク" and result.get("html_content"):
            components.html(result.get("html_content"), height=710)
            ai_legend_map = result.get("ai_legend")
            communities_map = result.get("communities")
            
            if ai_legend_map:
                st.markdown("##### AIによる推定トピック:")
                legend_items = []
                for group_id, topic in ai_legend_map.items():
                    color = COLOR_PALETTE[group_id % len(COLOR_PALETTE)]
                    legend_html = f"""
                    <span style="display: inline-block; margin: 4px; padding: 8px 12px; border-radius: 8px; background-color: #f0f2f6; border: 1px solid #e0e0e0;">
                        <span style='color:{color}; font-size: 20px; font-weight: bold; vertical-align: middle;'>■</span>
                        <span style="vertical-align: middle; margin-left: 8px; font-size: 14px;">{topic} (G{group_id})</span>
                    </span>
                    """
                    legend_items.append(legend_html.replace("\n", ""))
                st.markdown("<div style='line-height: 1.8;'>" + " ".join(legend_items) + "</div>", unsafe_allow_html=True)
            
            elif communities_map:
                st.markdown("##### 検出されたコミュニティ (色分け):")
                legend_items = []
                for group_id in communities_map.keys():
                    color = COLOR_PALETTE[group_id % len(COLOR_PALETTE)]
                    legend_html = f"""
                    <span style="display: inline-block; margin: 4px; padding: 8px 12px; border-radius: 8px; background-color: #f0f2f6; border: 1px solid #e0e0e0;">
                        <span style='color:{color}; font-size: 20px; font-weight: bold; vertical-align: middle;'>■</span>
                        <span style="vertical-align: middle; margin-left: 8px; font-size: 14px;">グループ {group_id}</span>
                    </span>
                    """
                    legend_items.append(legend_html.replace("\n", ""))
                st.markdown("<div style='line-height: 1.8;'>" + " ".join(legend_items) + "</div>", unsafe_allow_html=True)

        # その他の分析
        else:
            if result.get('image_base64'):
                st.image(base64.b64decode(result['image_base64']))
            
            if isinstance(result.get('data'), pd.DataFrame):
                st.dataframe(result['data'].head(10))
            elif isinstance(result.get('data'), dict):
                st.json(result['data'])
            elif isinstance(result.get('data'), str):
                st.markdown(result['data'])
            
        st.caption(f"サマリ: {result.get('summary', 'N/A')}")

        # (★) 個別実行エリア (汎用化対応)
        with st.expander(f"「{task_name}」のパラメータ修正・再実行"):
            
            st.markdown(f"**説明:** {suggestion_details.get('description', 'N/A')}")
            st.markdown("##### (オプション) パラメータの変更")
            
            try:
                # 1. 単純集計
                if task_name.startswith("単純集計:"):
                    default_col = suggestion_details['suitable_cols'][0]
                    new_col = st.selectbox(f"集計対象の列 ({task_name})", options=all_categorical_cols, index=all_categorical_cols.index(default_col) if default_col in all_categorical_cols else 0, key=f"sel_{task_name}")
                    suggestion_details['ui_selected_col'] = new_col
                
                # 2. クロス集計
                elif task_name.startswith("クロス集計"):
                    default_col1 = suggestion_details['suitable_cols'][0]
                    default_col2 = suggestion_details['suitable_cols'][1]
                    c1, c2 = st.columns(2)
                    new_col1 = c1.selectbox(f"列 1 (行) ({task_name})", options=all_categorical_cols, index=all_categorical_cols.index(default_col1) if default_col1 in all_categorical_cols else 0, key=f"sel_{task_name}_1")
                    new_col2 = c2.selectbox(f"列 2 (列) ({task_name})", options=all_categorical_cols, index=all_categorical_cols.index(default_col2) if default_col2 in all_categorical_cols else 1, key=f"sel_{task_name}_2")
                    suggestion_details['ui_selected_col1'] = new_col1
                    suggestion_details['ui_selected_col2'] = new_col2

                # 3. 時系列
                elif task_name == "時系列キーワード分析":
                    default_dt = suggestion_details['suitable_cols']['datetime'][0]
                    default_kw = suggestion_details['suitable_cols']['keywords'][0]
                    c1, c2 = st.columns(2)
                    new_dt = c1.selectbox(f"日時列 ({task_name})", options=date_cols, index=date_cols.index(default_dt) if default_dt in date_cols else 0, key=f"sel_{task_name}_dt")
                    new_kw = c2.selectbox(f"キーワード列 ({task_name})", options=all_categorical_cols, index=all_categorical_cols.index(default_kw) if default_kw in all_categorical_cols else 0, key=f"sel_{task_name}_kw")
                    suggestion_details['ui_selected_dt_col'] = new_dt
                    suggestion_details['ui_selected_kw_col'] = new_kw
                
                # 4. テキストマイニング
                elif task_name == "テキストマイニング（頻出単語）":
                    default_col = suggestion_details['suitable_cols'][0]
                    if not text_cols:
                         st.warning("分析対象のテキスト列が見つかりません。")
                         new_col = None
                    else:
                        new_col = st.selectbox(f"テキスト列 ({task_name})", options=text_cols, index=text_cols.index(default_col) if default_col in text_cols else 0, key=f"sel_{task_name}_txt")
                    suggestion_details['ui_selected_text_col'] = new_col

                    custom_sw = st.text_area(
                        "除外したい単語（カンマ, スペース, 改行区切り）:",
                        value=suggestion_details.get('ui_custom_stop_words', ''),
                        key=f"sw_{task_name}",
                        height=100,
                        placeholder="例: 弊社, 商品A, サービスB, ..."
                    )
                    suggestion_details['ui_custom_stop_words'] = custom_sw

                # 5. 共起ネットワーク
                elif task_name == "共起ネットワーク":
                    # 1. 絞り込み列
                    flag_col_options = all_categorical_cols
                    default_flag_col = suggestion_details.get('ui_selected_flag_col', location_cols[0] if location_cols else (flag_cols[0] if flag_cols else None))
                    
                    flag_col = st.selectbox(
                        "1. 絞り込みに使用するカテゴリ列:", flag_col_options,
                        index=flag_col_options.index(default_flag_col) if default_flag_col in flag_col_options else 0,
                        key=f"cn_filter_col_{task_name}",
                        help="ここで選んだ列のキーワードで、分析対象の投稿を絞り込みます。"
                    )
                    suggestion_details['ui_selected_flag_col'] = flag_col

                    # 2. 絞り込みキーワード
                    try:
                        s = df_B[flag_col].dropna().astype(str).str.split(',').explode().str.strip()
                        s = s[~s.isin(['', 'nan', 'Nan', 'NaN'])]
                        keyword_counts = s.value_counts()
                        options = keyword_counts.index.tolist()[:50]
                        default_options_kws = suggestion_details.get('ui_selected_keywords', keyword_counts.index.tolist()[:10])
                    except Exception:
                        options = []
                        default_options_kws = suggestion_details.get('ui_selected_keywords', [])

                    selected_keywords = st.multiselect(
                        f"2. 絞り込むキーワード（「{flag_col}」列 Top 50）:",
                        options,
                        default=default_options_kws,
                        key=f"cn_selected_keywords_{flag_col}",
                        help="分析対象とする投稿に含まれるキーワードを選択します。"
                    )
                    suggestion_details['ui_selected_keywords'] = selected_keywords
                    
                    # 3. テキスト列
                    default_text_col = suggestion_details.get('ui_selected_text_col', text_cols[0] if text_cols else None)
                    if not text_cols:
                         st.warning("分析対象のテキスト列が見つかりません。")
                         text_col = None
                    else:
                        text_col = st.selectbox(
                            "3. 分析対象の自由記述列:", text_cols,
                            index=text_cols.index(default_text_col) if default_text_col in text_cols else 0,
                            key=f"cn_text_col_{task_name}"
                        )
                    suggestion_details['ui_selected_text_col'] = text_col
                    
                    st.markdown("---")
                    st.markdown("**テキスト分析 設定**")
                    custom_sw_cn = st.text_area(
                        "除外したい単語（カンマ, スペース, 改行区切り）:",
                        value=suggestion_details.get('ui_custom_stop_words', ''),
                        key=f"sw_{task_name}",
                        height=100,
                        placeholder="例: 弊社, 商品A, サービスB, ..."
                    )
                    suggestion_details['ui_custom_stop_words'] = custom_sw_cn
                    
                    st.markdown("---")
                    st.markdown("**グラフ詳細設定**")
                    ui_cols = st.columns([0.5, 0.5])
                    
                    with ui_cols[0]:
                        # 4. レイアウト
                        st.markdown("**レイアウト・物理演算**")
                        solver = st.selectbox(
                            "レイアウト (layout)", ['barnesHut', 'fruchterman_reingold', 'repulsion'],
                            index=['barnesHut', 'fruchterman_reingold', 'repulsion'].index(suggestion_details.get('solver', 'barnesHut')),
                            key=f"cn_solver_{task_name}"
                        )
                        suggestion_details['solver'] = solver
                        gravity = st.slider(
                            "重力 (Gravity)", -50000, -1000, suggestion_details.get('gravity', -2000), step=1000,
                            key=f"cn_gravity_{task_name}"
                        )
                        suggestion_details['gravity'] = gravity
                        node_distance = st.slider(
                            "ノード間の反発力", 100, 500, suggestion_details.get('node_distance', 200),
                            key=f"cn_distance_{task_name}"
                        )
                        suggestion_details['node_distance'] = node_distance
                        spring_length = st.slider(
                            "エッジの長さ", 50, 500, suggestion_details.get('spring_length', 250),
                            key=f"cn_spring_{task_name}"
                        )
                        suggestion_details['spring_length'] = spring_length

                    with ui_cols[1]:
                        # 5. フィルタ
                        st.markdown("**フィルタ設定**")
                        top_n_words_limit = st.slider(
                            "分析対象の単語数 (Top N)", 50, 300, suggestion_details.get('top_n_words_limit', 100),
                            key=f"cn_top_n_{task_name}"
                        )
                        suggestion_details['top_n_words_limit'] = top_n_words_limit
                        max_degree_cutoff = st.slider(
                            "最大接続数 (Exclude Hubs)", 10, 100, suggestion_details.get('max_degree_cutoff', 50),
                            key=f"cn_max_degree_{task_name}"
                        )
                        suggestion_details['max_degree_cutoff'] = max_degree_cutoff
                        min_occurrence = st.slider(
                            "最小共起回数 (Min Freq)", 1, 30, suggestion_details.get('min_occurrence', 10),
                            key=f"cn_slider_v3_{task_name}"
                        )
                        suggestion_details['min_occurrence'] = min_occurrence
                        
                        # 6. デザイン
                        st.markdown("**デザイン設定**")
                        default_node_size = st.slider(
                            "基準ノードサイズ", 5, 50, suggestion_details.get('default_node_size', 15),
                            key=f"cn_node_size_v2_{task_name}"
                        )
                        suggestion_details['default_node_size'] = default_node_size
                        default_text_size = st.slider(
                            "テキストサイズ", 10, 100, suggestion_details.get('default_text_size', 50),
                            key=f"cn_text_size_v2_{task_name}"
                        )
                        suggestion_details['default_text_size'] = default_text_size
                    
                    # 7. AI凡例
                    st.markdown("---")
                    run_ai_legend = st.checkbox(
                        "🤖 AIで凡例を生成 (β) (実行に時間がかかります)",
                        value=suggestion_details.get('run_ai_legend', False),
                        key=f"cn_run_ai_legend_{task_name}"
                    )
                    suggestion_details['run_ai_legend'] = run_ai_legend

                # 6. 汎用カテゴリ深掘り
                elif task_name == "カテゴリ列の集計と深掘り":
                    defaults = suggestion_details['suitable_cols']
                    cat_options = defaults['category_cols']
                    if not cat_options:
                         st.warning("分析可能なカテゴリ列（...キーワード 等）が見つかりません。")
                         new_cat_col = None
                    else:
                        default_cat_col = get_generic_param(task_name, 'cat_col', cat_options[0])
                        new_cat_col = st.selectbox(
                            f"集計対象のカテゴリ列 ({task_name})", options=cat_options, 
                            index=cat_options.index(default_cat_col) if default_cat_col in cat_options else 0, 
                            key=f"sel_{task_name}_cat"
                        )
                    suggestion_details['ui_selected_category_col'] = new_cat_col
                    set_generic_param(task_name, 'cat_col', new_cat_col)

                # 7. 汎用 数値列TOP5
                elif task_name == "カテゴリ別 数値列TOP5分析":
                    defaults = suggestion_details['suitable_cols']
                    
                    cat_options = defaults['category_cols']
                    num_options = defaults['numeric_cols']
                    
                    if not cat_options or not num_options:
                        st.warning("分析に必要なカテゴリ列または数値列（いいね 等）が見つかりません。")
                    else:
                        c1, c2 = st.columns(2)
                        default_cat_col = get_generic_param(task_name, 'cat_col', cat_options[0])
                        default_num_col = get_generic_param(task_name, 'num_col', num_options[0])

                        new_cat_col = c1.selectbox(
                            f"カテゴリ列 ({task_name})", options=cat_options, 
                            index=cat_options.index(default_cat_col) if default_cat_col in cat_options else 0, 
                            key=f"sel_{task_name}_cat_top5"
                        )
                        new_num_col = c2.selectbox(
                            f"数値列（集計対象） ({task_name})", options=num_options, 
                            index=num_options.index(default_num_col) if default_num_col in num_options else 0, 
                            key=f"sel_{task_name}_num_top5"
                        )
                        
                        suggestion_details['ui_selected_category_col'] = new_cat_col
                        suggestion_details['ui_selected_numeric_col'] = new_num_col
                        suggestion_details['ui_selected_text_col'] = defaults['text_col'][0] # テキスト列は固定
                        
                        set_generic_param(task_name, 'cat_col', new_cat_col)
                        set_generic_param(task_name, 'num_col', new_num_col)

                # 8. A/B比較のUI
                elif task_name == "A/B 比較分析":
                    st.info("比較したい2つのグループ（AとB）を定義してください。")
                    ab_col_options = suggestion_details['suitable_cols']['category_cols']
                    
                    ab_params = st.session_state.step_b_ab_params
                    
                    c1, c2 = st.columns(2)
                    with c1:
                        st.markdown("##### グループ A")
                        a_col_key = f"ab_a_col_{task_name}"
                        a_val_key = f"ab_a_val_{task_name}"
                        
                        default_a_col = ab_params.get('a_col', ab_col_options[0] if ab_col_options else None)
                        a_col = st.selectbox(
                            "A: 比較列", ab_col_options, 
                            index=ab_col_options.index(default_a_col) if default_a_col in ab_col_options else 0, 
                            key=a_col_key
                        )
                        try:
                            a_val_options = sorted(list(df_B[a_col].astype(str).str.split(', ').explode().str.strip().unique()))
                        except Exception:
                            a_val_options = []
                        
                        default_a_val = ab_params.get('a_val', a_val_options[0] if a_val_options else None)
                        a_val = st.selectbox(
                            "A: 比較値", a_val_options, 
                            index=a_val_options.index(default_a_val) if default_a_val in a_val_options else 0, 
                            key=a_val_key
                        )
                    with c2:
                        st.markdown("##### グループ B")
                        b_col_key = f"ab_b_col_{task_name}"
                        b_val_key = f"ab_b_val_{task_name}"
                        
                        default_b_col = ab_params.get('b_col', ab_col_options[0] if ab_col_options else None)
                        b_col = st.selectbox(
                            "B: 比較列", ab_col_options, 
                            index=ab_col_options.index(default_b_col) if default_b_col in ab_col_options else 0, 
                            key=b_col_key
                        )
                        try:
                            b_val_options = sorted(list(df_B[b_col].astype(str).str.split(', ').explode().str.strip().unique()))
                        except Exception:
                            b_val_options = []
                            
                        default_b_val = ab_params.get('b_val', b_val_options[1] if len(b_val_options) > 1 else (b_val_options[0] if b_val_options else None))
                        b_val = st.selectbox(
                            "B: 比較値", b_val_options, 
                            index=b_val_options.index(default_b_val) if default_b_val in b_val_options else 0, 
                            key=b_val_key
                        )
                    
                    current_ab_params = {'a_col': a_col, 'a_val': a_val, 'b_col': b_col, 'b_val': b_val}
                    suggestion_details['ui_ab_params'] = current_ab_params
                    st.session_state.step_b_ab_params = current_ab_params
                
                # 9. (AIタスク)
                elif suggestion_details.get('type') == 'ai':
                    st.info("このタスクはAIによって実行されます。AIへの指示（説明）を変更できます。")
                    new_desc = st.text_area(
                        "AIへの指示 (description):",
                        value=suggestion_details.get('description', ''),
                        key=f"ai_desc_{task_name}",
                        height=100
                    )
                    suggestion_details['description'] = new_desc

            except Exception as e:
                st.error(f"パラメータUIの描画に失敗: {e}")
                logger.error(f"パラメータUI描画エラー ({task_name}): {e}", exc_info=True)

            # 個別実行ボタン
            if st.button(f"「{task_name}」を再実行/更新", key=f"run_{task_name}"):
                st.session_state.progress_text = f"「{task_name}」を個別に実行中..."
                with st.spinner(f"「{task_name}」を実行中..."):
                    try:
                        result_data = execute_analysis(task_name, df_B, suggestion_details)
                        st.session_state.step_b_results[task_name] = result_data # 結果を更新
                        st.session_state.suggestions_B[task_name] = suggestion_details # パラメータを保存
                        st.session_state.progress_text = f"「{task_name}」の実行が完了しました。"
                        st.rerun() # UIを更新してプレビューに反映
                    except Exception as e:
                         st.error(f"分析実行エラー: {e}")
                         logger.error(f"個別実行エラー ({task_name}): {e}", exc_info=True)
                         st.session_state.progress_text = f"「{task_name}」の実行に失敗しました。"


    # --- 6. (★) 最終エクスポート ---
    st.markdown("---")
    st.header("Step 6: 最終JSONのエクスポート")
    
    total_results = len(st.session_state.step_b_results)
    
    if total_results == 0:
        st.warning("Step 4 で分析を実行してください。")
    else:
        st.success(f"現在 {total_results} 件の分析結果がプレビューされています。")

    
    if st.button("StepC用 JSONを生成・エクスポート (Step 6)", type="primary", use_container_width=True):
        if total_results == 0:
            st.error("分析が1つも実行されていません。Step 4 で分析を実行してください。")
        else:
            with st.spinner("最終JSONファイルを生成中..."):
                try:
                    json_output_string = convert_results_to_json_string(st.session_state.step_b_results)
                    st.session_state.step_b_json_output = json_output_string
                    st.success("StepC用のJSONデータが生成されました！")
                except Exception as e:
                    logger.error(f"Step B 最終JSON出力変換エラー: {e}", exc_info=True)
                    st.error(f"分析結果のJSON変換中にエラー: {e}")

    # ダウンロードセクション
    if st.session_state.step_b_json_output:
        st.info(f"以下のJSONファイルには、Step 5 でプレビュー・実行された {len(st.session_state.step_b_results)} 件の分析結果がすべて含まれています。")
        
        st.download_button(
            label="分析データ (analysis_data.json) をダウンロード",
            data=st.session_state.step_b_json_output,
            file_name="analysis_data.json",
            mime="application/json",
            type="primary",
            use_container_width=True
        )
        
        st.markdown("---")
        st.subheader("出力データ (JSONL) プレビュー")
        
        preview_summaries = []
        try:
            for line in st.session_state.step_b_json_output.splitlines():
                line_data = json.loads(line)
                task_name = line_data.get("analysis_task")
                summary = line_data.get("summary", line_data.get("analysis_summaries", "No summary."))
                img_note = line_data.get("image_note", "No image.")
                
                if task_name == "OverallSummary":
                    preview_summaries.append(f"--- Overall Summary ---")
                    if isinstance(summary, dict):
                        for k, v in summary.items():
                             preview_summaries.append(f"  - {k}: {str(v)[:100]}...")
                    continue
                
                preview_summaries.append(f"[{task_name}] (Image: {img_note})")
                
        except Exception as e:
            preview_summaries = ["JSONプレビューの生成に失敗", str(e)]

        st.text_area(
            "JSONL (サマリープレビュー):",
            value="\n".join(preview_summaries),
            height=300,
            key="json_preview_B_summary",
            disabled=True
        )
        st.success("データをダウンロードし、Step C (AIレポート生成) に進んでください。")

# --- 9. (★) Step C: AIレポート生成 (Proモデル) ---
# (要件: Step Cは gemini-2.5-pro を使用)

def run_step_c_analysis(
    jsonl_data_string: str,
    model_name: str,
    progress_bar: st.delta_generator.DeltaGenerator,
    log_placeholder: st.delta_generator.DeltaGenerator,
    custom_instruction: str = "" # (★) 改善 C-4: UIからカスタム指示を受け取る
) -> str:
    """
    (★) Step C: AIレポート生成 (ハイブリッド・RateLimit対応版)
    
    [新ロジック] ハングアップ (504 Timeout) を回避するため、AIにBase64画像(巨大トークン)を
    渡すのを *やめ* 、「テキストデータ」のみを渡して考察を生成させる。
    Python側で、AIの考察(テキスト)と、元のBase64画像を「再結合」する。
    """
    logger.info(f"Step C AIレポート生成 (ハイブリッド処理) 開始... (Model: {model_name})")

    # (★) --- 1. モデルのRPM制限とスリープ時間を定義 ---
    if model_name == MODEL_PRO:
        rpm_limit = 2
        tpm_limit = 125000
    else: # (★) デフォルトは Flash
        model_name = MODEL_FLASH
        rpm_limit = 10
        tpm_limit = 250000
        
    sleep_time = (60 / rpm_limit) + 0.5 # (e.g., Pro: 30.5s, Flash: 6.5s)
    
    logger.info(f"モデル: {model_name}, RPM: {rpm_limit}, 待機: {sleep_time:.1f}秒")

    # (★) --- 2. チャンク生成用のAIプロンプトテンプレートを定義 (品質向上) ---
    
    # (★) [改善 C-4] カスタム指示が空でない場合、プロンプトに挿入するブロックを定義
    custom_instruction_block = ""
    if custom_instruction and custom_instruction.strip():
        custom_instruction_block = f"""
        # (重要) ユーザーからの追加指示:
        * {custom_instruction.strip()}
        * この指示を最優先で考慮してください。
        """

    # (★) [改善 C-1, C-2]
    ITERATIVE_SLIDE_PROMPT_TEMPLATE = """
    あなたはシニアデータアナリストであり、クライアント向けレポートの「スライド1枚」の
    【テキスト部分】を作成しています。
    提供される「分析タスクデータ」を読み、このタスク専用の
    スライドタイトルと考察（slide_content）を生成してください。

    # 分析タスクデータ (テキスト・数値データのみ):
    {task_data_text_only}
    
    # (★) [改善 C-2] 画像コンテキスト:
    {image_context}
    
    # (★) [改善 C-4] ユーザーの全体方針:
    {custom_instruction}

    # 指示:
    1.  **タイトル**: `task_data_text_only` の `analysis_task` 名に基づき、 professional な「slide_title」を考案してください。
    
    2.  **(★) [改善 C-1] 考察 (最重要)**: 
        `task_data_text_only` の `summary` と `data`（テーブルデータ）を解釈し、クライアントが知るべき【インサイト】を **Markdownの箇条書き** で記述してください。
        以下の3つの視点（何を・なぜ・だから何）で構成してください。
        
        * **何を（What）:** データが示す最も重要な「事実」や「傾向」は何か？ (例: `**〇〇** が **XX%** 増加...`)
        * **なぜ（Why）:** なぜその傾向が起きているのか？（背景や原因の「仮説」）
        * **だから何（So What）:** この事実から推測できる「次のアクションのヒント」は何か？

    3.  **書式**:
        - 回答は【Markdown形式】を使用してください。
        - 重要なキーワードや数値は `**太字**` で強調してください。

    # 出力形式 (厳守):
    * JSON以外のテキストは絶対に含めず、【単一のJSONオブジェクト】`{{ ... }}` のみを出力してください。
    * 以下の構造を厳格に守ってください。
        {{
          "slide_title": "（指示1で考案したタイトル）",
          "slide_content": [
            "（指示2, 3 に基づく Markdown 形式のインサイト1: **何を**...）",
            "（指示2, 3 に基づく Markdown 形式のインサイト2: **なぜ**...）",
            "（指示2, 3 に基づく Markdown 形式のインサイト3: **だから何**...）"
          ]
        }}

    # 回答 (単一のJSONオブジェクトのみ):
    """
    
    prompt = PromptTemplate.from_template(ITERATIVE_SLIDE_PROMPT_TEMPLATE)

    # (★) タイムアウトを 120秒 (2分) に設定
    llm = get_llm(model_name=model_name, temperature=0.2, timeout_seconds=120)
    if llm is None:
        st.error(f"AIモデル({model_name})が利用できません。")
        return "[]" # 空のJSONリスト
    
    chain = prompt | llm | StrOutputParser()

    # (★) --- 3. 逐次処理ループ (変更なし) ---
    report_slides_list = []
    log_messages_ui = []
    
    tasks_all = jsonl_data_string.strip().splitlines()
    
    # 3.1. OverallSummaryを抽出し、残りを処理対象タスクとする
    summary_line = "{}"
    tasks_to_process = []
    for line in tasks_all:
        if '"analysis_task": "OverallSummary"' in line:
            summary_line = line
        else:
            tasks_to_process.append(line)
            
    if not tasks_to_process:
        logger.warning("処理対象の分析タスクが0件です。")
        return "[]"

    total_tasks = len(tasks_to_process)
    logger.info(f"全 {total_tasks} タスクを逐次処理します。")
    
    # 3.2. 表紙スライドを追加
    report_slides_list.append({
        "slide_title": "SNSデータ分析レポート",
        "slide_layout": "title_only",
        "slide_content": ["AI-Generated Analysis (Powered by Gemini)"],
        "image_base64": None
    })
    
    # 3.3. 目次スライドを追加 (この時点ではタスク名のみ)
    try:
        agenda_items = []
        for i, task_line in enumerate(tasks_to_process):
            try:
                task_name = json.loads(task_line).get('analysis_task', f'分析タスク {i+1}')
            except json.JSONDecodeError:
                task_name = f'分析タスク {i+1} (読み込みエラー)'
            agenda_items.append(f"{i+1}. {task_name}")

        agenda_items.append(f"{len(tasks_to_process) + 1}. 結論と戦略的提言")
        report_slides_list.append({
            "slide_title": "本日のアジェンダ",
            "slide_layout": "title_and_content",
            "slide_content": agenda_items,
            "image_base64": None
        })
    except Exception as e:
        logger.error(f"目次スライドの生成に失敗: {e}")

    # 3.4. メインの分析スライドをループ処理
    for i, task_line in enumerate(tasks_to_process):
        
        task_name = f"Task {i+1}/{total_tasks}"
        original_task_json = {}
        
        try:
            # 3.4.1. タスクのパースと画像/テキストの分離
            original_task_json = json.loads(task_line)
            task_name = original_task_json.get('analysis_task', task_name)

            # 1. 画像をPython変数に退避
            image_to_pass_through = original_task_json.get("image_base64")
            
            # (★) [改善 C-2] 画像コンテキストを定義
            image_context_str = "（このスライドには画像は含まれません。）"
            if image_to_pass_through:
                image_context_str = (
                    "（(注) このスライドにはグラフやワードクラウド等の「データ可視化画像」が1枚含まれます。\n"
                    "   あなたに画像は見えませんが、`data` や `summary` を根拠に、"
                    "   その画像が「何を意味するのか」を解説する考察を記述してください。）"
                )
            
            # 2. AIに渡す「テキストのみ」のJSONを作成
            text_only_task_json = original_task_json.copy()
            text_only_task_json["image_base64"] = None
            if "data" in text_only_task_json and len(json.dumps(text_only_task_json["data"])) > 1000:
                text_only_task_json["data"] = f"（データプレビュー: {str(text_only_task_json['data'])[:1000]}...）"
            
            task_data_text_only_str = json.dumps(text_only_task_json, ensure_ascii=False) # (★) ensure_ascii=False
            
        except Exception as e:
            logger.error(f"タスク '{task_name}' のJSONパースに失敗: {e}")
            log_messages_ui.append(f"  -> ERROR: '{task_name}' のJSONパースに失敗。スキップします。")
            continue

        # 3.4.2. UI（進捗バー・ログ）の更新
        progress_percent = (i + 1) / (total_tasks + 1)
        progress_bar.progress(progress_percent, text=f"Step C (スライド生成中): {i+1}/{total_tasks} (モデル: {model_name})")
        log_messages_ui.append(f"[{i+1}/{total_tasks}] '{task_name}' の処理を開始...")
        log_placeholder.text_area("実行ログ:", "\n".join(log_messages_ui[::-1]), height=250, key=f"step_c_log_{i}")

        try:
            # 3.4.3. AIへのリクエスト (テキストのみ)
            log_messages_ui.append(f"  -> AI ({model_name}) にリクエストを送信... (Timeout: 120s)")
            log_placeholder.text_area("実行ログ:", "\n".join(log_messages_ui[::-1]), height=250, key=f"step_c_log_{i}_sending")
            
            response_str = chain.invoke({
                "task_data_text_only": task_data_text_only_str,
                "image_context": image_context_str, # (★) C-2
                "custom_instruction": custom_instruction_block # (★) C-4
            })
            
            log_messages_ui.append(f"  -> AI が応答しました。レスポンスを解析中...")
            log_placeholder.text_area("実行ログ:", "\n".join(log_messages_ui[::-1]), height=250, key=f"step_c_log_{i}_received")

            # (★) [改善 C-3] 堅牢なJSONパース
            start = response_str.find('{')
            end = response_str.rfind('}')
            
            if start != -1 and end != -1 and end > start:
                json_str = response_str[start:end+1]
                ai_response_json = json.loads(json_str)
                
                # 3.4.4. AIの考察と、退避させた画像を「再結合」
                final_slide_object = {
                    "slide_title": ai_response_json.get("slide_title", task_name),
                    "slide_layout": "text_and_image" if image_to_pass_through else "title_and_content",
                    "slide_content": ai_response_json.get("slide_content", ["AIによる考察の生成に失敗しました。"]),
                    "image_base64": image_to_pass_through # (★) ここで画像を戻す
                }
                report_slides_list.append(final_slide_object)
                log_messages_ui.append(f"  -> SUCCESS: スライド '{final_slide_object.get('slide_title')}' を生成しました。")
            else:
                raise Exception("AIがJSONオブジェクト `{{...}}` を返しませんでした。")
        
        except Exception as e:
            logger.error(f"タスク '{task_name}' の処理に失敗: {e}", exc_info=True)
            log_messages_ui.append(f"  -> ERROR: '{task_name}' の処理に失敗。{e}")
            report_slides_list.append({
                "slide_title": f"エラー: {task_name}",
                "slide_layout": "title_and_content",
                "slide_content": [f"このスライドの生成に失敗しました。", f"エラー: {e}"],
                "image_base64": None
            })
        
        # 3.4.5. Rate Limit のための待機
        if i < total_tasks:
            log_messages_ui.append(f"  -> Rate Limit (RPM) のため {sleep_time:.1f} 秒待機します...")
            log_placeholder.text_area("実行ログ:", "\n".join(log_messages_ui[::-1]), height=250, key=f"step_c_log_{i}_sleep")
            time.sleep(sleep_time)

    # (★) 4. 結論スライドの生成 (品質向上)
    try:
        chunk_name = f"結論スライド"
        progress_percent = 1.0
        progress_bar.progress(progress_percent, text=f"Step C (チャンク処理中): {chunk_name} (モデル: {model_name})")
        log_messages_ui.append(f"[{total_tasks+1}/{total_tasks+1}] {chunk_name} の処理を開始...")
        log_placeholder.text_area("実行ログ:", "\n".join(log_messages_ui[::-1]), height=250, key="step_c_log_final")

        conclusion_llm = get_llm(model_name=model_name, temperature=0.2, timeout_seconds=120)
        if conclusion_llm is None:
            raise Exception("結論スライド用AIモデルの取得に失敗")

        # (★) [改善 C-1, C-4] 結論プロンプトも強化
        CONCLUSION_PROMPT_TEMPLATE = """
        あなたはシニアデータアナリストです。
        以下の「分析サマリー」と「生成したスライドタイトル」に基づき、
        レポートの締めくくりとなる【結論と戦略的提言】のスライド1枚分のJSONオブジェクトを生成してください。

        # 分析サマリー (OverallSummary):
        {summary_data_line}
        
        # 生成済みスライドタイトル:
        {slide_titles}

        # (★) ユーザーの全体方針:
        {custom_instruction}

        # 指示:
        1.  タイトルは「結論と戦略的提言」とします。
        2.  レイアウトは「title_and_content」とします。
        3.  **(★) [改善 C-1] 内容**:
            分析全体から導かれる「結論（主要な発見）」と、クライアントが次に取るべき「具体的なアクション（提言）」を、Markdownの箇条書きで3〜5点にまとめてください。
            
            * **結論 (Key Findings):** （例: `**〇〇** が最も重要な課題であると判明...`）
            * **提言 (Recommendations):** （例: `**〇〇** にリソースを集中投下し、...`）
        
        4.  画像 (image_base64) は null とします。

        # 出力形式 (厳守):
        * JSON以外のテキストは絶対に含めず、【単一のJSONオブジェクト】`{{ ... }}` のみを出力してください。

        # 回答 (単一のJSONオブジェクトのみ):
        """
        
        conclusion_prompt = PromptTemplate.from_template(CONCLUSION_PROMPT_TEMPLATE)
        conclusion_chain = conclusion_prompt | conclusion_llm | StrOutputParser()
        
        log_messages_ui.append(f"  -> AI ({model_name}) にリクエストを送信... (Timeout: 120s)")
        log_placeholder.text_area("実行ログ:", "\n".join(log_messages_ui[::-1]), height=250, key="step_c_log_final_sending")
        
        response_str = conclusion_chain.invoke({
            "summary_data_line": summary_line,
            "slide_titles": json.dumps([s.get('slide_title') for s in report_slides_list], ensure_ascii=False),
            "custom_instruction": custom_instruction_block # (★) C-4
        })
        
        log_messages_ui.append(f"  -> AI が応答しました。レスポンスを解析中...")
        log_placeholder.text_area("実行ログ:", "\n".join(log_messages_ui[::-1]), height=250, key="step_c_log_final_received")

        # (★) [改善 C-3] 堅牢なJSONパース
        start = response_str.find('{')
        end = response_str.rfind('}')
            
        if start != -1 and end != -1 and end > start:
            json_str = response_str[start:end+1]
            report_slides_list.append(json.loads(json_str))
            log_messages_ui.append(f"  -> SUCCESS: 結論スライドを生成しました。")
        else:
            raise Exception("AIが結論スライドのJSONを返しませんでした。")
            
    except Exception as e:
         logger.error(f"結論スライドの生成に失敗: {e}")
         log_messages_ui.append(f"  -> ERROR: 結論スライドの生成に失敗。{e}")
         report_slides_list.append({
                "slide_title": "結論と戦略的提言 (生成失敗)",
                "slide_layout": "title_and_content",
                "slide_content": [f"結論スライドの自動生成に失敗しました。", f"エラー: {e}"],
                "image_base64": None
            })

    # (★) 5. 最終的なJSON文字列を返す
    progress_bar.progress(1.0, text="Step C: 完了！")
    log_placeholder.text_area("実行ログ:", "\n".join(log_messages_ui[::-1]), height=250, key="step_c_log_done")
    
    return json.dumps(report_slides_list, ensure_ascii=False, indent=2)

def render_step_c():
    """(Step C) AIレポート生成UIを描画する"""
    st.title(f"🖋️ Step C: AI分析レポート生成") 

    # Step C 固有のセッションステート
    if 'step_c_jsonl_data' not in st.session_state:
        st.session_state.step_c_jsonl_data = None
    
    # (★) 改善 C-4: UIで編集する指示 (Task ⑨ に相当)
    if 'step_c_custom_instruction' not in st.session_state:
        st.session_state.step_c_custom_instruction = ""
        
    if 'step_c_report_json' not in st.session_state:
        st.session_state.step_c_report_json = None
    if 'step_c_model' not in st.session_state:
        st.session_state.step_c_model = MODEL_FLASH 
    if 'current_file_id_C' not in st.session_state:
        st.session_state.current_file_id_C = None

    # --- 1. ファイルアップロード ---
    st.header("Step 1: 分析データ (JSON) のアップロード")
    st.info("Step B でエクスポートした `analysis_data.json` をアップロードしてください。")
    uploaded_report_file = st.file_uploader(
        "分析データファイル (analysis_data.json)",
        type=['json', 'jsonl', 'txt'],
        key="step_c_uploader"
    )

    if uploaded_report_file:
        try:
            current_file_id_C = f"{uploaded_report_file.name}_{uploaded_report_file.size}"
            if st.session_state.get('current_file_id_C') != current_file_id_C:
                logger.info(f"Step C: 新しいJSON {current_file_id_C} をロードします。")
                jsonl_data_string = uploaded_report_file.getvalue().decode('utf-8')
                st.session_state.step_c_jsonl_data = jsonl_data_string
                st.session_state.step_c_report_json = None # (★) 結果をリセット
                st.session_state.current_file_id_C = current_file_id_C
                st.success(f"ファイル「{uploaded_report_file.name}」読込完了")
            
        except Exception as e:
            logger.error(f"Step C ファイル読込エラー: {e}", exc_info=True)
            st.error(f"ファイル読み込み中にエラー: {e}")
            st.session_state.step_c_jsonl_data = None
            st.session_state.current_file_id_C = None
            return
    else:
        st.session_state.step_c_jsonl_data = None
        st.session_state.step_c_report_json = None
        st.session_state.current_file_id_C = None
        st.warning("分析を続けるには、Step B で生成した JSON ファイルをアップロードしてください。")
        return

    # --- 2. 分析レポートの実行 ---
    st.header("Step 2: AI分析レポートの実行")

    # (★) --- 修正: モデル選択UI ---
    st.markdown("分析に使用するAIモデルを選択してください。")
    
    model_options = [MODEL_FLASH, MODEL_PRO]
    try:
        default_index = model_options.index(st.session_state.step_c_model)
    except ValueError:
        default_index = 0
        
    selected_model_name = st.radio(
        "使用モデル",
        options=model_options,
        index=default_index,
        key="step_c_model_radio",
        horizontal=True,
    )
    st.session_state.step_c_model = selected_model_name

    if selected_model_name == MODEL_PRO:
        st.warning(
            f"**`{MODEL_PRO}` (無料枠) は 2 RPM (30秒/リクエスト) の制限があります。**\n"
            f"スライド10枚の生成には約5分かかります。ご注意ください。"
        )
    else:
        st.info(
            f"**`{MODEL_FLASH}` (無料枠) は 10 RPM (6秒/リクエスト) の制限があります。**\n"
            f"比較的 高速に生成できます。（推奨）"
        )
    
    st.markdown("---")
    st.subheader("（オプション）AIへの追加指示")
    st.info("レポート全体を通してAIに意識させたい「分析の視点」や「特に注目すべき点」があれば入力してください。")
    
    st.session_state.step_c_custom_instruction = st.text_area(
        "AIへの追加指示（全体の分析方針）:",
        value=st.session_state.step_c_custom_instruction,
        placeholder="例: 「競合A社と比較した際の、我々の強み」に焦点を当てて考察してください。\n例: 今回の分析の目的は「若年層向けの新規施策立案」です。その視点で提言をまとめてください。",
        height=100,
        key="step_c_custom_instruction_input"
    )
    st.markdown("---")
    
    if st.button(f"分析レポートを生成 (Step 2)", key="execute_button_C", type="primary", use_container_width=True):
        if not st.session_state.step_c_jsonl_data:
            st.error("データがありません。Step 1でファイルをアップロードしてください。")
            return
        
        progress_bar = st.progress(0.0, text="Step C: 分析待機中...")
        log_placeholder = st.empty()

        selected_model = st.session_state.step_c_model
        
        try:
            st.session_state.step_c_report_json = run_step_c_analysis(
                st.session_state.step_c_jsonl_data,
                selected_model,
                progress_bar, 
                log_placeholder,
                st.session_state.step_c_custom_instruction 
            )
            st.success("AIによる分析レポートが生成されました！")
            
        except Exception as e:
            logger.error(f"Step C 実行中に予期せぬエラー: {e}", exc_info=True)
            st.error(f"分析実行中に予期せぬエラーが発生しました: {e}")
            progress_bar.progress(1.0, text="エラーにより中断")


    # --- 3. 結果のプレビューとエクスポート ---
    if st.session_state.step_c_report_json:
        st.header("Step 3: 分析レポート（JSON）の確認とエクスポート")
        st.info("以下の構造化JSONは、Step D (PowerPoint生成) で使用します。")

        st.download_button(
            label="分析レポート (report_for_powerpoint.json) をダウンロード",
            data=st.session_state.step_c_report_json,
            file_name="report_for_powerpoint.json",
            mime="application/json",
            type="primary",
            use_container_width=True
        )

        st.markdown("---")
        st.subheader("生成されたレポート プレビュー")
        
        try:
            report_data = json.loads(st.session_state.step_c_report_json)
            if isinstance(report_data, list) and all(isinstance(item, dict) for item in report_data):
                st.text_area(
                    "AIが生成した構造化JSON (プレビュー: 先頭5000文字)",
                    value=st.session_state.step_c_report_json[:5000] + "...",
                    height=300,
                    key="json_preview_C",
                    disabled=True
                )
                
                st.markdown("---")
                st.subheader(f"スライド構成 プレビュー ({len(report_data)}枚)")
                for i, slide in enumerate(report_data):
                    title = slide.get('slide_title', '（タイトルなし）')
                    layout = slide.get('slide_layout', 'N/A')
                    
                    slide_content_list = slide.get('slide_content')
                    if isinstance(slide_content_list, list) and slide_content_list:
                        content_preview = str(slide_content_list[0]) if slide_content_list[0] else "（空のコンテンツ）"
                    else:
                        content_preview = "（コンテンツなし）"

                    has_image = "有り" if slide.get("image_base64") else "無し"
                    
                    expander_label = f"**{i+1}: {title}** (Layout: {layout}, Image: {has_image})"
                    with st.expander(expander_label):
                        # (★) プレビューで Markdown をレンダリング
                        st.markdown(f"**内容:**")
                        if isinstance(slide_content_list, list):
                            for content_line in slide_content_list:
                                st.markdown(content_line)
                        else:
                            st.markdown(str(slide_content_list)) # フォールバック
            else:
                st.error("AIの回答が期待したスライドのリスト形式ではありません。")
                st.text_area("AIの生回答 (JSON):", value=st.session_state.step_c_report_json, height=200, disabled=True)
                
        except Exception as e:
            st.error(f"レポートのプレビュー中にエラー: {e}")
            st.text_area("AIの生回答 (パース失敗):", value=st.session_state.step_c_report_json, height=200, disabled=True)
            
        st.success("データをダウンロードし、Step D (PowerPoint生成) に進んでください。")


# (★) ---Step D---
try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.text.text import _Run # (★) Markdown削除のため
    from pptx.table import _Cell    # (★) Markdown削除のため
except ImportError:
    st.error(
        "PowerPoint生成ライブラリ(python-pptx)が見つかりません。"
        "pip install python-pptx を実行してください。"
    )

def add_markdown_text(text_frame, content_list: List[str]):
    """
    TextFrame (pptx) に、Markdown (太字) を解釈しながらテキストを追加する
    """
    if not text_frame or not content_list:
        return

    try:
        # 既存の段落をクリア (最初の1つは残す)
        tf = text_frame
        tf.clear()
        
        is_first_paragraph = True
        
        for item in content_list:
            if not isinstance(item, str):
                item = str(item)

            # (★) Markdownテーブル形式の簡易サポート
            if item.strip().startswith('|') and item.strip().endswith('|'):
                try:
                    p = tf.add_paragraph()
                    p.text = item # (★) テーブルはそのまま（フォント変更）
                    p.font.name = 'Yu Gothic' # (★) 等幅フォント推奨だが、日本語環境を優先
                    p.font.size = Pt(10)
                    continue
                except Exception:
                    pass # フォールバックして通常のテキスト処理
            
            # (★) Markdown (太字) の処理
            if is_first_paragraph:
                p = tf.paragraphs[0]
                is_first_paragraph = False
            else:
                p = tf.add_paragraph()

            # `**` で文字列を分割
            parts = item.split('**')
            
            for i, part in enumerate(parts):
                if not part: continue
                
                run = p.add_run()
                run.text = part
                
                # `**` で挟まれた奇数番目の部分 (i=1, 3, 5...) を太字にする
                if i % 2 == 1:
                    run.font.bold = True
                
    except Exception as e:
        logger.error(f"add_markdown_text 処理中にエラー: {e}", exc_info=True)
        # エラー時もフォールバック
        if 'p' in locals():
            p = text_frame.add_paragraph()
            p.text = "[Markdownの解析に失敗しました]"

def run_step_d_ai_correction(
    current_json_str: str, 
    correction_prompt: str
) -> str:
    """
    (Step D) ユーザーの修正指示に基づき、AI (Pro) がスライド構成JSONを修正する
    """
    logger.info("Step D AIスライド修正 (Pro) 実行...")
    
    # (★) Step D の AI修正は、高品質な Pro モデルを使用
    llm = get_llm(model_name=MODEL_PRO, temperature=0.1, timeout_seconds=120)
    if llm is None:
        logger.error("run_step_d_ai_correction: LLM (Pro) が利用できません。")
        st.error("AIモデル(Pro)が利用できませんでした。")
        return current_json_str # (★) 失敗時は元のJSONを返す

    # (★) --- [改善 D-3] プロンプトの強化 ---
    prompt = PromptTemplate.from_template(
        """
        あなたはPowerPointの構成作家です。
        以下の「現在のスライド構成 (JSON)」を読み、「修正指示」に基づいてJSONを厳密に修正してください。

        # 修正指示:
        {user_prompt}
        
        # (★) 修正の目的:
        * ユーザーは「資料のストーリーの流れを改善」したり「メッセージをより強調」するために指示を出しています。
        * （例: 「削除して」は、そのスライドが不要と判断されたためです）

        # 現在のスライド構成 (JSON):
        {current_json}

        # 指示:
        1. 「修正指示」を【厳密に】実行してください。(例: 「削除して」なら、そのJSONオブジェクトを配列から取り除く)
        2. 「修正指示」にないスライドは、絶対に修正・削除しないでください。
        3. 出力は【JSON配列形式のみ】 (`[...]`) とします。
        4. JSON以外のテキスト（「承知しました」など）は【絶対に】含めないでください。

        # 回答 (修正後のJSON配列のみ):
        """
    )
    # (★) --- ここまで ---
    
    chain = prompt | llm | StrOutputParser()
    
    try:
        response_str = chain.invoke({
            "user_prompt": correction_prompt,
            "current_json": current_json_str
        })
        
        # (★) 堅牢なJSONパース
        start = response_str.find('[')
        end = response_str.rfind(']')
        
        if start != -1 and end != -1 and end > start:
            json_str = response_str[start:end+1]
            logger.info("Step D AIスライド修正 完了。")
            return json_str
        else:
            logger.error("Step D AIスライド修正: AIがJSON配列を返しませんでした。")
            st.error("AIがJSON配列を返しませんでした。修正はキャンセルされました。")
            return current_json_str
            
    except Exception as e:
        logger.error(f"run_step_d_ai_correction 実行エラー: {e}", exc_info=True)
        st.error(f"AIスライド修正中にエラーが発生しました: {e}")
        return current_json_str


def create_powerpoint_presentation(
    template_file: Optional[BytesIO],
    report_data: List[Dict[str, Any]],
    layout_map_names: Dict[str, str]
) -> BytesIO:
    """
    (Step D) テンプレート(.pptx)とスライド構成(JSON)に基づき、
    python-pptx を使用して最終的なPowerPointファイルを生成する。
    (★) 改善: テンプレートの既存スライドを保持し、Markdown(太字)を反映
    """
    logger.info("PowerPoint生成処理 開始...")
    
    try:
        # (★) 1. テンプレートの読み込み
        if template_file:
            template_file.seek(0)
            prs = Presentation(template_file)
            logger.info("アップロードされたテンプレートを使用してPPTXを生成します。")
            
            # (★) --- [改善 D-2] テンプレートスライドの削除ロジックを削除 ---
            # (L2822-L2828 の削除)
            logger.info(f"テンプレートの既存スライド {len(prs.slides)} 枚を保持します。")
            # (★) --- ここまで ---

        else:
            prs = Presentation()
            logger.info("デフォルトのテンプレートを使用してPPTXを生成します。")

        # (★) 2. ユーザーが選択したレイアウトをマッピング (変更なし)
        layout_map = {
            "title_only": find_layout_by_name(prs, layout_map_names.get("title")),
            "agenda": find_layout_by_name(prs, layout_map_names.get("agenda")),
            "title_and_content": find_layout_by_name(prs, layout_map_names.get("content_text")),
            "text_and_image": find_layout_by_name(prs, layout_map_names.get("content_image")),
        }
        
        fallback_layout = prs.slide_layouts[1] # 「タイトルとコンテンツ」
        fallback_title_layout = prs.slide_layouts[0] # 「タイトル スライド」
        
        if layout_map["title_only"] is None:
             layout_map["title_only"] = fallback_title_layout
             logger.warning("「表紙」レイアウトが見つからないため、デフォルトの「タイトル スライド」を使用します。")
        if layout_map["agenda"] is None:
             layout_map["agenda"] = fallback_layout
             logger.warning("「目次」レイアウトが見つからないため、デフォルトの「タイトルとコンテンツ」を使用します。")
        if layout_map["title_and_content"] is None:
             layout_map["title_and_content"] = fallback_layout
             logger.warning("「テキスト」レイアウトが見つからないため、デフォルトの「タイトルとコンテンツ」を使用します。")
        if layout_map["text_and_image"] is None:
             layout_map["text_and_image"] = fallback_layout
             logger.warning("「画像+テキスト」レイアウトが見つからないため、デフォルトの「タイトルとコンテンツ」を使用します。")

        logger.info(f"使用レイアウトマッピング: {layout_map_names}")

        # (★) --- 3. スライドの生成 (JSONデータをループ) ---
        
        # (★) 3.1. 表紙スライド (変更なし)
        first_slide_data = report_data[0]
        if first_slide_data.get("slide_layout") == "title_only":
            slide = prs.slides.add_slide(layout_map["title_only"])
            try:
                slide.shapes.title.text = first_slide_data.get("slide_title", "分析レポート")
            except: pass
            try:
                if len(slide.placeholders) > 1 and slide.placeholders[1]:
                     # (★) [改善 D-1] 表紙のサブタイトルも Markdown ヘルパー経由に変更
                     add_markdown_text(
                         slide.placeholders[1].text_frame, 
                         first_slide_data.get("slide_content", [""])
                     )
            except: pass
            
            report_data = report_data[1:] # (★) 表紙をリストから削除
        
        # (★) 3.2. 目次(Agenda)スライドの自動生成 (Step C で生成済のため変更なし)

        # (★) --- 3.3. コンテンツスライド (残り) ---
        for i, slide_data in enumerate(report_data):
            slide_title = slide_data.get("slide_title", f"スライド {i+2}") 
            slide_layout_key = slide_data.get("slide_layout", "title_and_content")
            slide_content = slide_data.get("slide_content", ["（コンテンツなし）"])
            
            image_base64 = slide_data.get("image_base64")

            if image_base64 and slide_layout_key == "title_and_content":
                slide_layout_key = "text_and_image"
            
            if slide_title == "本日のアジェンダ":
                layout_to_use = layout_map["agenda"]
            elif image_base64:
                layout_to_use = layout_map["text_and_image"]
            else:
                layout_to_use = layout_map["title_and_content"]
            
            slide = prs.slides.add_slide(layout_to_use)
            
            try:
                slide.shapes.title.text = slide_title
            except Exception as e:
                logger.warning(f"スライド {i+2} のタイトル設定失敗: {e}")

            # (★) --- コンテンツと画像の配置 (ロジックを堅牢化) ---
            try:
                text_placeholders = []
                image_placeholders = []
                
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 0: continue
                    if shape.has_text_frame:
                        text_placeholders.append(shape)
                    elif shape.placeholder_format.idx > 100: 
                        image_placeholders.append(shape)

                # (★) 画像がある場合の処理 (text_and_image)
                if image_base64:
                    # (★) 1. テキストを挿入
                    if text_placeholders:
                        tf = text_placeholders[0].text_frame
                        # (★) --- [改善 D-1] Markdown ヘルパーを呼び出す ---
                        add_markdown_text(tf, slide_content)
                        # (★) --- (L2898 の re.sub を削除) ---
                    
                    # (★) 2. 画像を挿入
                    image_ph = None
                    if image_placeholders:
                        image_ph = image_placeholders[0]
                    elif len(text_placeholders) > 1:
                        image_ph = text_placeholders[1] 

                    if image_ph:
                        try:
                            img_bytes = base64.b64decode(image_base64)
                            img_stream = BytesIO(img_bytes)
                            image_ph.insert_picture(img_stream)
                            logger.info(f"スライド '{slide_title}': グラフ画像の挿入に成功。")
                        except Exception as e:
                            logger.error(f"スライド '{slide_title}': グラフ画像の挿入に失敗: {e}")
                            if image_ph.has_text_frame:
                                image_ph.text_frame.text = f"（画像挿入エラー: {e}）"
                    else:
                         logger.warning(f"スライド '{slide_title}': 画像用プレースホルダが見つかりません。")

                # (★) 画像がない場合の処理 (title_and_content)
                else:
                    if not text_placeholders:
                         logger.warning(f"スライド '{slide_title}': コンテンツプレースホルダが見つかりません。")
                         continue
                    tf = text_placeholders[0].text_frame
                    # (★) --- [改善 D-1] Markdown ヘルパーを呼び出す ---
                    add_markdown_text(tf, slide_content)
                    # (★) --- (L2928 の re.sub を削除) ---

            except Exception as e:
                logger.error(f"スライド {i+2} ('{slide_title}') のコンテンツ/画像設定中にエラー: {e}", exc_info=True)

        logger.info("PowerPoint生成処理 完了。")
        file_stream = BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream

    except Exception as e:
        logger.error(f"create_powerpoint_presentation 全体でエラー: {e}", exc_info=True)
        st.error(f"PowerPointの生成に失敗しました: {e}")
        return None

def render_step_d():
    """(Step D) PowerPoint生成UIを描画する"""
    st.title(f"プレゼンテーション (PowerPoint) 生成 (Step D)")

    if 'step_d_template_file' not in st.session_state:
        st.session_state.step_d_template_file = None
    if 'step_d_template_file_id' not in st.session_state:
        st.session_state.step_d_template_file_id = None
    if 'step_d_report_data' not in st.session_state:
        st.session_state.step_d_report_data = []
    if 'current_report_file_id_D' not in st.session_state:
        st.session_state.current_report_file_id_D = None
        
    if 'step_d_generated_pptx' not in st.session_state:
        st.session_state.step_d_generated_pptx = None
    if 'step_d_layout_map' not in st.session_state:
        st.session_state.step_d_layout_map = {}
    if 'tips_list' not in st.session_state:
        st.session_state.tips_list = []
    if 'current_tip_index' not in st.session_state:
        st.session_state.current_tip_index = 0
    if 'last_tip_time' not in st.session_state:
        st.session_state.last_tip_time = time.time()

    # --- 1. テンプレートのアップロード (★) ---
    st.header("Step 1: テンプレート PowerPoint のアップロード")
    st.info(
        "（オプション）使用したい .pptx テンプレートがあればアップロードしてください。\n"
        "AIが生成したスライドは、テンプレート内の既存スライド（表紙など）の「後」に追加されます。"
    )
    
    template_file = st.file_uploader(
        "PowerPoint テンプレート (.pptx)",
        type=['pptx'],
        key="step_d_template_uploader"
    )
    
    template_layout_names = []
    default_layouts = {
        "title": "タイトル スライド",
        "agenda": "セクション見出し",
        "content_text": "タイトルとコンテンツ",
        "content_image": "2つのコンテンツ"
    }
    
    if template_file:
        try:
            template_file_id = f"{template_file.name}_{template_file.size}"
            
            if st.session_state.get('step_d_template_file_id') != template_file_id:
                logger.info(f"Step D: 新しいテンプレート {template_file_id} をロードします。")
                template_file.seek(0)
                template_bytes = template_file.getvalue()
                prs = Presentation(BytesIO(template_bytes))
                template_layout_names = [layout.name for layout in prs.slide_layouts]
                
                st.success(f"テンプレート「{template_file.name}」を読み込みました。")
                st.session_state.step_d_template_file = BytesIO(template_bytes)
                st.session_state.step_d_template_file_id = template_file_id
                st.session_state.step_d_layout_map = {} 
            else:
                st.session_state.step_d_template_file.seek(0)
                prs = Presentation(st.session_state.step_d_template_file)
                template_layout_names = [layout.name for layout in prs.slide_layouts]

        except Exception as e:
            st.error(f"テンプレートの読み込みに失敗: {e}")
            template_layout_names = []
            st.session_state.step_d_template_file = None
            st.session_state.step_d_template_file_id = None
            
    else:
        if st.session_state.step_d_template_file is not None:
             st.session_state.step_d_template_file = None
             st.session_state.step_d_template_file_id = None
             st.session_state.step_d_layout_map = {}


    # --- 2. Step C 分析結果のアップロード ---
    st.header("Step 2: Step C 分析レポート (JSON) のアップロード")
    st.info("Step C でエクスポートした `report_for_powerpoint.json` をアップロードしてください。")
    report_file = st.file_uploader(
        "分析レポートファイル (report_for_powerpoint.json)",
        type=['json'],
        key="step_d_report_uploader"
    )

    if report_file:
        try:
            current_report_file_id = f"{report_file.name}_{report_file.size}"
            
            if ('step_d_report_data' not in st.session_state or 
                not st.session_state.step_d_report_data or 
                st.session_state.get('current_report_file_id_D') != current_report_file_id):
                
                logger.info(f"Step D: 新しいレポート {current_report_file_id} をロードします。")
                report_json_string = report_file.getvalue().decode('utf-8')
                report_data = json.loads(report_json_string)
            
                if isinstance(report_data, list) and all(isinstance(item, dict) for item in report_data):
                    st.success(f"分析レポート「{report_file.name}」を読み込みました ({len(report_data)}スライド)。")
                    st.session_state.step_d_report_data = report_data
                    st.session_state.current_report_file_id_D = current_report_file_id
                    st.session_state.step_d_generated_pptx = None
                else:
                    st.error("アップロードされたJSONが期待する形式（スライドのリスト）ではありません。")
                    st.session_state.step_d_report_data = []
                    st.session_state.current_report_file_id_D = None

        except Exception as e:
            logger.error(f"Step D JSONレポート読込エラー: {e}", exc_info=True)
            st.error(f"分析レポートの読み込み中にエラー: {e}")
            st.session_state.step_d_report_data = []
            st.session_state.current_report_file_id_D = None
    
    if not st.session_state.step_d_report_data:
        st.session_state.step_d_report_data = []
        st.session_state.step_d_generated_pptx = None
        st.session_state.current_report_file_id_D = None
        st.warning("PowerPointを生成するには、Step C で生成した JSON レポートをアップロードしてください。")
        return

    # --- 3. (★) テンプレートのレイアウト割り当て---
    st.header("Step 3: テンプレートレイアウトの割り当て")
    
    if not st.session_state.step_d_template_file:
        st.info("Step 1 でテンプレートをアップロードすると、レイアウト名を選択できます。（現在デフォルト設定）")
        layout_options = list(default_layouts.values())
    else:
        st.info("テンプレートから読み込んだレイアウト名を、各スライドタイプに割り当ててください。")
        layout_options = template_layout_names
        
    if not layout_options:
         st.error("レイアウトの読み込みに失敗しました。デフォルト設定を使用します。")
         layout_options = list(default_layouts.values())

    def get_default_index(default_name_key):
        if default_name_key in st.session_state.step_d_layout_map:
            saved_name = st.session_state.step_d_layout_map[default_name_key]
            if saved_name in layout_options:
                return layout_options.index(saved_name)
        
        target_name = default_layouts[default_name_key]
        if target_name in layout_options:
            return layout_options.index(target_name)
            
        for i, opt in enumerate(layout_options):
            if default_name_key in opt.lower():
                return i
            if target_name.split(' ')[0] in opt:
                return i
        return 0

    layout_map = {}
    col1, col2 = st.columns(2)
    with col1:
        layout_map["title"] = st.selectbox(
            "1. 表紙 (Title) スライド:", layout_options, 
            index=get_default_index("title"), key="layout_select_title"
        )
        layout_map["agenda"] = st.selectbox(
            "2. 目次 (Agenda) スライド:", layout_options, 
            index=get_default_index("agenda"), key="layout_select_agenda"
        )
    with col2:
        layout_map["content_text"] = st.selectbox(
            "3. 分析 (テキストのみ) スライド:", layout_options, 
            index=get_default_index("content_text"), key="layout_select_text"
        )
        layout_map["content_image"] = st.selectbox(
            "4. 分析 (テキスト+画像) スライド:", layout_options, 
            index=get_default_index("content_image"), key="layout_select_image"
        )
    
    st.session_state.step_d_layout_map = layout_map


    # --- 4. スライド構成の編集 ---
    st.header("Step 4: スライド構成の確認・編集")
    st.info("（(★) マウスのドラッグ＆ドロップでスライドの順番を入れ替えることができます）")

    try:
        headers_list = []
        header_to_item_map = {}
        
        if not st.session_state.step_d_report_data:
             st.warning("JSONデータが空か、正しく読み込まれていません。")
             return

        for i, item in enumerate(st.session_state.step_d_report_data):
            if not isinstance(item, dict):
                st.error(f"データ形式エラー: {item} は辞書ではありません。")
                continue
            
            title = item.get('slide_title', '（タイトルなし）')
            layout = item.get('slide_layout', 'N/A')
            has_image = "🖼️" if (item.get("image_base64")) else "📄"
            header_str = f"**{i+1}: {title}** (Layout: `{layout}`, {has_image})"
            headers_list.append(header_str)
            header_to_item_map[header_str] = item

        if not all(isinstance(h, str) for h in headers_list):
            st.error("内部エラー: ヘッダーリストの作成に失敗しました。")
            return

        sorted_headers = sort_items(
            items=headers_list,
            key="sortable_slides_v4"
        )
        
        cleaned_sorted_data = []
        for header in sorted_headers:
            if header in header_to_item_map:
                cleaned_sorted_data.append(header_to_item_map[header])
            else:
                logger.error(f"マッピングエラー: ソート後のヘッダー '{header}' が見つかりません。")
            
        st.session_state.step_d_report_data = cleaned_sorted_data
        
    except Exception as e:
        logger.error(f"streamlit-sortables 実行エラー: {e}", exc_info=True)
        st.error(f"スライド編集UIの描画に失敗: {e}。")


    # --- 5. AIによる修正指示  ---
    st.header("Step 5: (オプション) AIによる内容の修正指示")
    st.markdown(f"（(★) 使用モデル: `{MODEL_PRO}`）")
    
    with st.expander("AIにスライド内容の修正を指示する"):
        correction_prompt = st.text_area(
            "修正内容を具体的に指示してください:",
            placeholder=(
                "例: 「エグゼクティブ・サマリー」スライドの箇条書きを3点に要約して。\n"
                "例: 「共起ネットワーク」スライドを削除して。"
            ),
            key="step_d_correction_prompt"
        )
        
        if st.button("AIでスライド構成を修正", key="run_ai_correction_D", type="secondary"):
            if correction_prompt.strip():
                with st.spinner(f"AI ({MODEL_PRO}) がスライド構成 (JSON) を修正中..."):
                    current_json_str = json.dumps(st.session_state.step_d_report_data, ensure_ascii=False)
                    corrected_json_str = run_step_d_ai_correction(current_json_str, correction_prompt)
                    
                    try:
                        corrected_data = json.loads(corrected_json_str)
                        if isinstance(corrected_data, list):
                            st.session_state.step_d_report_data = corrected_data
                            st.success("AIによるスライド構成の修正が完了しました。Step 4 の構成が更新されています。")
                            st.rerun()
                        else:
                            st.error("AIがリスト形式でないデータを返しました。修正はキャンセルされました。")
                    except Exception as e:
                        st.error(f"AIの回答のパースに失敗: {e}。修正はキャンセルされました。")
            else:
                st.warning("修正指示を入力してください。")

    # --- 6. PowerPoint生成 ---
    st.header("Step 6: PowerPointの生成とエクスポート")
    
    tip_placeholder_d = st.empty()
    
    if st.button("PowerPointを生成 (Step 6)", key="generate_pptx_D", type="primary", use_container_width=True):
        st.session_state.step_d_generated_pptx = None
        
        if not st.session_state.tips_list or len(st.session_state.tips_list) <= 1:
            with st.spinner("分析TIPSをAIで生成中..."):
                st.session_state.tips_list = get_analysis_tips_list_from_ai()
                st.session_state.current_tip_index = random.randint(0, len(st.session_state.tips_list) - 1) if st.session_state.tips_list else 0
                st.session_state.last_tip_time = time.time()
        
        with st.spinner("PowerPointファイルを生成中..."):
            now = time.time()
            if (now - st.session_state.last_tip_time > 10):
                if len(st.session_state.tips_list) > 1:
                    st.session_state.current_tip_index = (st.session_state.current_tip_index + 1) % len(st.session_state.tips_list)
                st.session_state.last_tip_time = now
            if st.session_state.tips_list:
                try:
                    current_tip = st.session_state.tips_list[st.session_state.current_tip_index]
                    tip_placeholder_d.info(f"💡 データ分析TIPS: {current_tip}")
                except IndexError:
                    st.session_state.current_tip_index = 0

            generated_file_stream = create_powerpoint_presentation(
                st.session_state.step_d_template_file,
                st.session_state.step_d_report_data,
                st.session_state.step_d_layout_map
            )
            
            tip_placeholder_d.empty()
            
            if generated_file_stream:
                st.session_state.step_d_generated_pptx = generated_file_stream.getvalue()
                st.success("PowerPointファイルの生成が完了しました。")
            else:
                st.error("PowerPointファイルの生成に失敗しました。")

    if st.session_state.step_d_generated_pptx:
        st.download_button(
            label="生成された PowerPoint をダウンロード",
            data=st.session_state.step_d_generated_pptx,
            file_name="AI_Analysis_Report_v3.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
        st.balloons()

# --- 11. (★) Main関数 (アプリケーション実行) ---
def main():
    """Streamlitアプリケーションのメイン実行関数"""
    
    # (★) --- st.set_page_config() を最初に実行 ---
    st.set_page_config(page_title="AI Data Analysis App", layout="wide")
    
    if 'log_messages' not in st.session_state:
        st.session_state.log_messages = []
    if 'current_step' not in st.session_state:
        st.session_state.current_step = 'A' # 初期ステップ
    if 'tips_list' not in st.session_state:
        st.session_state.tips_list = []
    if 'current_tip_index' not in st.session_state:
        st.session_state.current_tip_index = 0
    if 'last_tip_time' not in st.session_state:
        st.session_state.last_tip_time = time.time()

    # (★) --- .envファイルから環境変数を読み込む ---
    try:
        load_dotenv()
        logger.info(".env ファイルの読み込み試行完了。")
    except Exception as e:
        logger.warning(f".env ファイルの読み込みに失敗: {e}") 
    
    # --- サイドバー (ナビゲーション) ---
    with st.sidebar:
        st.title("AI レポーティング App")
        st.markdown("---")
        
        st.header("⚙️ AI 設定")
        
        if not os.getenv("GOOGLE_API_KEY"):
            st.warning(
                "Google APIキーが.envファイルに設定されていないか、読み込めませんでした。\n\n"
                "(.envファイルに `GOOGLE_API_KEY='あなたのAPIキー'` と記載してください)"
            )
        else:
            st.success("Google APIキー 読込完了")
            # (★) アプリ起動時にLLMとspaCyのロードを試みる
            if 'llm_checked' not in st.session_state:
                if get_llm(MODEL_FLASH_LITE) is None: 
                    st.error("LLMの初期化に失敗。APIキーが正しいか確認してください。")
                if load_spacy_model() is None:
                    st.error("spaCyモデルのロードに失敗。")
                st.session_state.llm_checked = True # (★) 毎リラン時にチェックしないよう
        
        st.markdown("---")
        
        # (★) --- Step A〜D のナビゲーションボタン ---
        st.header("🔄 ナビゲーション")
        current_step = st.session_state.current_step
        
        if st.button(
            "Step A: AIタグ付け", key="nav_A", use_container_width=True,
            type=("primary" if current_step == 'A' else "secondary")
        ):
            if st.session_state.current_step != 'A':
                st.session_state.current_step = 'A'; st.rerun()

        if st.button(
            "Step B: 分析実行", key="nav_B", use_container_width=True,
            type=("primary" if current_step == 'B' else "secondary")
        ):
            if st.session_state.current_step != 'B':
                st.session_state.current_step = 'B'; st.rerun()

        if st.button(
            "Step C: AIレポート生成", key="nav_C", use_container_width=True,
            type=("primary" if current_step == 'C' else "secondary")
        ):
            if st.session_state.current_step != 'C':
                st.session_state.current_step = 'C'; st.rerun()

        if st.button(
            "Step D: PowerPoint生成", key="nav_D", use_container_width=True,
            type=("primary" if current_step == 'D' else "secondary")
        ):
            if st.session_state.current_step != 'D':
                st.session_state.current_step = 'D'; st.rerun()
                
    # (★) 既存の main() 関数のロジック (セッションステート初期化) を移動
    if 'llm_checked' not in st.session_state:
        if os.getenv("GOOGLE_API_KEY"):
            if get_llm(MODEL_FLASH_LITE) is None: 
                pass # (★) エラーはサイドバーで表示
            if load_spacy_model() is None:
                pass # (★) エラーはサイドバーで表示
        st.session_state.llm_checked = True


    # --- メインコンテンツ (ステップに応じて描画) ---
    if st.session_state.current_step == 'A':
        render_step_a()
    elif st.session_state.current_step == 'B':
        render_step_b()
    elif st.session_state.current_step == 'C':
        render_step_c()
    elif st.session_state.current_step == 'D':
        render_step_d()
    else:
        st.error("不明なステップです。Step Aに戻ります。")
        st.session_state.current_step = 'A'; st.rerun()

if __name__ == "__main__":
    main()